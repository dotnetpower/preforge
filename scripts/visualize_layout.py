#!/usr/bin/env python3
"""PPT 슬라이드의 그리드 레이아웃을 시각화하는 스크립트"""

from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
import subprocess
import tempfile
import os
import sys
from dataclasses import dataclass
from typing import List, Dict, Tuple
from contextlib import nullcontext
import numpy as np
from skimage.metrics import structural_similarity as ssim
from skimage import filters
import cv2

# 프로젝트 경로 추가
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from preforge.parsers.pptx_parser import PptxParser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu
import copy


def remove_headers_footers_from_pptx(pptx_path: Path, output_path: Path) -> Path:
    """PPTX에서 헤더, 푸터, 페이지번호를 제거한 복사본 생성
    
    Args:
        pptx_path: 원본 PPTX 파일 경로
        output_path: 출력 PPTX 파일 경로
        
    Returns:
        출력 파일 경로
    """
    prs = Presentation(str(pptx_path))
    slide_height = prs.slide_height
    
    def should_remove_shape(shape) -> bool:
        """shape가 헤더/푸터/페이지번호인지 판단"""
        # 플레이스홀더 타입 확인
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 헤더, 푸터, 슬라이드 번호, 날짜 플레이스홀더 제거
            if ph_type in [
                PP_PLACEHOLDER.FOOTER,           # 푸터
                PP_PLACEHOLDER.SLIDE_NUMBER,     # 슬라이드 번호
                PP_PLACEHOLDER.DATE,             # 날짜
                PP_PLACEHOLDER.HEADER,           # 헤더 (노트/유인물용)
            ]:
                return True
        
        # 위치 기반 추가 검사 (하단 8%)
        if shape.top is not None and shape.height is not None:
            bottom_ratio = (shape.top + shape.height) / slide_height
            
            # 하단 영역의 작은 텍스트 박스 (푸터/페이지번호일 가능성)
            if bottom_ratio > 0.92 and hasattr(shape, 'text'):
                text = shape.text.strip()
                # 페이지 번호 패턴 또는 짧은 텍스트 (‹#› 포함)
                if text.isdigit() or text == '‹#›' or len(text) < 20:
                    return True
        
        return False
    
    def remove_shapes_from_container(container, name=""):
        """컨테이너(슬라이드, 마스터, 레이아웃)에서 헤더/푸터 shape 제거"""
        shapes_to_remove = []
        for shape in container.shapes:
            if should_remove_shape(shape):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        if shapes_to_remove:
            print(f"    {name}: {len(shapes_to_remove)}개 요소 제거")
    
    # 1. 슬라이드 마스터에서 제거
    for i, master in enumerate(prs.slide_masters):
        remove_shapes_from_container(master, f"마스터 {i}")
        
        # 2. 슬라이드 레이아웃에서 제거
        for j, layout in enumerate(master.slide_layouts):
            remove_shapes_from_container(layout, f"레이아웃 {i}-{j}")
    
    # 3. 각 슬라이드에서 제거
    for i, slide in enumerate(prs.slides):
        remove_shapes_from_container(slide, f"슬라이드 {i+1}")
    
    # 수정된 PPTX 저장
    prs.save(str(output_path))
    return output_path


@dataclass
class ContentBox:
    """컨텐츠 영역 정보"""
    left: int
    top: int
    width: int
    height: int
    content_type: str  # 'text', 'table', 'image'


def is_section_title(shape, slide_height: int) -> bool:
    """섹션 제목인지 판별 (예: '1. 질병', '2. 제품')
    
    특징:
    - 상단 영역에 위치 (상위 12%)
    - '숫자. 한글' 패턴
    - 좌측에 위치
    """
    import re
    
    # 상단 영역 확인
    if shape.top > slide_height * 12 // 100:
        return False
    
    # 텍스트 확인
    if not shape.has_text_frame:
        return False
    
    text = shape.text_frame.text.strip()
    
    # 섹션 제목 패턴: "숫자. 한글" (예: "1. 질병", "2. 제품")
    if re.match(r'^\d+\.\s*[가-힣]+$', text):
        return True
    
    return False


def is_header_line(shape, slide_height: int, slide_width: int) -> bool:
    """헤더 아래 구분선인지 판별
    
    특징:
    - 상단 영역에 위치 (상위 12%)
    - 너비가 슬라이드의 80% 이상
    - 높이가 매우 작음 (선)
    """
    # 상단 영역 확인
    if shape.top > slide_height * 12 // 100:
        return False
    
    # 넓고 얇은 선 확인
    if shape.width >= slide_width * 80 // 100 and shape.height < slide_height * 1 // 100:
        return True
    
    return False


def get_slide_content_boxes(pptx_path: Path, page_number: int, exclude_master: bool = True) -> Tuple[List[ContentBox], int, int]:
    """슬라이드의 모든 컨텐츠 영역 추출
    
    Args:
        pptx_path: PPTX 파일 경로
        page_number: 페이지 번호 (1부터 시작)
        exclude_master: True이면 슬라이드 마스터 요소 제외 (로고, Confidential, 페이지 번호 등)
    """
    prs = Presentation(str(pptx_path))
    slide = prs.slides[page_number - 1]
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 마스터 요소 제외를 위한 임계값
    section_header_threshold = slide_height * 12 // 100  # 상단 12% (섹션 제목 영역)
    footer_threshold = slide_height * 92 // 100   # 하단 8%
    left_margin = slide_width * 3 // 100          # 좌측 3%
    right_margin = slide_width * 85 // 100        # 우측 15%
    
    boxes = []
    
    for shape in slide.shapes:
        # 마스터 요소 필터링
        if exclude_master:
            # 플레이스홀더 타입으로 필터링 (페이지 번호, 날짜, 푸터 등)
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    # 슬라이드 번호, 날짜, 푸터는 제외
                    if ph_type in (PP_PLACEHOLDER.SLIDE_NUMBER, 
                                   PP_PLACEHOLDER.DATE,
                                   PP_PLACEHOLDER.FOOTER):
                        continue
                except:
                    pass
            
            # 섹션 제목 제외 (예: "1. 질병")
            if is_section_title(shape, slide_height):
                continue
            
            # 헤더 구분선 제외
            if is_header_line(shape, slide_height, slide_width):
                continue
            
            # 위치로 필터링 - 하단 영역 (푸터, 로고, 페이지 번호)
            if shape.top >= footer_threshold:
                continue
            
            # 오른쪽 상단 작은 요소 (Confidential 등)
            if (shape.top < section_header_threshold and 
                shape.left >= right_margin and
                shape.width < slide_width * 20 // 100):
                continue
            
            # 오른쪽 하단 작은 요소 (로고 등)
            if (shape.top >= slide_height * 80 // 100 and
                shape.left >= right_margin):
                continue
        
        # 컨텐츠 타입 판별
        content_type = None
        if shape.has_table:
            content_type = 'table'
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            content_type = 'image'
        elif shape.has_text_frame and shape.text_frame.text.strip():
            content_type = 'text'
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            content_type = 'group'
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            content_type = 'shape'
        
        if content_type:
            boxes.append(ContentBox(
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=shape.height,
                content_type=content_type
            ))
    
    return boxes, slide_width, slide_height


def get_main_content_bounding_box(pptx_path: Path, page_number: int) -> Tuple[int, int, int, int, int, int]:
    """슬라이드의 메인 컨텐츠 전체를 감싸는 단일 바운딩 박스 반환
    
    슬라이드 마스터 요소(로고, Confidential, 페이지 번호 등)를 제외하고
    실제 컨텐츠만의 바운딩 박스를 계산합니다.
    
    Returns:
        (left, top, right, bottom, slide_width, slide_height)
    """
    boxes, slide_width, slide_height = get_slide_content_boxes(pptx_path, page_number, exclude_master=True)
    
    if not boxes:
        return (0, 0, slide_width, slide_height, slide_width, slide_height)
    
    min_left = min(box.left for box in boxes)
    min_top = min(box.top for box in boxes)
    max_right = max(box.left + box.width for box in boxes)
    max_bottom = max(box.top + box.height for box in boxes)
    
    return (min_left, min_top, max_right, max_bottom, slide_width, slide_height)


def draw_single_content_box(
    image: Image.Image,
    pptx_path: Path,
    page_number: int,
    box_color: str = "red",
    line_width: int = 3,
    padding: int = 10
) -> Image.Image:
    """메인 컨텐츠 전체를 하나의 박스로 그리기 (마스터 요소 제외)"""
    draw = ImageDraw.Draw(image)
    img_width, img_height = image.size
    
    # 바운딩 박스 가져오기
    left, top, right, bottom, slide_width, slide_height = get_main_content_bounding_box(pptx_path, page_number)
    
    # EMU를 이미지 픽셀로 변환
    scale_x = img_width / slide_width
    scale_y = img_height / slide_height
    
    # 픽셀 좌표로 변환
    x1 = int(left * scale_x) - padding
    y1 = int(top * scale_y) - padding
    x2 = int(right * scale_x) + padding
    y2 = int(bottom * scale_y) + padding
    
    # 경계 처리
    x1 = max(line_width, x1)
    y1 = max(line_width, y1)
    x2 = min(img_width - line_width, x2)
    y2 = min(img_height - line_width, y2)
    
    # 박스 그리기
    draw.rectangle([x1, y1, x2, y2], outline=box_color, width=line_width)
    
    return image


def group_boxes_by_grid(
    boxes: List[ContentBox], 
    slide_width: int, 
    slide_height: int,
    rows: int = 1,
    cols: int = 2
) -> Dict[Tuple[int, int], List[ContentBox]]:
    """박스들을 그리드 셀로 그룹화
    
    이미 마스터 요소가 필터링된 상태의 박스들을 받아서
    그리드 셀 위치에 따라 그룹화합니다.
    """
    mid_x = slide_width // 2
    
    grouped = {}
    
    for box in boxes:
        # 열 결정 (좌우 분할)
        center_x = box.left + box.width // 2
        col = 0 if center_x < mid_x else 1
        
        # 행 결정 (단순화: 1행으로 가정)
        row = 0
        
        cell = (row, col) if cols > 1 else (row, 0)
        
        if cell not in grouped:
            grouped[cell] = []
        grouped[cell].append(box)
    
    return grouped


def find_content_bounds_pixel(
    image: Image.Image,
    region: Tuple[int, int, int, int],
    bg_threshold: int = 250,
    margin: int = 5
) -> Tuple[int, int, int, int]:
    """픽셀 기반으로 실제 컨텐츠 영역을 찾기
    
    Args:
        image: 슬라이드 이미지
        region: 검사할 영역 (x1, y1, x2, y2)
        bg_threshold: 배경으로 간주할 밝기 임계값 (흰색에 가까움)
        margin: 컨텐츠 경계에 추가할 여백
    
    Returns:
        (x1, y1, x2, y2): 실제 컨텐츠의 바운딩 박스
    """
    import numpy as np
    
    x1, y1, x2, y2 = region
    
    # 영역 유효성 검사
    img_width, img_height = image.size
    x1 = max(0, min(x1, img_width - 1))
    y1 = max(0, min(y1, img_height - 1))
    x2 = max(x1 + 1, min(x2, img_width))
    y2 = max(y1 + 1, min(y2, img_height))
    
    # 영역 크롭
    cropped = image.crop((x1, y1, x2, y2))
    
    # RGB로 변환
    if cropped.mode != 'RGB':
        cropped = cropped.convert('RGB')
    
    # numpy 배열로 변환
    pixels = np.array(cropped)
    
    # 배경이 아닌 픽셀 찾기 (R, G, B 모두 threshold 이하인 픽셀이 하나라도 있으면 컨텐츠)
    # 흰색 배경: (255, 255, 255) 또는 그에 가까운 값
    is_content = np.any(pixels < bg_threshold, axis=2)
    
    # 컨텐츠가 있는 행과 열 찾기
    content_rows = np.any(is_content, axis=1)
    content_cols = np.any(is_content, axis=0)
    
    if not np.any(content_rows) or not np.any(content_cols):
        # 컨텐츠가 없으면 원래 영역 반환
        return region
    
    # 컨텐츠의 바운딩 박스 찾기
    row_indices = np.where(content_rows)[0]
    col_indices = np.where(content_cols)[0]
    
    content_y1 = row_indices[0]
    content_y2 = row_indices[-1] + 1
    content_x1 = col_indices[0]
    content_x2 = col_indices[-1] + 1
    
    # 원래 좌표계로 변환하고 여백 추가
    # 단, 원래 검색 영역(region)을 벗어나지 않도록 제한
    orig_x1, orig_y1, orig_x2, orig_y2 = region
    
    result = (
        max(orig_x1, x1 + content_x1 - margin),  # 검색 영역 왼쪽 경계 이상
        max(orig_y1, y1 + content_y1 - margin),  # 검색 영역 위쪽 경계 이상
        min(orig_x2, x1 + content_x2 + margin),  # 검색 영역 오른쪽 경계 이하
        min(orig_y2, y1 + content_y2 + margin)   # 검색 영역 아래쪽 경계 이하
    )
    
    return result


def adjust_box_to_avoid_content(
    image: Image.Image,
    box: Tuple[int, int, int, int],
    line_width: int = 4,
    bg_threshold: int = 250,
    max_adjust: int = 20
) -> Tuple[int, int, int, int]:
    """박스 테두리가 컨텐츠와 겹치지 않도록 조정
    
    Args:
        image: 슬라이드 이미지
        box: 박스 좌표 (x1, y1, x2, y2)
        line_width: 박스 선 두께
        bg_threshold: 배경으로 간주할 밝기 임계값
        max_adjust: 최대 조정 거리 (픽셀)
    
    Returns:
        조정된 박스 좌표
    """
    import numpy as np
    
    x1, y1, x2, y2 = box
    img_width, img_height = image.size
    
    # RGB로 변환
    if image.mode != 'RGB':
        img_rgb = image.convert('RGB')
    else:
        img_rgb = image
    
    pixels = np.array(img_rgb)
    
    # 컨텐츠 감지 (배경이 아닌 픽셀)
    is_content = np.any(pixels < bg_threshold, axis=2)
    
    # 각 테두리를 확인하고 조정
    check_width = line_width + 2  # 선 두께 + 여유
    
    # 상단 테두리 확인 (y1 주변)
    for offset in range(max_adjust):
        y_check = max(0, y1 - offset)
        y_start = max(0, y_check - check_width // 2)
        y_end = min(img_height, y_check + check_width // 2)
        x_start = max(0, x1)
        x_end = min(img_width, x2)
        
        if y_start >= y_end or x_start >= x_end:
            continue
        
        # 이 영역에 컨텐츠가 있는지 확인
        region_content = is_content[y_start:y_end, x_start:x_end]
        if not np.any(region_content):
            y1 = y_check
            break
    
    # 하단 테두리 확인 (y2 주변)
    for offset in range(max_adjust):
        y_check = min(img_height, y2 + offset)
        y_start = max(0, y_check - check_width // 2)
        y_end = min(img_height, y_check + check_width // 2)
        x_start = max(0, x1)
        x_end = min(img_width, x2)
        
        if y_start >= y_end or x_start >= x_end:
            continue
        
        region_content = is_content[y_start:y_end, x_start:x_end]
        if not np.any(region_content):
            y2 = y_check
            break
    
    # 좌측 테두리 확인 (x1 주변)
    for offset in range(max_adjust):
        x_check = max(0, x1 - offset)
        x_start = max(0, x_check - check_width // 2)
        x_end = min(img_width, x_check + check_width // 2)
        y_start = max(0, y1)
        y_end = min(img_height, y2)
        
        if x_start >= x_end or y_start >= y_end:
            continue
        
        region_content = is_content[y_start:y_end, x_start:x_end]
        if not np.any(region_content):
            x1 = x_check
            break
    
    # 우측 테두리 확인 (x2 주변)
    for offset in range(max_adjust):
        x_check = min(img_width, x2 + offset)
        x_start = max(0, x_check - check_width // 2)
        x_end = min(img_width, x_check + check_width // 2)
        y_start = max(0, y1)
        y_end = min(img_height, y2)
        
        if x_start >= x_end or y_start >= y_end:
            continue
        
        region_content = is_content[y_start:y_end, x_start:x_end]
        if not np.any(region_content):
            x2 = x_check
            break
    
    return (x1, y1, x2, y2)


def detect_content_with_ssim(
    image: Image.Image, 
    region: Tuple[int, int, int, int],
    bg_threshold: int = 245,  # 더 엄격하게 (250 → 245)
    edge_threshold: float = 0.02  # 엣지 임계값
) -> np.ndarray:
    """SSIM을 사용하여 컨텐츠 영역 감지
    
    배경과의 구조적 유사도를 비교하여 컨텐츠를 더 정확하게 감지합니다.
    
    Args:
        image: PIL 이미지
        region: 분석할 영역 (x1, y1, x2, y2)
        bg_threshold: 배경으로 간주할 픽셀 밝기 임계값 (더 낮을수록 엄격)
        edge_threshold: 엣지 검출 임계값
    
    Returns:
        컨텐츠 영역을 나타내는 boolean numpy 배열
    """
    x1, y1, x2, y2 = region
    
    # 영역 추출
    region_img = image.crop((x1, y1, x2, y2))
    gray = np.array(region_img.convert('L'))
    
    # 1. 기본 임계값 기반 감지 (엄격하게)
    basic_content = gray < bg_threshold
    
    # 2. Edge 검출로 컨텐츠 경계 강조 (텍스트, 이미지 테두리)
    edges = filters.sobel(gray)
    edge_content = edges > edge_threshold
    
    # 3. 통합 감지 (OR 연산)
    combined_content = basic_content | edge_content
    
    # 4. 작은 노이즈 제거 (매우 작은 영역만)
    kernel = np.ones((2, 2), np.uint8)
    combined_content = cv2.morphologyEx(
        combined_content.astype(np.uint8), 
        cv2.MORPH_OPEN,  # 작은 점 제거
        kernel, 
        iterations=1
    )
    
    return combined_content.astype(bool)


def find_content_bounds_ssim(
    image: Image.Image,
    region: Tuple[int, int, int, int],
    margin: int = 0,
    bg_threshold: int = 250,
    ssim_threshold: float = 0.95
) -> Tuple[int, int, int, int]:
    """SSIM 기반으로 컨텐츠 바운딩 박스 찾기"""
    x1, y1, x2, y2 = region
    
    # SSIM 기반 컨텐츠 감지
    content_mask = detect_content_with_ssim(image, region, bg_threshold, ssim_threshold)
    
    if not np.any(content_mask):
        return region
    
    # 컨텐츠가 있는 행/열 찾기
    rows_with_content = np.any(content_mask, axis=1)
    cols_with_content = np.any(content_mask, axis=0)
    
    if not np.any(rows_with_content) or not np.any(cols_with_content):
        return region
    
    # 첫 번째와 마지막 컨텐츠 위치
    top_idx = np.argmax(rows_with_content)
    bottom_idx = len(rows_with_content) - 1 - np.argmax(rows_with_content[::-1])
    left_idx = np.argmax(cols_with_content)
    right_idx = len(cols_with_content) - 1 - np.argmax(cols_with_content[::-1])
    
    # 전역 좌표로 변환
    content_x1 = x1 + left_idx
    content_y1 = y1 + top_idx
    content_x2 = x1 + right_idx + 1
    content_y2 = y1 + bottom_idx + 1
    
    # 마진 적용
    content_x1 = max(x1, content_x1 - margin)
    content_y1 = max(y1, content_y1 - margin)
    content_x2 = min(x2, content_x2 + margin)
    content_y2 = min(y2, content_y2 + margin)
    
    return (content_x1, content_y1, content_x2, content_y2)


def find_content_bounds_pixel_ssim(
    image: Image.Image,
    region: Tuple[int, int, int, int],
    bg_threshold: int = 250,
    margin: int = 5
) -> Tuple[int, int, int, int]:
    """SSIM 기반으로 실제 컨텐츠 영역을 찾기
    
    detect_content_with_ssim을 사용하여 더 정교하게 컨텐츠를 감지합니다.
    
    Args:
        image: 슬라이드 이미지
        region: 검사할 영역 (x1, y1, x2, y2)
        bg_threshold: 배경으로 간주할 밝기 임계값
        margin: 컨텐츠 경계에 추가할 여백
    
    Returns:
        (x1, y1, x2, y2): 실제 컨텐츠의 바운딩 박스
    """
    x1, y1, x2, y2 = region
    
    # SSIM 기반 컨텐츠 감지
    content_mask = detect_content_with_ssim(image, region, bg_threshold=bg_threshold)
    
    # 컨텐츠가 있는 행과 열 찾기
    content_rows = np.any(content_mask, axis=1)
    content_cols = np.any(content_mask, axis=0)
    
    if not np.any(content_rows) or not np.any(content_cols):
        # 컨텐츠가 없으면 원래 영역 반환
        return region
    
    # 컨텐츠의 바운딩 박스 찾기
    row_indices = np.where(content_rows)[0]
    col_indices = np.where(content_cols)[0]
    
    content_y1 = row_indices[0]
    content_y2 = row_indices[-1] + 1
    content_x1 = col_indices[0]
    content_x2 = col_indices[-1] + 1
    
    # 원래 좌표계로 변환하고 여백 추가
    orig_x1, orig_y1, orig_x2, orig_y2 = region
    
    result = (
        max(orig_x1, x1 + content_x1 - margin),
        max(orig_y1, y1 + content_y1 - margin),
        min(orig_x2, x1 + content_x2 + margin),
        min(orig_y2, y1 + content_y2 + margin)
    )
    
    return result


def adjust_box_to_avoid_content_ssim(
    image: Image.Image,
    box: Tuple[int, int, int, int],
    line_width: int = 4,
    max_adjust: int = 60,
    bg_threshold: int = 250,
    debug: bool = False,
    safety_edges: Tuple[str, ...] = ("top", "right", "bottom", "left"),
    bounds: Tuple[int, int, int, int] = None,  # (min_x, min_y, max_x, max_y) 조정 허용 범위
) -> Tuple[int, int, int, int]:
    """SSIM 기반으로 박스 테두리가 컨텐츠와 겹치지 않도록 조정
    
    기존 픽셀 기반 방식보다 더 정교하게 컨텐츠 경계를 감지합니다.
    """
    x1, y1, x2, y2 = box
    img_width, img_height = image.size
    original_box = (x1, y1, x2, y2)
    
    # 조정 경계 설정 (지정되지 않으면 이미지 전체 범위)
    if bounds is None:
        bounds = (0, 0, img_width, img_height)
    min_x, min_y, max_x, max_y = bounds

    def edge_content_ratio(side: str, bx1: int, by1: int, bx2: int, by2: int) -> float:
        """테두리 영역의 컨텐츠 비율 (0~1)"""
        if side == "top":
            region = (bx1, max(0, by1 - line_width), bx2, min(img_height, by1 + line_width))
        elif side == "bottom":
            region = (bx1, max(0, by2 - line_width), bx2, min(img_height, by2 + line_width))
        elif side == "left":
            region = (max(0, bx1 - line_width), by1, min(img_width, bx1 + line_width), by2)
        else:  # right
            region = (max(0, bx2 - line_width), by1, min(img_width, bx2 + line_width), by2)

        if region[2] <= region[0] or region[3] <= region[1]:
            return 0.0

        content_mask = detect_content_with_ssim(image, region, bg_threshold=bg_threshold)
        return float(np.mean(content_mask))

    # 최소 크기: 선 두께 기준으로 계산 (하드코딩 금지)
    min_width = max(line_width * 3, 8)
    min_height = max(line_width * 3, 8)
    content_threshold = 0.01  # 1% 이상일 때만 수축

    # 1) 테두리에 컨텐츠가 남아있으면 안쪽으로 수축 (최대 max_adjust 회)
    for _ in range(max_adjust):
        moved = False

        if edge_content_ratio("top", x1, y1, x2, y2) > content_threshold and (y2 - y1) > min_height:
            y1 = max(min_y, min(y1 + 1, y2 - min_height))
            moved = True

        if edge_content_ratio("bottom", x1, y1, x2, y2) > content_threshold and (y2 - y1) > min_height:
            y2 = min(max_y, max(y1 + min_height, y2 - 1))
            moved = True

        if edge_content_ratio("left", x1, y1, x2, y2) > content_threshold and (x2 - x1) > min_width:
            x1 = max(min_x, min(x1 + 1, x2 - min_width))
            moved = True

        if edge_content_ratio("right", x1, y1, x2, y2) > content_threshold and (x2 - x1) > min_width:
            x2 = min(max_x, max(x1 + min_width, x2 - 1))
            moved = True

        if not moved:
            break

    # 2) 배경이 깨끗한 방향으로는 다시 외측으로 확장 (라인이 배경 위에 오도록)
    # 상단/좌측을 우선 확장해 헤딩이 박스 안에 들어오게 함
    for _ in range(max_adjust):
        expanded = False

        if edge_content_ratio("top", x1, y1, x2, y2) == 0.0 and y1 > min_y:
            y1 = max(min_y, y1 - 1)
            expanded = True

        if edge_content_ratio("left", x1, y1, x2, y2) == 0.0 and x1 > min_x:
            x1 = max(min_x, x1 - 1)
            expanded = True

        if expanded:
            continue

        if edge_content_ratio("bottom", x1, y1, x2, y2) == 0.0 and y2 < max_y:
            y2 = min(max_y, y2 + 1)
            expanded = True

        if edge_content_ratio("right", x1, y1, x2, y2) == 0.0 and x2 < max_x:
            x2 = min(max_x, x2 + 1)
            expanded = True

        if not expanded:
            break

    # 3) 최종 안전 검사: RGB 기반으로 테두리에 어두운 픽셀이 남아있으면 한 픽셀씩 안쪽으로 수축
    # 특히 우측/상단을 우선 검사해 우측 박스 겹침을 방지
    if image.mode != "RGB":
        rgb_image = image.convert("RGB")
    else:
        rgb_image = image

    pixels = np.array(rgb_image)

    def edge_has_dark_pixels(side: str) -> bool:
        if side == "top":
            y_start = max(0, y1 - line_width)
            y_end = min(img_height, y1 + line_width)
            region = pixels[y_start:y_end, x1:x2]
        elif side == "bottom":
            y_start = max(0, y2 - line_width)
            y_end = min(img_height, y2 + line_width)
            region = pixels[y_start:y_end, x1:x2]
        elif side == "left":
            x_start = max(0, x1 - line_width)
            x_end = min(img_width, x1 + line_width)
            region = pixels[y1:y2, x_start:x_end]
        else:  # right
            x_start = max(0, x2 - line_width)
            x_end = min(img_width, x2 + line_width)
            region = pixels[y1:y2, x_start:x_end]

        if region.size == 0:
            return False

        return bool(np.any(region < bg_threshold))

    for _ in range(max_adjust):
        moved = False

        if "top" in safety_edges and edge_has_dark_pixels("top") and (y2 - y1) > min_height:
            y1 = max(min_y, min(y1 + 1, y2 - min_height))
            moved = True

        if "right" in safety_edges and edge_has_dark_pixels("right") and (x2 - x1) > min_width:
            x2 = min(max_x, max(x1 + min_width, x2 - 1))
            moved = True

        if "left" in safety_edges and edge_has_dark_pixels("left") and (x2 - x1) > min_width:
            x1 = max(min_x, min(x1 + 1, x2 - min_width))
            moved = True

        if "bottom" in safety_edges and edge_has_dark_pixels("bottom") and (y2 - y1) > min_height:
            y2 = min(max_y, max(y1 + min_height, y2 - 1))
            moved = True

        if not moved:
            break

    if debug:
        print(f"SSIM 박스 조정: {original_box} → {(x1, y1, x2, y2)}")

    return (x1, y1, x2, y2)


def draw_content_boxes_pixel(
    image: Image.Image,
    pptx_path: Path,
    page_number: int,
    slide_width: int,
    slide_height: int,
    box_color: str = "red",
    line_width: int = 4,
    padding: int = 15,
    use_ssim: bool = False
) -> Image.Image:
    """픽셀 기반으로 실제 컨텐츠 영역에 맞춰 박스 그리기
    
    이미지를 그리드 셀별로 나누고, 각 영역 내에서 픽셀 분석으로
    실제 컨텐츠 바운딩 박스를 찾습니다.
    
    Args:
        use_ssim: True이면 SSIM 기반 정교한 감지 사용
    """
    draw = ImageDraw.Draw(image)
    img_width, img_height = image.size
    
    # 파서로 레이아웃 정보 가져오기
    parser = PptxParser()
    result = parser.parse(pptx_path)
    
    page_layout = None
    for layout in result.page_layouts:
        if layout.page_number == page_number:
            page_layout = layout
            break
    
    if page_layout is None:
        return image
    
    cols = page_layout.cols
    rows = page_layout.rows
    
    # 마스터 요소 영역 계산 (제외할 영역)
    header_height = int(img_height * 0.08)  # 상단 영역 축소: 헤더 제외 최소화
    footer_height = int(img_height * 0.08)  # 하단 8% (로고, 페이지 번호)
    
    # 컨텐츠 영역 (마스터 요소 제외)
    content_top = header_height
    content_bottom = img_height - footer_height
    
    if cols == 1:
        # 1x1 레이아웃: 전체 영역에서 컨텐츠 기반 박스 계산
        region = (padding, content_top, img_width - padding, content_bottom)
        
        if use_ssim:
            bounds = find_content_bounds_pixel_ssim(image, region, margin=padding)
        else:
            bounds = find_content_bounds_pixel(image, region, margin=padding)
        
        # RGB 기반으로 정확한 콘텐츠 경계 찾기
        if image.mode != "RGB":
            rgb_image = image.convert("RGB")
        else:
            rgb_image = image
        pixels = np.array(rgb_image)
        
        content_threshold = 200
        margin_from_content = line_width * 4 + 5  # 박스 선 두께 고려하여 마진 확대
        
        # 4방향 콘텐츠 경계 탐색
        # 좌측: 왼쪽에서 오른쪽으로 스캔
        leftmost_x = img_width
        for x in range(padding, img_width - padding):
            col = pixels[content_top:content_bottom, x, :]
            brightness = np.mean(col, axis=1)
            if np.sum(brightness < content_threshold) > 3:
                leftmost_x = x
                break
        
        # 우측: 오른쪽에서 왼쪽으로 스캔
        rightmost_x = 0
        for x in range(img_width - padding - 1, padding, -1):
            col = pixels[content_top:content_bottom, x, :]
            brightness = np.mean(col, axis=1)
            if np.sum(brightness < content_threshold) > 3:
                rightmost_x = x
                break
        
        # 상단: 위에서 아래로 스캔 (헤더 제외 영역에서)
        topmost_y = img_height
        for y in range(content_top, content_bottom):
            row = pixels[y, padding:img_width-padding, :]
            brightness = np.mean(row, axis=1)
            if np.sum(brightness < content_threshold) > 5:
                topmost_y = y
                break
        
        # 하단: 아래에서 위로 스캔 (푸터 제외 영역에서)
        bottommost_y = 0
        for y in range(content_bottom - 1, content_top, -1):
            row = pixels[y, padding:img_width-padding, :]
            brightness = np.mean(row, axis=1)
            if np.sum(brightness < content_threshold) > 5:
                bottommost_y = y
                break
        
        # 콘텐츠 범위에 여백 추가
        # 헤더/푸터 영역은 무시: content_top/content_bottom 범위 내로 제한
        final_x1 = max(padding, leftmost_x - margin_from_content)
        final_x2 = min(img_width - padding, rightmost_x + margin_from_content)
        final_y1 = max(content_top, topmost_y - margin_from_content)  # 헤더 영역 무시
        final_y2 = min(content_bottom, bottommost_y + margin_from_content)  # 푸터 영역 무시
        
        adjusted_bounds = (final_x1, final_y1, final_x2, final_y2)
        
        # 디버그 출력
        print(f"  1x1 박스 경계: x={final_x1}-{final_x2}, y={final_y1}-{final_y2}")
        
        draw.rectangle(adjusted_bounds, outline=box_color, width=line_width)
        _draw_debug_corners(draw, adjusted_bounds, "C", font_size=12)  # C = Center (단일 박스)
    else:
        # 1x2 레이아웃: 좌우 영역을 분리하여 각각 컨텐츠 찾기
        mid_x = img_width // 2
        gap = 10  # 좌우 영역 사이 여백

        # 왼쪽 영역 - 먼저 계산
        left_region = (padding, content_top, mid_x - gap, content_bottom)
        if use_ssim:
            left_bounds = find_content_bounds_pixel_ssim(image, left_region, margin=padding)
        else:
            left_bounds = find_content_bounds_pixel(image, left_region, margin=padding)

        # 좌측 박스 조정
        if use_ssim:
            left_adjusted = adjust_box_to_avoid_content_ssim(
                image,
                left_bounds,
                line_width=line_width,
                safety_edges=(),  # 좌측 박스는 안전 수축 패스 비활성화
                bounds=(padding, content_top, mid_x, content_bottom),
            )
        else:
            left_adjusted = adjust_box_to_avoid_content(image, left_bounds, line_width=line_width)
        
        # 좌측 박스의 오른쪽 경계 기준으로 우측 박스 최소 시작점 계산
        lx1, ly1, lx2, ly2 = left_adjusted
        min_right_x = lx2 + line_width + gap  # 좌측 박스 끝 + 최소 간격
        
        # 오른쪽 영역 - 화면 중앙부터 검색하여 모든 우측 컨텐츠 포함
        right_region = (mid_x + gap, content_top, img_width - padding, content_bottom)
        if use_ssim:
            right_bounds = find_content_bounds_pixel_ssim(image, right_region, margin=padding)
        else:
            right_bounds = find_content_bounds_pixel(image, right_region, margin=padding)
        
        # 우측 박스 경계 확장 - 콘텐츠를 모두 포함하도록
        rx1, ry1, rx2, ry2 = right_bounds
        
        # RGB 기반으로 경계 확장
        if image.mode != "RGB":
            rgb_image = image.convert("RGB")
        else:
            rgb_image = image
        pixels = np.array(rgb_image)
        
        max_expand = 200  # 최대 확장 픽셀
        bg_threshold = 240  # 배경색 임계값
        content_threshold = 200  # 콘텐츠 픽셀 판단 임계값
        
        # === 좌측 경계 확장 (우측 영역 전체에서 가장 왼쪽 콘텐츠 찾기) ===
        # 전체 content 영역에서 콘텐츠를 탐색
        search_y1 = content_top
        search_y2 = content_bottom
        
        # 우측 영역 경계 (좌측 박스와 겹치지 않도록)
        search_x_start = lx2 + gap  # 좌측 박스 끝 + 간격
        
        # 가장 왼쪽에 있는 콘텐츠의 x 좌표 찾기 (전체 우측 영역 스캔)
        leftmost_content_x = img_width  # 초기값: 오른쪽 끝
        
        for check_x in range(search_x_start, img_width - padding):
            # 해당 x 좌표의 세로줄 검사 (전체 콘텐츠 영역)
            col = pixels[search_y1:search_y2, check_x, :]
            brightness = np.mean(col, axis=1)  # 각 픽셀의 밝기
            dark_pixels = np.sum(brightness < content_threshold)
            
            if dark_pixels > 3:  # 어두운 픽셀이 3개 이상이면 콘텐츠 있음
                leftmost_content_x = check_x
                break  # 가장 왼쪽 콘텐츠를 찾았으므로 중단
        
        # 콘텐츠 시작점보다 약간 왼쪽에 경계 설정 (여백)
        margin_from_content = line_width * 2 + 5
        final_rx1 = max(search_x_start, leftmost_content_x - margin_from_content)
        
        # === 하단 경계 확장 (우측 영역 전체에서 가장 아래 콘텐츠 찾기) ===
        # 우측 영역 전체에서 하단 콘텐츠 탐색
        bottommost_content_y = content_top  # 초기값: 위쪽
        
        for check_y in range(search_y1, search_y2):
            # 해당 y 좌표의 가로줄 검사 (우측 전체 영역)
            row = pixels[check_y, final_rx1:img_width - padding, :]
            brightness = np.mean(row, axis=1)
            dark_pixels = np.sum(brightness < content_threshold)
            
            if dark_pixels > 5:  # 어두운 픽셀이 5개 이상이면 콘텐츠 있음
                bottommost_content_y = check_y  # 계속 업데이트 (마지막 콘텐츠 위치)
        
        # 콘텐츠 끝보다 약간 아래에 경계 설정 (여백)
        final_ry2 = min(content_bottom, bottommost_content_y + margin_from_content)
        
        # 상단 경계도 콘텐츠에 맞춤
        topmost_content_y = content_bottom
        for check_y in range(search_y1, search_y2):
            row = pixels[check_y, final_rx1:img_width - padding, :]
            brightness = np.mean(row, axis=1)
            dark_pixels = np.sum(brightness < content_threshold)
            if dark_pixels > 5:
                topmost_content_y = check_y
                break
        final_ry1 = max(content_top, topmost_content_y - margin_from_content)
        
        right_bounds = (final_rx1, final_ry1, rx2, final_ry2)
        
        # 디버그: 경계 계산 결과 출력
        print(f"  우측 박스 경계: x={final_rx1}-{rx2}, y={final_ry1}-{final_ry2}")
        
        # 우측 박스: 직접 계산한 경계를 그대로 사용 (SSIM 조정 스킵)
        # 이미 콘텐츠 기반으로 정확히 계산했으므로 추가 조정 불필요
        right_adjusted = right_bounds

        # 박스끼리 살짝 겹치는 것은 허용 (enforce_gap 제거)

        # 박스 그리기
        if left_adjusted[2] > left_adjusted[0]:  # 유효한 영역인지 확인
            draw.rectangle(left_adjusted, outline=box_color, width=line_width)
            _draw_debug_corners(draw, left_adjusted, "L", font_size=12)
        if right_adjusted[2] > right_adjusted[0]:
            draw.rectangle(right_adjusted, outline=box_color, width=line_width)
            _draw_debug_corners(draw, right_adjusted, "R", font_size=12)
    
    return image


def _draw_debug_corners(
    draw: ImageDraw.Draw,
    box: Tuple[int, int, int, int],
    key: str,
    font_size: int = 12,
    text_color: str = "blue",
    bg_color: str = "white"
) -> None:
    """박스 모서리에 디버그 정보(키, 좌표) 표시
    
    Args:
        draw: ImageDraw 객체
        box: (x1, y1, x2, y2) 박스 좌표
        key: 박스 식별 키 (예: "L", "R")
        font_size: 폰트 크기
        text_color: 텍스트 색상
        bg_color: 배경 색상
    """
    x1, y1, x2, y2 = box
    
    # 코너 정보: (위치, 텍스트, 앵커)
    corners = [
        ((x1, y1), f"{key}-TL ({x1},{y1})", "lt"),  # 좌상단
        ((x2, y1), f"{key}-TR ({x2},{y1})", "rt"),  # 우상단
        ((x1, y2), f"{key}-BL ({x1},{y2})", "lb"),  # 좌하단
        ((x2, y2), f"{key}-BR ({x2},{y2})", "rb"),  # 우하단
    ]
    
    try:
        from PIL import ImageFont
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
    except:
        font = None
    
    for pos, text, anchor in corners:
        # 텍스트 크기 계산
        if font:
            bbox = draw.textbbox(pos, text, font=font, anchor=anchor)
        else:
            # 기본 폰트 사용시 대략적 크기
            text_width = len(text) * 6
            text_height = 10
            if anchor[0] == 'r':
                bbox = (pos[0] - text_width, pos[1], pos[0], pos[1] + text_height)
            else:
                bbox = (pos[0], pos[1], pos[0] + text_width, pos[1] + text_height)
            if anchor[1] == 'b':
                bbox = (bbox[0], bbox[1] - text_height, bbox[2], bbox[3] - text_height)
        
        # 배경 사각형 그리기 (텍스트 가독성)
        padding = 2
        draw.rectangle(
            (bbox[0] - padding, bbox[1] - padding, bbox[2] + padding, bbox[3] + padding),
            fill=bg_color
        )
        
        # 텍스트 그리기
        if font:
            draw.text(pos, text, fill=text_color, font=font, anchor=anchor)
        else:
            # anchor 위치 조정
            text_x, text_y = pos
            if anchor[0] == 'r':
                text_x -= len(text) * 6
            if anchor[1] == 'b':
                text_y -= 10
            draw.text((text_x, text_y), text, fill=text_color)


def find_vertical_gap(
    image: Image.Image,
    x_start: int,
    x_end: int,
    y_start: int,
    y_end: int,
    bg_threshold: int = 250
) -> int:
    """세로 방향으로 컨텐츠가 없는 X 위치 찾기
    
    두 컨텐츠 사이의 빈 공간(세로 줄)을 찾습니다.
    """
    import numpy as np
    
    # 영역 유효성 검사
    img_width, img_height = image.size
    x_start = max(0, x_start)
    x_end = min(img_width, x_end)
    y_start = max(0, y_start)
    y_end = min(img_height, y_end)
    
    if x_start >= x_end or y_start >= y_end:
        return None
    
    # 영역 크롭
    cropped = image.crop((x_start, y_start, x_end, y_end))
    
    if cropped.mode != 'RGB':
        cropped = cropped.convert('RGB')
    
    pixels = np.array(cropped)
    
    # 각 열(X 위치)에서 컨텐츠가 있는지 확인
    is_content = np.any(pixels < bg_threshold, axis=2)
    
    # 각 열의 컨텐츠 픽셀 수
    content_per_column = np.sum(is_content, axis=0)
    
    # 컨텐츠가 가장 적은 열 찾기 (완전히 비어있으면 0)
    min_content = np.min(content_per_column)
    
    # 컨텐츠가 거의 없는 열들 중 중앙에 가까운 것 선택
    threshold = min_content + (y_end - y_start) * 0.05  # 5% 이하 허용
    gap_columns = np.where(content_per_column <= threshold)[0]
    
    if len(gap_columns) == 0:
        return None
    
    # 중앙에 가까운 갭 선택
    center = len(content_per_column) // 2
    best_gap = gap_columns[np.argmin(np.abs(gap_columns - center))]
    
    return x_start + int(best_gap)
    
    return grouped


def calculate_bounding_box(boxes: List[ContentBox], padding: int = 0) -> Tuple[int, int, int, int]:
    """박스들의 통합 바운딩 박스 계산"""
    if not boxes:
        return (0, 0, 0, 0)
    
    min_left = min(box.left for box in boxes)
    min_top = min(box.top for box in boxes)
    max_right = max(box.left + box.width for box in boxes)
    max_bottom = max(box.top + box.height for box in boxes)
    
    return (
        min_left - padding,
        min_top - padding,
        max_right + padding,
        max_bottom + padding
    )


def convert_pptx_to_images(
    pptx_path: Path,
    output_dir: Path,
    reuse: bool = True,
    dpi: int = 150,
    remove_headers_footers: bool = False,
) -> dict:
    """PPT를 이미지로 변환 (LibreOffice 사용)

    기존 변환 결과가 있고 PPTX보다 최신이라면 변환을 건너뜁니다.

    Args:
        pptx_path: PPTX 파일 경로
        output_dir: 출력 디렉토리
        reuse: 기존 결과 재사용 여부
        dpi: 이미지 해상도
        remove_headers_footers: 헤더/푸터/페이지번호 제거 여부

    Returns:
        dict: {page_number: image_path} 형태의 딕셔너리
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    slide_count = len(Presentation(str(pptx_path)).slides)
    ppt_mtime = pptx_path.stat().st_mtime

    # 기존 결과 재사용 (헤더/푸터 제거 옵션이 변경되면 재생성)
    cache_marker = output_dir / ".no_headers_footers" if remove_headers_footers else output_dir / ".with_headers_footers"
    opposite_marker = output_dir / ".with_headers_footers" if remove_headers_footers else output_dir / ".no_headers_footers"
    
    if reuse and not opposite_marker.exists():
        existing_images = {}
        for img_path in output_dir.glob("slide-*.png"):
            try:
                page_num = int(img_path.stem.split("-")[-1])
            except ValueError:
                continue
            existing_images[page_num] = img_path

        have_all_pages = len(existing_images) >= slide_count and all(
            page in existing_images for page in range(1, slide_count + 1)
        )
        images_are_fresh = existing_images and min(
            img.stat().st_mtime for img in existing_images.values()
        ) >= ppt_mtime

        if have_all_pages and images_are_fresh:
            return existing_images
    
    # 기존 이미지 삭제 (옵션 변경 시)
    if opposite_marker.exists():
        for img_path in output_dir.glob("slide-*.png"):
            img_path.unlink()
        opposite_marker.unlink()

    try:
        # 헤더/푸터 제거 옵션 적용
        if remove_headers_footers:
            print("  헤더/푸터/페이지번호 제거 중...")
            clean_pptx_path = output_dir / f"{pptx_path.stem}_clean.pptx"
            remove_headers_footers_from_pptx(pptx_path, clean_pptx_path)
            source_pptx = clean_pptx_path
        else:
            source_pptx = pptx_path
        
        # LibreOffice로 PDF 변환
        pdf_dir = output_dir / "pdf_temp"
        pdf_dir.mkdir(parents=True, exist_ok=True)

        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(pdf_dir),
            str(source_pptx)
        ], check=True, capture_output=True)

        # PDF 파일 찾기
        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            return {}

        pdf_path = pdf_files[0]

        # pdftoppm으로 이미지 변환
        subprocess.run([
            "pdftoppm", "-png", "-r", str(dpi),
            str(pdf_path),
            str(output_dir / "slide")
        ], check=True, capture_output=True)

        # 생성된 이미지 매핑 (slide-1.png -> page 1)
        images = {}
        for img_path in output_dir.glob("slide-*.png"):
            try:
                page_num = int(img_path.stem.split("-")[-1])
            except ValueError:
                continue
            images[page_num] = img_path

        # 캐시 마커 생성
        cache_marker.touch()
        
        return images

    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"변환 오류: {e}")
        return {}


def create_blank_slide_image(width: int, height: int, bg_color: str = "white") -> Image.Image:
    """빈 슬라이드 이미지 생성"""
    return Image.new("RGB", (width, height), bg_color)


def draw_content_boxes(
    image: Image.Image,
    pptx_path: Path,
    page_number: int,
    slide_width: int,
    slide_height: int,
    box_color: str = "red",
    line_width: int = 4,
    padding: int = 15
) -> Image.Image:
    """실제 컨텐츠 영역에 맞춰 박스 그리기 (컨텐츠와 겹치지 않음)"""
    draw = ImageDraw.Draw(image)
    img_width, img_height = image.size
    
    # EMU를 이미지 픽셀로 변환하는 비율
    scale_x = img_width / slide_width
    scale_y = img_height / slide_height
    
    # 슬라이드의 컨텐츠 박스들 추출
    content_boxes, sw, sh = get_slide_content_boxes(pptx_path, page_number)
    
    if not content_boxes:
        return image
    
    # 파서로 레이아웃 정보 가져오기
    parser = PptxParser()
    result = parser.parse(pptx_path)
    
    page_layout = None
    for layout in result.page_layouts:
        if layout.page_number == page_number:
            page_layout = layout
            break
    
    if page_layout is None:
        return image
    
    cols = page_layout.cols
    rows = page_layout.rows
    
    # 박스들을 그리드 셀로 그룹화
    grouped = group_boxes_by_grid(content_boxes, slide_width, slide_height, rows, cols)
    
    # 각 그룹의 바운딩 박스에 테두리 그리기
    for cell_key, cell_boxes in grouped.items():
        # 바운딩 박스 계산 (EMU 단위)
        bbox = calculate_bounding_box(cell_boxes, padding=int(padding / scale_x))
        
        # 이미지 좌표로 변환
        x1 = int(bbox[0] * scale_x)
        y1 = int(bbox[1] * scale_y)
        x2 = int(bbox[2] * scale_x)
        y2 = int(bbox[3] * scale_y)
        
        # 이미지 경계 내로 제한
        x1 = max(5, x1)
        y1 = max(5, y1)
        x2 = min(img_width - 5, x2)
        y2 = min(img_height - 5, y2)
        
        # 박스 그리기
        draw.rectangle([x1, y1, x2, y2], outline=box_color, width=line_width)
    
    return image


def draw_grid_boxes(
    image: Image.Image,
    grid_cells: list,
    slide_width: int,
    slide_height: int,
    box_color: str = "red",
    line_width: int = 3,
    padding: int = 5
) -> Image.Image:
    """그리드 셀에 박스 그리기 (기존 방식 - 셀 전체 영역)"""
    draw = ImageDraw.Draw(image)
    img_width, img_height = image.size
    
    # EMU를 이미지 픽셀로 변환하는 비율
    scale_x = img_width / slide_width
    scale_y = img_height / slide_height
    
    for cell in grid_cells:
        # 셀 좌표를 이미지 좌표로 변환
        x1 = int(cell.left * scale_x) + padding
        y1 = int(cell.top * scale_y) + padding
        x2 = int((cell.left + cell.width) * scale_x) - padding
        y2 = int((cell.top + cell.height) * scale_y) - padding
        
        # 박스 그리기
        draw.rectangle([x1, y1, x2, y2], outline=box_color, width=line_width)
    
    return image


def visualize_page_layout(
    pptx_path: Path,
    page_number: int,
    output_path: Path = None,
    image_width: int = 960,
    image_height: int = 540,
    box_color: str = "red",
    show: bool = True,
    use_slide_image: bool = True,
    fit_content: bool = True,  # 컨텐츠에 맞춰 박스 그리기
    single_box: bool = False,  # True: 전체 컨텐츠를 하나의 박스로, False: 그리드 셀별 박스
    pixel_based: bool = True,  # True: 픽셀 기반 컨텐츠 감지 (더 정확함)
    use_ssim: bool = False,    # True: SSIM 기반 정교한 감지
    remove_headers_footers: bool = False  # 헤더/푸터/페이지번호 제거
):
    """특정 페이지의 레이아웃을 시각화
    
    Args:
        single_box: True이면 마스터 요소(로고, Confidential 등)를 제외하고
                   메인 컨텐츠 전체를 하나의 박스로 그림
        pixel_based: True이면 실제 렌더링된 픽셀을 분석하여 컨텐츠 영역 감지
        use_ssim: True이면 SSIM 기반으로 더 정교한 컨텐츠 감지
        remove_headers_footers: True이면 PPTX에서 헤더/푸터/페이지번호 제거 후 변환
    """
    
    # PPT 파싱
    parser = PptxParser()
    result = parser.parse(pptx_path)
    
    # 해당 페이지의 레이아웃 찾기
    page_layout = None
    for layout in result.page_layouts:
        if layout.page_number == page_number:
            page_layout = layout
            break
    
    if page_layout is None:
        print(f"페이지 {page_number}를 찾을 수 없습니다.")
        return
    
    print(f"페이지 {page_number}: {page_layout.rows}x{page_layout.cols} 레이아웃")
    print(f"그리드 셀 수: {len(page_layout.grid_cells)}")
    
    # 슬라이드 이미지 준비
    image = None
    if use_slide_image:
        cache_context = (
            tempfile.TemporaryDirectory() if output_path is None
            else nullcontext(str(output_path.parent / "slide_cache"))
        )

        with cache_context as temp_dir:
            cache_dir = Path(temp_dir)
            cache_dir.mkdir(parents=True, exist_ok=True)
            slide_images = convert_pptx_to_images(
                pptx_path, cache_dir, reuse=True,
                remove_headers_footers=remove_headers_footers
            )

            if page_number in slide_images:
                # 임시 파일이 삭제되기 전에 메모리로 복사
                with Image.open(slide_images[page_number]) as img:
                    image = img.copy()
                image_width, image_height = image.size
                print(f"슬라이드 이미지 로드: {image_width}x{image_height}")
    
    if image is None:
        # 빈 슬라이드 이미지 생성
        image = create_blank_slide_image(image_width, image_height, "white")
    
    # 박스 그리기
    if single_box:
        # 마스터 요소 제외, 전체 컨텐츠를 하나의 박스로
        image = draw_single_content_box(
            image,
            pptx_path,
            page_number,
            box_color=box_color,
            line_width=3,
            padding=10
        )
    elif fit_content and pixel_based:
        # 픽셀 기반으로 실제 컨텐츠 영역에 맞춰 박스 (가장 정확함)
        image = draw_content_boxes_pixel(
            image,
            pptx_path,
            page_number,
            page_layout.slide_width,
            page_layout.slide_height,
            box_color=box_color,
            line_width=4,
            padding=15,
            use_ssim=use_ssim
        )
    elif fit_content:
        # shape 기반으로 컨텐츠에 맞춰 박스
        image = draw_content_boxes(
            image,
            pptx_path,
            page_number,
            page_layout.slide_width,
            page_layout.slide_height,
            box_color=box_color,
            line_width=4,
            padding=20
        )
    else:
        # 그리드 셀 전체에 박스
        image = draw_grid_boxes(
            image,
            page_layout.grid_cells,
            page_layout.slide_width,
            page_layout.slide_height,
            box_color=box_color,
            line_width=4,
            padding=8
        )
    
    # 레이아웃 정보 텍스트 추가
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except:
        font = ImageFont.load_default()
        small_font = font
    
    # 반투명 배경의 정보 박스
    info_text = f"Page {page_number}: {page_layout.rows}x{page_layout.cols}"
    bbox = draw.textbbox((0, 0), info_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    
    # 배경 사각형
    draw.rectangle([5, 5, text_width + 20, text_height + 15], fill=(255, 255, 255, 200))
    draw.text((10, 8), info_text, fill="black", font=font)
    
    # 저장 또는 표시
    if output_path:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        image.save(output_path)
        print(f"저장됨: {output_path}")
    
    if show:
        image.show()
    
    return image


def visualize_all_pages(
    pptx_path: Path,
    output_dir: Path,
    pages: list = None
):
    """모든 페이지 (또는 지정된 페이지들)의 레이아웃 시각화"""
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # PPT 파싱
    parser = PptxParser()
    result = parser.parse(pptx_path)
    
    for layout in result.page_layouts:
        if pages and layout.page_number not in pages:
            continue
        
        output_path = output_dir / f"page_{layout.page_number:02d}_layout.png"
        visualize_page_layout(
            pptx_path,
            layout.page_number,
            output_path=output_path,
            show=False
        )


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="PPT 레이아웃 시각화")
    parser.add_argument("pptx_path", type=Path, help="PPTX 파일 경로")
    parser.add_argument("-p", "--page", type=int, help="시각화할 페이지 번호")
    parser.add_argument("-o", "--output", type=Path, help="출력 파일/디렉토리 경로")
    parser.add_argument("--all", action="store_true", help="모든 페이지 시각화")
    parser.add_argument("--no-show", action="store_true", help="이미지 표시 안 함")
    parser.add_argument("--use-ssim", action="store_true", help="SSIM 기반 정교한 컨텐츠 감지 사용")
    parser.add_argument("--no-headers-footers", action="store_true", 
                        help="헤더/푸터/페이지번호 제거 후 변환")
    
    args = parser.parse_args()
    
    if args.all:
        output_dir = args.output or Path("layout_images")
        visualize_all_pages(args.pptx_path, output_dir)
    elif args.page:
        visualize_page_layout(
            args.pptx_path,
            args.page,
            output_path=args.output,
            show=not args.no_show,
            use_ssim=args.use_ssim,
            remove_headers_footers=args.no_headers_footers
        )
    else:
        # 기본: 6페이지 시각화
        visualize_page_layout(
            args.pptx_path,
            6,
            output_path=args.output,
            show=not args.no_show,
            remove_headers_footers=args.no_headers_footers
        )
