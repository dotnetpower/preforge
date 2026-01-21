"""테스트용 Office 문서 생성 스크립트"""
from pathlib import Path
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt


def create_test_docx():
    """테스트용 DOCX 파일 생성"""
    doc = Document()
    
    # 문서 제목
    title = doc.add_heading('Preforge 테스트 문서', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 메타데이터 설정
    doc.core_properties.title = 'Preforge 테스트 문서'
    doc.core_properties.author = 'Preforge Team'
    doc.core_properties.subject = '파서 테스트'
    doc.core_properties.keywords = '테스트, 파싱, 검증'
    
    # 1. 개요
    doc.add_heading('1. 개요', 1)
    doc.add_paragraph(
        'Preforge는 다양한 문서 형식을 파싱하고 분석하는 Python 라이브러리입니다. '
        '이 문서는 DOCX 파서의 기능을 테스트하기 위한 샘플 문서입니다.'
    )
    
    # 2. 주요 기능
    doc.add_heading('2. 주요 기능', 1)
    
    doc.add_heading('2.1 문서 파싱', 2)
    doc.add_paragraph(
        'Preforge는 다음과 같은 문서 형식을 지원합니다:'
    )
    
    # 글머리 기호 목록
    doc.add_paragraph('Word 문서 (.docx)', style='List Bullet')
    doc.add_paragraph('PowerPoint 프레젠테이션 (.pptx)', style='List Bullet')
    doc.add_paragraph('Excel 스프레드시트 (.xlsx)', style='List Bullet')
    doc.add_paragraph('PDF 문서 (.pdf)', style='List Bullet')
    doc.add_paragraph('HTML 파일 (.html)', style='List Bullet')
    
    doc.add_heading('2.2 데이터 추출', 2)
    doc.add_paragraph(
        '각 문서에서 다음 정보를 추출할 수 있습니다:'
    )
    
    # 번호 목록
    doc.add_paragraph('텍스트 및 제목', style='List Number')
    doc.add_paragraph('테이블 데이터', style='List Number')
    doc.add_paragraph('이미지 및 메타데이터', style='List Number')
    doc.add_paragraph('문서 속성', style='List Number')
    
    # 3. 테이블 예시
    doc.add_heading('3. 테이블 예시', 1)
    doc.add_paragraph('다음은 지원되는 문서 형식에 대한 정보입니다:')
    
    # 테이블 생성
    table = doc.add_table(rows=6, cols=4)
    table.style = 'Light Grid Accent 1'
    
    # 헤더 행
    headers = ['문서 형식', '확장자', '라이브러리', '상태']
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.paragraphs[0].runs[0].font.bold = True
    
    # 데이터 행
    data = [
        ['Word', '.docx', 'python-docx', '완료'],
        ['PowerPoint', '.pptx', 'python-pptx', '완료'],
        ['Excel', '.xlsx', 'openpyxl', '계획'],
        ['PDF', '.pdf', 'pdfplumber', '완료'],
        ['HTML', '.html', 'beautifulsoup4', '완료'],
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_data in enumerate(row_data):
            table.rows[i].cells[j].text = cell_data
    
    # 4. 코드 예시
    doc.add_heading('4. 사용 예시', 1)
    doc.add_paragraph('Python 코드 예시:')
    
    code = doc.add_paragraph(
        'from preforge.parsers import DocxParser\n\n'
        'parser = DocxParser()\n'
        'doc = parser.parse("sample.docx")\n\n'
        'print(f"제목: {doc.metadata.title}")\n'
        'print(f"텍스트 블록 수: {len(doc.text_contents)}")\n'
        'print(f"테이블 수: {len(doc.tables)}")',
        style='No Spacing'
    )
    code.runs[0].font.name = 'Courier New'
    code.runs[0].font.size = Pt(9)
    
    # 5. 결론
    doc.add_heading('5. 결론', 1)
    doc.add_paragraph(
        'Preforge는 다양한 문서 형식을 통합적으로 처리할 수 있는 강력한 도구입니다. '
        '이 테스트 문서를 통해 DOCX 파서의 텍스트, 제목, 테이블 추출 기능을 검증할 수 있습니다.'
    )
    
    # 파일 저장
    output_path = Path(__file__).parent.parent / 'private' / 'test_document.docx'
    doc.save(output_path)
    print(f"✅ 테스트 DOCX 파일 생성: {output_path}")
    return output_path


def create_test_pptx():
    """테스트용 PPTX 파일 생성"""
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    # 슬라이드 1: 제목 슬라이드
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Preforge 테스트 프레젠테이션"
    subtitle.text = "문서 파싱 라이브러리 소개\nPreforge Team"
    
    # 슬라이드 2: 개요
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "1. 개요"
    
    tf = content2.text_frame
    tf.text = "Preforge는 다양한 문서 형식을 파싱하는 Python 라이브러리입니다."
    
    p = tf.add_paragraph()
    p.text = "주요 특징:"
    p.level = 0
    
    for feature in ["다중 형식 지원", "구조화된 데이터 추출", "메타데이터 보존", "확장 가능한 아키텍처"]:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
    
    # 슬라이드 3: 지원 형식
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.title
    title3.text = "2. 지원하는 문서 형식"
    
    # 테이블 추가
    rows, cols = 6, 3
    left = PptxInches(1.5)
    top = PptxInches(2)
    width = PptxInches(7)
    height = PptxInches(4)
    
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 헤더
    table.cell(0, 0).text = "형식"
    table.cell(0, 1).text = "확장자"
    table.cell(0, 2).text = "상태"
    
    # 데이터
    formats = [
        ("Word", ".docx", "완료"),
        ("PowerPoint", ".pptx", "완료"),
        ("Excel", ".xlsx", "계획"),
        ("PDF", ".pdf", "완료"),
        ("HTML", ".html", "완료"),
    ]
    
    for i, (fmt, ext, status) in enumerate(formats, start=1):
        table.cell(i, 0).text = fmt
        table.cell(i, 1).text = ext
        table.cell(i, 2).text = status
    
    # 슬라이드 4: 아키텍처
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "3. 시스템 아키텍처"
    
    tf4 = content4.text_frame
    tf4.text = "계층적 구조"
    
    layers = [
        "Parser Layer: 문서 형식별 파싱",
        "Extractor Layer: 데이터 추출",
        "Model Layer: NER 및 AI 분석",
        "Knowledge Layer: 그래프 기반 검색"
    ]
    
    for layer in layers:
        p = tf4.add_paragraph()
        p.text = layer
        p.level = 1
    
    # 슬라이드 5: 결론
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "4. 결론"
    
    tf5 = content5.text_frame
    tf5.text = "Preforge는 다음을 제공합니다:"
    
    conclusions = [
        "통합된 문서 처리 인터페이스",
        "정확한 데이터 추출",
        "확장 가능한 아키텍처",
        "실무 적용 가능한 성능"
    ]
    
    for conclusion in conclusions:
        p = tf5.add_paragraph()
        p.text = conclusion
        p.level = 1
    
    # 파일 저장
    output_path = Path(__file__).parent.parent / 'private' / 'test_presentation.pptx'
    prs.save(output_path)
    print(f"✅ 테스트 PPTX 파일 생성: {output_path}")
    return output_path


if __name__ == "__main__":
    print("테스트용 Office 문서 생성 중...\n")
    
    docx_path = create_test_docx()
    pptx_path = create_test_pptx()
    
    print(f"\n생성 완료!")
    print(f"- DOCX: {docx_path.name}")
    print(f"- PPTX: {pptx_path.name}")
