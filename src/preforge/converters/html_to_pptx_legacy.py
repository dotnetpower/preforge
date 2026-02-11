"""
Converter for HTML(.html) documents to PowerPoint(.pptx)

Analyzes HTML structure and generates slides by section.
Supports various elements including tables, text, lists, images, etc.
"""
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
from io import BytesIO
import logging
import re
import base64
from urllib.parse import urlparse

from bs4 import BeautifulSoup, Tag, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

logger = logging.getLogger(__name__)


class HtmlToPptxConverter:
    """Converter to transform HTML to PowerPoint"""
    
    def __init__(self):
        """Initialize converter"""
        self.prs = None
        self.current_slide = None
        self.html_path = None  # HTML path for screenshots
        
        # Table split settings
        self.max_rows_per_slide = 8  # Maximum rows per slide (excluding header) - improved readability
        
        # Color definitions (extracted from HTML CSS)
        self.colors = {
            'primary_red': RGBColor(220, 38, 38),  # #dc2626
            'primary_red_light': RGBColor(254, 242, 242),  # #fef2f2
            'primary_red_dark': RGBColor(153, 27, 27),  # #991b1b
            'gray_50': RGBColor(249, 250, 251),
            'gray_100': RGBColor(243, 244, 246),
            'gray_200': RGBColor(229, 231, 235),
            'gray_600': RGBColor(75, 85, 99),
            'gray_800': RGBColor(31, 41, 55),
            'gray_900': RGBColor(17, 24, 39),
            'white': RGBColor(255, 255, 255),
            'black': RGBColor(0, 0, 0),
        }
        
        # Slide size settings (16:9 ratio)
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        
        # Margin settings (adjusted narrower)
        self.margin_left = Inches(0.3)
        self.margin_right = Inches(0.3)
        self.margin_top = Inches(0.5)
        self.margin_bottom = Inches(0.3)
        
        self.content_width = self.slide_width - self.margin_left - self.margin_right
        self.content_height = self.slide_height - self.margin_top - self.margin_bottom
    
    def _extract_cell_styles(self, cell_elem: Tag) -> Dict[str, Any]:
        """Extract styles (Bold, Color, etc.) from HTML cell"""
        styles = {
            'bold': False,
            'color': None,
            'background': None,
            'link': None,
        }
        
        # Check cell's own style
        style_attr = cell_elem.get('style', '')
        
        # Extract color
        color_match = re.search(r'color:\s*(#[a-fA-F0-9]{6}|#[a-fA-F0-9]{3}|rgb\([^)]+\))', style_attr)
        if color_match:
            color_str = color_match.group(1)
            styles['color'] = self._parse_color(color_str)
        
        # Extract background-color
        bg_match = re.search(r'background(?:-color)?:\s*(#[a-fA-F0-9]{6}|#[a-fA-F0-9]{3}|rgb\([^)]+\))', style_attr)
        if bg_match:
            styles['background'] = self._parse_color(bg_match.group(1))
        
        # Check font-weight: bold or 700
        if 'font-weight' in style_attr:
            weight_match = re.search(r'font-weight:\s*(\w+)', style_attr)
            if weight_match:
                weight = weight_match.group(1)
                if weight in ('bold', '700', '800', '900'):
                    styles['bold'] = True
        
        # Check inner bold tags (b, strong)
        if cell_elem.find(['b', 'strong']):
            styles['bold'] = True
        
        # Check inner color style (span etc.)
        colored_elem = cell_elem.find(style=True)
        if colored_elem and not styles['color']:
            inner_style = colored_elem.get('style', '')
            inner_color = re.search(r'color:\s*(#[a-fA-F0-9]{6}|#[a-fA-F0-9]{3}|rgb\([^)]+\))', inner_style)
            if inner_color:
                styles['color'] = self._parse_color(inner_color.group(1))
        
        # Check link
        link = cell_elem.find('a')
        if link:
            styles['link'] = link.get('href', '')
        
        return styles
    
    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """Convert color string to RGBColor"""
        if not color_str:
            return None
        
        try:
            # hex color (#rrggbb or #rgb)
            if color_str.startswith('#'):
                hex_color = color_str[1:]
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
            
            # rgb(r, g, b)
            if color_str.startswith('rgb'):
                match = re.search(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
                if match:
                    r, g, b = int(match.group(1)), int(match.group(2)), int(match.group(3))
                    return RGBColor(r, g, b)
        except:
            pass
        
        return None
    
    def convert(self, html_path: Path, output_path: Path) -> None:
        """
        Convert HTML file to PPTX
        
        Args:
            html_path: Input HTML file path
            output_path: Output PPTX file path
        """
        logger.info(f"HTML -> PPTX conversion started: {html_path} -> {output_path}")
        
        # Save HTML file path (for screenshots)
        self.html_path = Path(html_path).absolute()
        
        # Read HTML file
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'lxml')
        
        # Initialize presentation
        self.prs = Presentation()
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height
        
        # Create title slide
        self._create_title_slide(soup)
        
        # Create Analysis Summary section
        self._create_analysis_summary_slides(soup)
        
        # Process main content
        self._process_main_content(soup)
        
        # Save PPTX
        self.prs.save(str(output_path))
        logger.info(f"Conversion complete: {output_path} (total {len(self.prs.slides)} slides)")
    
    def _create_title_slide(self, soup: BeautifulSoup) -> None:
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Empty layout
        
        # Extract title
        title_elem = soup.find('div', class_='header-title')
        subtitle_elem = soup.find('div', class_='header-subtitle')
        
        title_text = title_elem.get_text(strip=True) if title_elem else "GeneSeq Vista AI Agent"
        subtitle_text = subtitle_elem.get_text(strip=True) if subtitle_elem else ""
        
        # Background rectangle
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.slide_width, self.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['gray_50']
        background.line.fill.background()
        
        # Top red box
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(2),
            self.slide_width, Inches(1.5)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['primary_red']
        header_box.line.fill.background()
        
        # Title text
        title_frame = header_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']
        
        # Subtitle text
        if subtitle_text:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4),
                self.slide_width - Inches(2), Inches(1.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = self.colors['gray_800']
            subtitle_frame.word_wrap = True
    
    def _create_analysis_summary_slides(self, soup: BeautifulSoup) -> None:
        """Create Analysis Summary section slides"""
        analysis_div = soup.find('div', class_='analysis-summary')
        if not analysis_div:
            return
        
        # Full summary slide
        summary_section = analysis_div.find('div', class_='summary-section')
        if summary_section:
            self._create_summary_slide(summary_section)
        
        # Target Gene Ranking slide
        summary_sections = analysis_div.find_all('div', class_='summary-section')
        if len(summary_sections) > 1:
            self._create_ranking_slide(summary_sections[1])
    
    def _create_summary_slide(self, section: Tag) -> None:
        """Create full summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Extract header
        header = section.find('div', class_='section-header')
        header_text = header.get_text(strip=True) if header else "Analysis Summary"
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin_left, self.margin_top - Inches(0.2),
            self.content_width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = header_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Extract table - display key-value format in card style
        table_elem = section.find('table')
        if table_elem:
            # Check if key-value format (no thead, 2 columns)
            thead = table_elem.find('thead')
            tbody = table_elem.find('tbody')
            
            if not thead and tbody:
                rows = tbody.find_all('tr')
                first_row = rows[0] if rows else None
                if first_row:
                    cells = first_row.find_all(['th', 'td'])
                    if len(cells) == 2 and len(rows) <= 5:
                        # Display in card style
                        self._add_key_value_cards(
                            slide, table_elem,
                            self.margin_left, self.margin_top + Inches(0.6),
                            self.content_width, Inches(5)
                        )
                        return
            
            # Regular table
            self._add_table_to_slide(
                slide, table_elem,
                self.margin_left, self.margin_top + Inches(0.6),
                self.content_width, Inches(5)
            )
    
    def _create_ranking_slide(self, section: Tag) -> None:
        """Create Target Gene Ranking slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Extract header
        header = section.find('div', class_='section-header')
        header_text = header.get_text(strip=True) if header else "Target Gene Ranking"
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin_left, self.margin_top - Inches(0.2),
            self.content_width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = header_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Extract and add table
        table_elem = section.find('table')
        if table_elem:
            self._add_table_to_slide(
                slide, table_elem,
                self.margin_left, self.margin_top + Inches(0.6),
                self.content_width, Inches(5.5)
            )
    
    def _process_main_content(self, soup: BeautifulSoup) -> None:
        """Process all main content"""
        # Find all sections in content-container
        content_container = soup.find('div', class_='content-container')
        if not content_container:
            return
        
        # Gene title (main title)
        gene_title_elem = content_container.find('h1', class_='gene-title')
        main_gene_title = gene_title_elem.get_text(strip=True) if gene_title_elem else "Gene Analysis"
        
        # Process all gene-sections
        gene_sections = content_container.find_all('div', class_='gene-section')
        seq_viewer_index = 0  # SeqViewerApp dedicated index
        
        for idx, gene_section in enumerate(gene_sections, 1):
            # Extract section title
            subsection_title = gene_section.find('h2', class_='subsection-title')
            section_title = subsection_title.get_text(strip=True) if subsection_title else f"Section {idx}"
            
            # Capture SeqViewerApp (using Playwright) - process all elements with find_all
            seq_viewers = gene_section.find_all('div', class_='SeqViewerApp', recursive=False)
            for sv_idx, seq_viewer in enumerate(seq_viewers):
                if self.html_path:
                    viewer_title = section_title + f" - Sequence Viewer"
                    if len(seq_viewers) > 1:
                        viewer_title += f" ({sv_idx + 1})"
                    self._capture_element_screenshot(
                        '.SeqViewerApp', 
                        viewer_title,
                        seq_viewer_index
                    )
                    seq_viewer_index += 1
            
            # Process if image exists
            image_placeholder = gene_section.find('div', class_='image-placeholder')
            if image_placeholder:
                img = image_placeholder.find('img')
                if img and img.get('src', '').startswith('data:image'):
                    self._create_image_slide(img, section_title, main_gene_title)
            
            # If tables exist - combine small tables into one slide
            tables = gene_section.find_all('table', class_='data-table')
            if tables:
                # Collect table information
                table_infos = []
                for table in tables:
                    row_count = len(table.find_all('tr'))
                    table_infos.append({'table': table, 'rows': row_count})
                
                # Group small tables (combined 8 rows or less)
                i = 0
                while i < len(table_infos):
                    current_group = [table_infos[i]]
                    total_rows = table_infos[i]['rows']
                    
                    # Check if next table fits (dynamic margin calculation)
                    # Calculate available height in slide
                    title_space = Inches(0.85)  # Title space (main_title + section_title)
                    table_gap = Inches(0.3)  # Gap between tables
                    available_height = self.slide_height - self.margin_top - self.margin_bottom - title_space
                    row_height = Inches(0.28)  # Height per row (considering font size + margin)
                    
                    # Calculate current table height
                    current_height = row_height * total_rows
                    
                    while i + 1 < len(table_infos):
                        next_rows = table_infos[i + 1]['rows']
                        next_table_height = row_height * next_rows + table_gap
                        
                        # Check if next table fits in remaining space
                        remaining_space = available_height - current_height
                        
                        if next_table_height <= remaining_space:
                            # Can add next table
                            i += 1
                            current_group.append(table_infos[i])
                            total_rows += next_rows
                            current_height += next_table_height
                        else:
                            break
                    
                    # Process group
                    if len(current_group) == 1:
                        # Single table
                        table_title = f"{section_title}"
                        self._create_data_table_slide(current_group[0]['table'], table_title, main_gene_title)
                    else:
                        # Combine multiple tables into one slide
                        self._create_combined_table_slide(
                            [info['table'] for info in current_group], 
                            section_title, 
                            main_gene_title
                        )
                    
                    i += 1
            
            # Process subsections (h3)
            h3_sections = gene_section.find_all('h3')
            for h3 in h3_sections:
                h3_title = h3.get_text(strip=True)
                # Find table after h3
                next_table = h3.find_next('table')
                if next_table and next_table.parent == gene_section:
                    self._create_data_table_slide(next_table, h3_title, main_gene_title)
        
        # Also process sections with only subsection-title
        all_subsections = content_container.find_all('h2', class_='subsection-title')
        for subsection in all_subsections:
            section_title = subsection.get_text(strip=True)
            
            # Process Evidence table
            next_elem = subsection.find_next_sibling()
            while next_elem:
                if next_elem.name == 'h2' or next_elem.get('class') and 'gene-section' in next_elem.get('class', []):
                    break
                
                if next_elem.name == 'div' and 'evidence-table' in next_elem.get('class', []):
                    self._create_evidence_table_slide(next_elem, section_title)
                
                next_elem = next_elem.find_next_sibling()
    
    def _create_data_table_slide(self, table_elem: Tag, section_title: str, main_title: str = "") -> None:
        """Create data table slide (improved readability, auto-split support)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Main title (small)
        if main_title:
            main_box = slide.shapes.add_textbox(
                self.margin_left, Inches(0.1),
                self.content_width, Inches(0.25)
            )
            main_frame = main_box.text_frame
            main_frame.text = main_title
            main_para = main_frame.paragraphs[0]
            main_para.font.size = Pt(12)
            main_para.font.color.rgb = self.colors['gray_600']
        
        # Section title
        title_top = Inches(0.35) if main_title else self.margin_top - Inches(0.1)
        title_box = slide.shapes.add_textbox(
            self.margin_left, title_top,
            self.content_width, Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Add table (larger area)
        table_top = title_top + Inches(0.5)
        table_height = self.slide_height - table_top - self.margin_bottom
        
        # Add table (may return multiple slides)
        created_slides = self._add_improved_table(
            slide, table_elem,
            self.margin_left, table_top,
            self.content_width, table_height
        )
        
        # Add title to additional slides
        if created_slides and len(created_slides) > 1:
            for idx, extra_slide in enumerate(created_slides[1:], 2):
                # Also add title to additional slides
                if extra_slide:
                    extra_title_box = extra_slide.shapes.add_textbox(
                        self.margin_left, Inches(0.1),
                        self.content_width, Inches(0.4)
                    )
                    extra_title_frame = extra_title_box.text_frame
                    extra_title_frame.text = f"{section_title} (continued {idx})"
                    extra_title_para = extra_title_frame.paragraphs[0]
                    extra_title_para.font.size = Pt(18)
                    extra_title_para.font.bold = True
                    extra_title_para.font.color.rgb = self.colors['primary_red']
    
    def _create_combined_table_slide(self, tables: List[Tag], section_title: str, main_title: str = "") -> None:
        """Combine multiple small tables into one slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Main title (small)
        if main_title:
            main_box = slide.shapes.add_textbox(
                self.margin_left, Inches(0.1),
                self.content_width, Inches(0.25)
            )
            main_frame = main_box.text_frame
            main_frame.text = main_title
            main_para = main_frame.paragraphs[0]
            main_para.font.size = Pt(12)
            main_para.font.color.rgb = self.colors['gray_600']
        
        # Section title
        title_top = Inches(0.35) if main_title else self.margin_top - Inches(0.1)
        title_box = slide.shapes.add_textbox(
            self.margin_left, title_top,
            self.content_width, Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Add tables sequentially
        current_top = title_top + Inches(0.5)
        available_height = self.slide_height - current_top - self.margin_bottom
        
        for table_idx, table_elem in enumerate(tables):
            # Calculate row count for each table
            rows = table_elem.find_all('tr')
            row_count = len(rows)
            
            # Estimate table height (0.25 inches per row)
            table_height = min(Inches(0.25) * row_count, available_height * 0.4)
            
            if table_idx > 0:
                # Gap between tables
                current_top += Inches(0.2)
            
            # Add table
            self._add_improved_table(
                slide, table_elem,
                self.margin_left, current_top,
                self.content_width, table_height
            )
            
            # Calculate next table position
            current_top += table_height + Inches(0.1)

    def _create_gene_overview_slide(self, gene_section: Tag, gene_title: str) -> None:
        """Create Gene overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.margin_left, self.margin_top - Inches(0.2),
            self.content_width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = gene_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Find Background section
        background_div = gene_section.find('div', class_='background-text')
        if background_div:
            y_position = self.margin_top + Inches(0.7)
            
            # "Background" subtitle
            subtitle_box = slide.shapes.add_textbox(
                self.margin_left, y_position,
                self.content_width, Inches(0.3)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = "Background"
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.bold = True
            subtitle_para.font.color.rgb = self.colors['gray_800']
            
            # Background text
            background_text = background_div.get_text(strip=True)
            text_box = slide.shapes.add_textbox(
                self.margin_left, y_position + Inches(0.4),
                self.content_width, Inches(4)
            )
            text_frame = text_box.text_frame
            text_frame.text = background_text
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = self.colors['gray_800']
                paragraph.line_spacing = 1.5
    
    def _create_image_slide(self, img_tag: Tag, section_title: str, main_title: str = "") -> None:
        """Create slide with image"""
        import base64
        from io import BytesIO
        
        src = img_tag.get('src', '')
        if not src.startswith('data:image'):
            return
        
        try:
            # Decode base64
            header, data = src.split(',', 1)
            img_bytes = base64.b64decode(data)
            
            # Open image with PIL
            pil_img = Image.open(BytesIO(img_bytes))
            img_width, img_height = pil_img.size
            
            # Create new slide
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            
            # Add title
            title_box = slide.shapes.add_textbox(
                self.margin_left, Inches(0.2),
                self.content_width, Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = section_title if section_title else "Analysis Chart"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['primary_red']
            
            # Calculate image size (adjusted to fit slide)
            available_width = self.content_width
            available_height = self.slide_height - Inches(1.2)  # Excluding title space
            
            # Resize while maintaining ratio
            img_ratio = img_width / img_height
            available_ratio = available_width / available_height
            
            if img_ratio > available_ratio:
                # Based on width
                final_width = available_width
                final_height = available_width / img_ratio
            else:
                # Based on height
                final_height = available_height
                final_width = available_height * img_ratio
            
            # Center alignment
            img_left = self.margin_left + (available_width - final_width) / 2
            img_top = Inches(0.7) + (available_height - final_height) / 2
            
            # Save image to BytesIO
            img_stream = BytesIO()
            pil_img.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            # Add image to slide
            slide.shapes.add_picture(
                img_stream,
                img_left, img_top,
                final_width, final_height
            )
            
            logger.info(f"Image slide created: {section_title} ({img_width}x{img_height})")
            
        except Exception as e:
            logger.error(f"Image slide creation failed: {e}")
    
    def _create_gene_content_slides(self, gene_section: Tag, gene_title: str) -> None:
        """Create detailed content slides for Gene section"""
        
        # Major institution recommendation status table
        major_table = gene_section.find('table', class_='major-institution-table')
        if major_table:
            self._create_table_slide(major_table, f"{gene_title} - Major Institution Recommendations")
        
        # Manufacturer-specific commercial kit table
        company_table = gene_section.find('table', class_='company-product-table')
        if company_table:
            self._create_table_slide(company_table, f"{gene_title} - Commercial Kits by Manufacturer")
        
        # Reference cards
        reference_cards = gene_section.find_all('div', class_='reference-card')
        for i, card in enumerate(reference_cards, 1):
            self._create_reference_slide(card, f"{gene_title} - Reference {i}")
    
    def _create_table_slide(self, table_elem: Tag, slide_title: str) -> None:
        """Create table-only slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.margin_left, self.margin_top - Inches(0.2),
            self.content_width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Add table
        self._add_table_to_slide(
            slide, table_elem,
            self.margin_left, self.margin_top + Inches(0.6),
            self.content_width, Inches(5.5)
        )
    
    def _create_reference_slide(self, reference_card: Tag, slide_title: str) -> None:
        """Create Reference slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        y_position = self.margin_top - Inches(0.2)
        
        # Reference number
        ref_number = reference_card.find('div', class_='reference-number')
        if ref_number:
            ref_box = slide.shapes.add_textbox(
                self.margin_left, y_position,
                self.content_width, Inches(0.4)
            )
            ref_frame = ref_box.text_frame
            ref_frame.text = ref_number.get_text(strip=True)
            ref_para = ref_frame.paragraphs[0]
            ref_para.font.size = Pt(24)
            ref_para.font.bold = True
            ref_para.font.color.rgb = self.colors['primary_red']
            y_position += Inches(0.5)
        
        # Reference title
        ref_title = reference_card.find('div', class_='reference-title')
        if ref_title:
            title_box = slide.shapes.add_textbox(
                self.margin_left, y_position,
                self.content_width, Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.text = ref_title.get_text(strip=True)
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['gray_800']
            y_position += Inches(0.7)
        
        # Reference meta information
        ref_meta = reference_card.find('div', class_='reference-meta')
        if ref_meta:
            meta_items = ref_meta.find_all('div', class_='reference-meta-item')
            meta_text = " | ".join([item.get_text(strip=True) for item in meta_items])
            
            meta_box = slide.shapes.add_textbox(
                self.margin_left, y_position,
                self.content_width, Inches(0.4)
            )
            meta_frame = meta_box.text_frame
            meta_frame.text = meta_text
            meta_frame.word_wrap = True
            meta_para = meta_frame.paragraphs[0]
            meta_para.font.size = Pt(11)
            meta_para.font.color.rgb = self.colors['gray_600']
            y_position += Inches(0.5)
        
        # Reference summary
        ref_summary = reference_card.find('div', class_='reference-summary')
        if ref_summary:
            summary_box = slide.shapes.add_textbox(
                self.margin_left, y_position,
                self.content_width, Inches(3)
            )
            summary_frame = summary_box.text_frame
            summary_frame.text = ref_summary.get_text(strip=True)
            summary_frame.word_wrap = True
            
            for paragraph in summary_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = self.colors['gray_800']
                paragraph.line_spacing = 1.4
            y_position += Inches(3.2)
        
        # Evidence table
        evidence_table = reference_card.find('div', class_='evidence-table')
        if evidence_table:
            # Convert table to HTML table format and add
            evidence_rows = evidence_table.find_all('div', class_='evidence-row')
            if evidence_rows:
                self._add_evidence_to_slide(slide, evidence_rows, y_position)
    
    def _add_evidence_to_slide(self, slide, evidence_rows: List[Tag], y_position: float) -> None:
        """Add Evidence information to slide"""
        # Evidence title
        evidence_title_box = slide.shapes.add_textbox(
            self.margin_left, y_position,
            self.content_width, Inches(0.3)
        )
        evidence_title_frame = evidence_title_box.text_frame
        evidence_title_frame.text = "Evidence Details"
        evidence_title_para = evidence_title_frame.paragraphs[0]
        evidence_title_para.font.size = Pt(14)
        evidence_title_para.font.bold = True
        evidence_title_para.font.color.rgb = self.colors['gray_800']
        
        y_position += Inches(0.4)
        
        # Convert each evidence row to text
        for i, row in enumerate(evidence_rows[:2]):  # Display maximum 2
            evidence_header = row.find('div', class_='evidence-header')
            evidence_cell = row.find('div', class_='evidence-cell')
            
            if evidence_cell:
                # Combine all text into one
                texts = [elem.strip() for elem in evidence_cell.stripped_strings]
                combined_text = " | ".join(texts)
                
                text_box = slide.shapes.add_textbox(
                    self.margin_left, y_position,
                    self.content_width, Inches(0.8)
                )
                text_frame = text_box.text_frame
                text_frame.text = combined_text
                text_frame.word_wrap = True
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(9)
                    paragraph.font.color.rgb = self.colors['gray_800']
                
                y_position += Inches(0.9)
    
    def _add_key_value_cards(
        self,
        slide,
        table_elem: Tag,
        left: float,
        top: float,
        width: float,
        height: float
    ) -> List[Any]:
        """Display key-value format table in simple card style"""
        
        tbody = table_elem.find('tbody')
        if not tbody:
            return []
        
        rows = tbody.find_all('tr')
        if not rows:
            return []
        
        # Simple solid color (dark gray)
        label_bg_color = RGBColor(55, 65, 81)      # Gray-700
        value_bg_color = RGBColor(249, 250, 251)   # Gray-50
        border_color = RGBColor(209, 213, 219)     # Gray-300
        
        y_position = top
        card_height = Inches(1.3)
        card_spacing = Inches(0.12)
        label_width = Inches(1.6)
        
        for i, tr in enumerate(rows):
            cells = tr.find_all(['th', 'td'])
            if len(cells) < 2:
                continue
            
            label = self._clean_text(cells[0].get_text(strip=True))
            value = self._clean_text(cells[1].get_text(strip=True))
            
            # Label area (angular rectangle)
            label_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,  # Angular rectangle
                left, y_position,
                label_width, card_height
            )
            label_shape.fill.solid()
            label_shape.fill.fore_color.rgb = label_bg_color
            label_shape.line.fill.background()  # No border
            
            # Label text
            label_tf = label_shape.text_frame
            label_tf.word_wrap = True
            label_tf.paragraphs[0].text = label
            label_tf.paragraphs[0].font.size = Pt(13)
            label_tf.paragraphs[0].font.bold = True
            label_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            label_shape.text_frame.margin_left = Pt(10)
            label_shape.text_frame.margin_right = Pt(10)
            label_shape.text_frame.margin_top = Pt(10)
            label_shape.text_frame.margin_bottom = Pt(10)
            
            # Vertical center alignment
            from pptx.enum.text import MSO_ANCHOR
            label_tf.auto_size = None
            
            # Value area (angular rectangle)
            value_left = left + label_width
            value_width = width - label_width
            
            value_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,  # Angular rectangle
                value_left, y_position,
                value_width, card_height
            )
            value_shape.fill.solid()
            value_shape.fill.fore_color.rgb = value_bg_color
            value_shape.line.color.rgb = border_color
            value_shape.line.width = Pt(1)
            
            # Value text
            value_tf = value_shape.text_frame
            value_tf.word_wrap = True
            value_tf.paragraphs[0].text = value
            value_tf.paragraphs[0].font.size = Pt(11)
            value_tf.paragraphs[0].font.color.rgb = self.colors['gray_800']
            value_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            value_shape.text_frame.margin_left = Pt(12)
            value_shape.text_frame.margin_right = Pt(12)
            value_shape.text_frame.margin_top = Pt(10)
            value_shape.text_frame.margin_bottom = Pt(10)
            
            y_position += card_height + card_spacing
        
        return [slide]
    
    def _add_improved_table(
        self, 
        slide, 
        table_elem: Tag, 
        left: float, 
        top: float, 
        width: float, 
        height: float
    ) -> List[Any]:
        """Add table with improved readability (auto-split support, academic style)"""
        
        # Check if key-value format table (no thead, 2 columns, max 5 rows)
        thead = table_elem.find('thead')
        tbody = table_elem.find('tbody')
        
        if not thead and tbody:
            rows = tbody.find_all('tr')
            if len(rows) <= 5:  # 5 rows or less
                first_row = rows[0] if rows else None
                if first_row:
                    cells = first_row.find_all(['th', 'td'])
                    if len(cells) == 2:  # 2 columns (key-value format)
                        # Display in key-value card style
                        return self._add_key_value_cards(slide, table_elem, left, top, width, height)
        
        # Extract table data and width info (including colspan)
        rows_data = []
        header_rows = []
        body_rows = []
        col_widths_html = []
        has_header = False
        merge_info = []  # [(row_idx, col_idx, colspan, rowspan), ...]
        cell_styles = {}  # {(row_idx, col_idx): {'bold': bool, 'color': RGBColor, 'link': str}, ...}
        
        def extract_row_data(tr, row_idx, is_header=False):
            """Extract row data (including colspan handling)"""
            cells = tr.find_all(['th', 'td'])
            row_data = []
            col_idx = 0
            
            for cell in cells:
                text = self._clean_text(cell.get_text(strip=True))
                colspan = int(cell.get('colspan', 1))
                rowspan = int(cell.get('rowspan', 1))
                
                # Extract styles
                styles = self._extract_cell_styles(cell)
                if styles['bold'] or styles['color'] or styles['link']:
                    cell_styles[(row_idx, col_idx)] = styles
                
                row_data.append(text)
                # Add empty cells if colspan exists
                for _ in range(colspan - 1):
                    row_data.append('')
                
                # Save merge info
                if colspan > 1 or rowspan > 1:
                    merge_info.append((row_idx, col_idx, colspan, rowspan))
                
                col_idx += colspan
            
            return row_data
        
        # Process thead
        if thead:
            has_header = True
            header_trs = thead.find_all('tr')
            for idx, tr in enumerate(header_trs):
                row_data = extract_row_data(tr, len(rows_data))
                header_rows.append(row_data)
                rows_data.append(row_data)
                
                # Extract width info from first header row
                if not col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    col_widths_html = self._extract_column_widths(cells)
        
        # Process tbody
        tbody = table_elem.find('tbody')
        if tbody:
            body_trs = tbody.find_all('tr')
            for idx, tr in enumerate(body_trs):
                row_data = extract_row_data(tr, len(rows_data))
                body_rows.append(row_data)
                rows_data.append(row_data)
                
                # Extract width from first row if no thead
                if not has_header and idx == 0 and not col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    col_widths_html = self._extract_column_widths(cells)
        
        # If no thead and no tbody (use tr directly)
        if not has_header and not tbody:
            all_rows = table_elem.find_all('tr')
            for idx, tr in enumerate(all_rows):
                row_data = extract_row_data(tr, len(rows_data))
                body_rows.append(row_data)
                rows_data.append(row_data)
                
                # Extract width info from first row
                if idx == 0 and not col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    col_widths_html = self._extract_column_widths(cells)
        
        if not rows_data:
            return []
        
        # Determine column count
        max_cols = max(len(row) for row in rows_data)
        
        # Make all rows have same column count
        for row in rows_data:
            while len(row) < max_cols:
                row.append("")
        
        # Split table into multiple slides if too large
        header_count = len(header_rows)
        body_count = len(body_rows)
        
        created_slides = []
        
        # If splitting is required
        if body_count > self.max_rows_per_slide:
            # Split into multiple slides
            num_chunks = (body_count + self.max_rows_per_slide - 1) // self.max_rows_per_slide
            
            logger.info(f"Table split: {body_count} rows split into {num_chunks} slides")
            
            for chunk_idx in range(num_chunks):
                start_idx = chunk_idx * self.max_rows_per_slide
                end_idx = min(start_idx + self.max_rows_per_slide, body_count)
                
                # Prepare chunk data
                chunk_data = header_rows + body_rows[start_idx:end_idx]
                
                # Filter merge_info for this chunk
                chunk_merge_info = []
                for row_idx, col_idx, colspan, rowspan in merge_info:
                    # Adjust for rows after header_count
                    if row_idx < header_count:
                        chunk_merge_info.append((row_idx, col_idx, colspan, rowspan))
                    elif row_idx - header_count >= start_idx and row_idx - header_count < end_idx:
                        new_row_idx = header_count + (row_idx - header_count - start_idx)
                        chunk_merge_info.append((new_row_idx, col_idx, colspan, rowspan))
                
                # Filter cell_styles for this chunk
                chunk_cell_styles = {}
                for (r, c), styles in cell_styles.items():
                    if r < header_count:
                        chunk_cell_styles[(r, c)] = styles
                    elif r - header_count >= start_idx and r - header_count < end_idx:
                        new_r = header_count + (r - header_count - start_idx)
                        chunk_cell_styles[(new_r, c)] = styles
                
                # Create new slide (use existing slide for first chunk)
                chunk_slide = slide if chunk_idx == 0 else None
                
                # Create table
                result = self._create_ppt_table(
                    chunk_slide,
                    chunk_data,
                    header_count,
                    col_widths_html,
                    left, top, width, height,
                    chunk_idx + 1, num_chunks,
                    chunk_merge_info,
                    max_cols,
                    chunk_cell_styles
                )
                
                if result:
                    created_slides.append(result)
            
            return created_slides
        else:
            # Display on single slide
            result = self._create_ppt_table(
                slide,
                rows_data,
                header_count,
                col_widths_html,
                left, top, width, height,
                1, 1,
                merge_info,
                max_cols,
                cell_styles
            )
            return [result] if result else []
    
    def _extract_column_widths(self, cells: List[Tag]) -> List[Optional[int]]:
        """Extract width attribute from HTML table cells"""
        widths = []
        for cell in cells:
            width = None
            
            # Extract width from style attribute
            style = cell.get('style', '')
            if 'width:' in style:
                import re
                match = re.search(r'width:\s*(\d+)(?:px|%)?', style)
                if match:
                    width = int(match.group(1))
            
            # Check width attribute directly
            elif cell.get('width'):
                try:
                    width = int(cell.get('width').replace('px', '').replace('%', ''))
                except:
                    pass
            
            widths.append(width)
        
        return widths
    
    def _create_ppt_table(
        self,
        slide,
        rows_data: List[List[str]],
        header_count: int,
        col_widths_html: List[Optional[int]],
        left: float,
        top: float,
        width: float,
        height: float,
        chunk_num: int = 1,
        total_chunks: int = 1,
        merge_info: List[tuple] = None,
        max_cols_override: int = None,
        cell_styles: Dict[Tuple[int, int], Dict[str, Any]] = None
    ) -> Any:
        """Create actual PowerPoint table (academic style)"""
        
        if merge_info is None:
            merge_info = []
        if cell_styles is None:
            cell_styles = {}
        
        # If new slide is needed
        if slide is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            # Adjust table start position for additional slides
            top = Inches(0.6)  # Space for title
            height = self.slide_height - top - self.margin_bottom
        
        if not rows_data:
            return None
        
        max_cols = max_cols_override if max_cols_override else len(rows_data[0])
        
        # Adjust font size based on rows/columns
        base_font_size = 8 if len(rows_data) > 15 or max_cols > 6 else 9
        header_font_size = base_font_size + 1
        
        # Calculate table height - minimize to fit content
        row_count = len(rows_data)
        # Calculate minimum height per row (set smaller)
        min_row_height = Inches(0.22)  # Reduce minimum row height
        required_height = min_row_height * row_count
        
        # Use only needed table height (less than slide height)
        height = min(required_height, height)
        
        if row_count > 20:
            base_font_size = 7
            header_font_size = 8
        
        try:
            # Create PowerPoint table
            ppt_table = slide.shapes.add_table(
                row_count, max_cols,
                left, top, width, height
            ).table
            
            # Academic style: set borders for all cells first
            from pptx.oxml.ns import qn
            from pptx.oxml import parse_xml
            
            # Fill table data and apply academic style
            for i, row_data in enumerate(rows_data):
                for j, cell_data in enumerate(row_data):
                    if j >= max_cols:
                        continue
                    cell = ppt_table.cell(i, j)
                    
                    # Set text
                    cell.text = str(cell_data) if j < len(row_data) else ""
                    
                    # Vertical center alignment (MIDDLE)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Minimize cell margins
                    cell.margin_left = Pt(4)
                    cell.margin_right = Pt(4)
                    cell.margin_top = Pt(2)
                    cell.margin_bottom = Pt(2)
                    
                    # Academic style: no background color (transparent)
                    cell.fill.background()
                    
                    # Get styles extracted from HTML
                    html_style = cell_styles.get((i, j), {})
                    has_custom_bold = html_style.get('bold', False)
                    custom_color = html_style.get('color')
                    has_link = html_style.get('link')
                    
                    # Set paragraph format
                    for paragraph in cell.text_frame.paragraphs:
                        # Header row
                        if i < header_count:
                            paragraph.font.size = Pt(header_font_size)
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = self.colors['black']
                            paragraph.alignment = PP_ALIGN.CENTER
                            # Header no word wrap (to prevent Gene etc. from wrapping)
                            cell.text_frame.word_wrap = False
                        else:
                            paragraph.font.size = Pt(base_font_size)
                            
                            # Apply HTML style (Bold)
                            if has_custom_bold:
                                paragraph.font.bold = True
                            
                            # Apply HTML style (Color)
                            if custom_color:
                                paragraph.font.color.rgb = custom_color
                            else:
                                paragraph.font.color.rgb = self.colors['gray_800']
                            
                            # Blue + underline if link exists
                            if has_link:
                                paragraph.font.color.rgb = RGBColor(0, 102, 204)
                                paragraph.font.underline = True
                            
                            # Columns requiring left alignment (long text)
                            if len(cell_data) > 30 or '\n' in cell_data:
                                paragraph.alignment = PP_ALIGN.LEFT
                            else:
                                paragraph.alignment = PP_ALIGN.CENTER
                            
                            # Data rows allow word wrap
                            cell.text_frame.word_wrap = True
                        
                        # Line spacing
                        paragraph.line_spacing = 1.1
            
            # Apply academic style borders (lines only at top, header bottom, bottom)
            self._apply_academic_table_borders(ppt_table, header_count, row_count, max_cols)
            
            # Apply cell merge
            for row_idx, col_idx, colspan, rowspan in merge_info:
                try:
                    if row_idx < row_count and col_idx < max_cols:
                        start_cell = ppt_table.cell(row_idx, col_idx)
                        end_row = min(row_idx + rowspan - 1, row_count - 1)
                        end_col = min(col_idx + colspan - 1, max_cols - 1)
                        
                        end_cell = ppt_table.cell(end_row, end_col)
                        start_cell.merge(end_cell)
                        
                        # Center alignment for merged cells
                        for paragraph in start_cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                except Exception as merge_err:
                    logger.debug(f"Cell merge failed: {merge_err}")
            
            # Adjust column width based on HTML width attribute
            if col_widths_html and any(w is not None for w in col_widths_html):
                self._apply_html_column_widths(ppt_table, col_widths_html, width)
            else:
                # Auto adjustment
                self._adjust_column_widths(ppt_table, rows_data)
            
            return slide
            
        except Exception as e:
            logger.error(f"Table addition failed: {e}")
    
    def _apply_academic_table_borders(self, ppt_table, header_count: int, row_count: int, col_count: int) -> None:
        """Apply table borders (thick lines for header, thin horizontal lines for data rows)"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import nsmap
        
        # Define line thickness
        thick_line = Pt(1.5)  # Thick line (header top/bottom)
        thin_line = Pt(0.5)   # Thin line (data rows)
        no_line = Pt(0)       # No line
        
        black = RGBColor(0, 0, 0)
        gray_line = RGBColor(200, 200, 200)  # Light gray line
        
        for i in range(row_count):
            for j in range(col_count):
                try:
                    cell = ppt_table.cell(i, j)
                    
                    # Set border for each cell
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    # Top line
                    if i == 0:
                        # First row: thick line at top
                        self._set_cell_border(cell, 'top', thick_line, black)
                    elif i == header_count and header_count > 0:
                        # First data row: thick line at header bottom already exists
                        self._set_cell_border(cell, 'top', no_line, black)
                    else:
                        # Data row top: thin gray line
                        self._set_cell_border(cell, 'top', thin_line, gray_line)
                    
                    # Bottom line
                    if i == row_count - 1:
                        # Last row: thick line at bottom
                        self._set_cell_border(cell, 'bottom', thick_line, black)
                    elif i == header_count - 1 and header_count > 0:
                        # Header last row: thick line at bottom
                        self._set_cell_border(cell, 'bottom', thick_line, black)
                    else:
                        # Data row: thin gray line at bottom
                        self._set_cell_border(cell, 'bottom', thin_line, gray_line)
                    
                    # No left/right borders
                    self._set_cell_border(cell, 'left', no_line, black)
                    self._set_cell_border(cell, 'right', no_line, black)
                    
                except Exception as e:
                    pass  # Ignore if border setting fails
    
    def _set_cell_border(self, cell, side: str, width, color: RGBColor) -> None:
        """Set specific border for cell (improved version)"""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Border element name
        border_map = {
            'top': 'a:lnT',
            'bottom': 'a:lnB', 
            'left': 'a:lnL',
            'right': 'a:lnR'
        }
        
        border_elem_name = border_map.get(side)
        if not border_elem_name:
            return
        
        # Remove existing border element
        for existing in list(tcPr):
            if existing.tag == qn(border_elem_name):
                tcPr.remove(existing)
        
        # Convert to EMU units (1 pt = 12700 EMU)
        width_emu = int(width) if width > 0 else 0
        
        # Create new border element
        ln = etree.Element(qn(border_elem_name))
        
        if width_emu > 0:
            ln.set('w', str(width_emu))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')
            
            # Set color
            solidFill = etree.SubElement(ln, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
            
            # Preset dash
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'solid')
        else:
            ln.set('w', '0')
            noFill = etree.SubElement(ln, qn('a:noFill'))
        
        # Add as first child of tcPr (order matters)
        tcPr.insert(0, ln)
    
    def _apply_html_column_widths(
        self, 
        ppt_table, 
        col_widths_html: List[Optional[int]], 
        total_width: float
    ) -> None:
        """Apply width attribute extracted from HTML to PowerPoint table
        
        HTML width is usually in pixels, specified columns take that width,
        unspecified columns take remaining space.
        """
        try:
            col_count = len(col_widths_html)
            if col_count == 0:
                return
            
            # Separate columns with specified width and unspecified
            specified_widths = [w for w in col_widths_html if w is not None]
            unspecified_count = col_widths_html.count(None)
            
            if not specified_widths:
                return
            
            # Convert width specified in HTML to PowerPoint units
            # Generally HTML tables are based on 800-1000px
            # PowerPoint slide is about 9.4 inches = about 900px (at 96dpi)
            html_to_ppt_ratio = total_width / 900  # 1px = approximately this ratio of EMU
            
            # Calculate width of specified columns first
            specified_total_ppt = 0
            for html_width in col_widths_html:
                if html_width is not None:
                    ppt_width = int(html_width * html_to_ppt_ratio)
                    specified_total_ppt += ppt_width
            
            # Calculate remaining width
            remaining_width = total_width - specified_total_ppt
            
            # Adjust ratio if insufficient space for unspecified columns
            if remaining_width < 0 or (unspecified_count > 0 and remaining_width < total_width * 0.3):
                # Limit specified columns to 30% of total, rest to unspecified
                specified_portion = 0.3 if unspecified_count > 0 else 1.0
                total_specified_html = sum(specified_widths)
                
                for j, html_width in enumerate(col_widths_html):
                    if html_width is not None:
                        proportion = html_width / total_specified_html
                        ppt_table.columns[j].width = int(total_width * specified_portion * proportion)
                
                if unspecified_count > 0:
                    remaining = total_width * (1 - specified_portion)
                    equal_width = int(remaining / unspecified_count)
                    for j, html_width in enumerate(col_widths_html):
                        if html_width is None:
                            ppt_table.columns[j].width = equal_width
            else:
                # Apply as is if sufficient space
                for j, html_width in enumerate(col_widths_html):
                    if html_width is not None:
                        ppt_table.columns[j].width = int(html_width * html_to_ppt_ratio)
                
                # Distribute remaining width equally to unspecified columns
                if unspecified_count > 0:
                    equal_width = int(remaining_width / unspecified_count)
                    for j, html_width in enumerate(col_widths_html):
                        if html_width is None:
                            ppt_table.columns[j].width = equal_width
        
        except Exception as e:
            logger.debug(f"HTML width application failed, switching to auto adjustment: {e}")
    
    def _adjust_column_widths(self, ppt_table, rows_data: List[List[str]]) -> None:
        """Auto adjust column width (when no HTML width) - based on text length"""
        try:
            col_count = len(rows_data[0]) if rows_data else 0
            if col_count == 0:
                return
            
            # Calculate current total table width
            total_table_width = sum(col.width for col in ppt_table.columns)
            
            # Calculate maximum text length for each column (weighted)
            max_lengths = [0] * col_count
            for row in rows_data:
                for j, cell in enumerate(row):
                    cell_text = str(cell)
                    # Korean needs wider space (1.5x)
                    korean_count = len([c for c in cell_text if ord(c) >= 0xAC00 and ord(c) <= 0xD7A3])
                    english_count = len(cell_text) - korean_count
                    weighted_length = english_count + (korean_count * 1.8)
                    max_lengths[j] = max(max_lengths[j], weighted_length)
            
            # Ensure minimum width (5% minimum per column)
            min_proportion = 0.05
            
            # Total weighted length
            total_length = sum(max_lengths)
            if total_length == 0:
                # Distribute equally if all empty cells
                equal_width = total_table_width // col_count
                for j in range(col_count):
                    ppt_table.columns[j].width = equal_width
                return
            
            # Allocate width proportionally to each column
            for j in range(col_count):
                proportion = max_lengths[j] / total_length
                # Ensure minimum width
                proportion = max(proportion, min_proportion)
                col_width = int(total_table_width * proportion)
                ppt_table.columns[j].width = col_width
                
            logger.debug(f"Column width auto adjustment completed: {[ppt_table.columns[j].width for j in range(col_count)]}")
        
        except Exception as e:
            logger.debug(f"Column width adjustment failed (ignored): {e}")
    
    def _create_evidence_table_slide(self, evidence_div: Tag, section_title: str) -> None:
        """Create Evidence table slide (maintain original header, include links)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.margin_left, self.margin_top - Inches(0.2),
            self.content_width, Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Find Evidence rows
        evidence_rows = evidence_div.find_all('div', class_='evidence-row')
        
        if not evidence_rows:
            # Also try evidence-cell
            evidence_rows = evidence_div.find_all('div', class_='evidence-cell')
        
        if not evidence_rows:
            return
        
        # Process maximum 10 rows (to fit slide)
        max_rows = min(len(evidence_rows), 10)
        
        # Extract table data
        table_data = []
        link_data = []  # Save link info for each cell [(row, col, url), ...]
        
        # Extract header (maintain original)
        header_div = evidence_div.find('div', class_='evidence-header')
        if header_div:
            headers = [elem.strip() for elem in header_div.stripped_strings]
            # Use original header (max 8 columns)
            if len(headers) > 0:
                table_data.append(headers[:8])
        
        # Extract data rows
        for row_idx, row in enumerate(evidence_rows[:max_rows]):
            row_texts = []
            text_elements = row.find_all('div', class_='evidence-text')
            
            for col_idx, elem in enumerate(text_elements[:8]):
                # Check link
                link = elem.find('a')
                if link:
                    link_text = link.get_text(strip=True)
                    link_url = link.get('href', '')
                    row_texts.append(link_text)
                    if link_url:
                        link_data.append((len(table_data), col_idx, link_url))
                else:
                    text = self._clean_text(elem.get_text(strip=True))
                    # Truncate text if too long
                    if len(text) > 80:
                        text = text[:77] + "..."
                    row_texts.append(text)
            
            if row_texts:
                table_data.append(row_texts)
        
        if len(table_data) <= 1:
            return
        
        # Unify column count
        max_cols = max(len(row) for row in table_data)
        for row in table_data:
            while len(row) < max_cols:
                row.append('')
        
        # Create table
        try:
            table_top = self.margin_top + Inches(0.4)
            table_height = self.slide_height - table_top - self.margin_bottom
            
            ppt_table = slide.shapes.add_table(
                len(table_data), max_cols,
                self.margin_left, table_top,
                self.content_width, table_height
            ).table
            
            # Fill data
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    if j >= max_cols:
                        continue
                    cell = ppt_table.cell(i, j)
                    cell.text = str(cell_data)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Cell margins
                    cell.margin_left = Pt(3)
                    cell.margin_right = Pt(3)
                    cell.margin_top = Pt(2)
                    cell.margin_bottom = Pt(2)
                    
                    # No background color
                    cell.fill.background()
                    
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(7)
                        
                        if i == 0:  # Header
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = self.colors['black']
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.text_frame.word_wrap = False
                        else:
                            paragraph.font.color.rgb = self.colors['gray_800']
                            cell.text_frame.word_wrap = True
                            # Left align document title and AI summary columns (index 2, 5)
                            if j in [2, 5]:
                                paragraph.alignment = PP_ALIGN.LEFT
                            else:
                                paragraph.alignment = PP_ALIGN.CENTER
                            
                            # Display Link column in blue
                            if cell_data == 'Link':
                                paragraph.font.color.rgb = RGBColor(0, 102, 204)
                                paragraph.font.underline = True
            
            # Add hyperlinks
            for row_idx, col_idx, url in link_data:
                try:
                    cell = ppt_table.cell(row_idx, col_idx)
                    # Adding hyperlinks directly to table cells is difficult in PowerPoint
                    # Instead keep underline and blue color on text
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 102, 204)
                            run.font.underline = True
                except:
                    pass
            
            # Apply borders
            self._apply_academic_table_borders(ppt_table, 1, len(table_data), max_cols)
        
        except Exception as e:
            logger.error(f"Evidence table creation failed: {e}")
    
    def _add_table_to_slide(
        self, 
        slide, 
        table_elem: Tag, 
        left: float, 
        top: float, 
        width: float, 
        height: float
    ) -> None:
        """Convert HTML table to PowerPoint table and add"""
        
        # Extract table data
        rows_data = []
        col_widths_html = []
        header_count = 0
        merge_info = []  # [(row_idx, col_idx, colspan, rowspan), ...]
        
        # Process both thead and tbody
        thead = table_elem.find('thead')
        tbody = table_elem.find('tbody')
        
        if thead:
            header_rows = thead.find_all('tr')
            for idx, tr in enumerate(header_rows):
                cells = tr.find_all(['th', 'td'])
                row_data = []
                col_idx = 0
                for cell in cells:
                    text = cell.get_text(strip=True)
                    colspan = int(cell.get('colspan', 1))
                    rowspan = int(cell.get('rowspan', 1))
                    
                    row_data.append(text)
                    # Add empty cells if colspan exists
                    for _ in range(colspan - 1):
                        row_data.append('')
                    
                    # Save merge info
                    if colspan > 1 or rowspan > 1:
                        merge_info.append((len(rows_data), col_idx, colspan, rowspan))
                    
                    col_idx += colspan
                
                rows_data.append(row_data)
                header_count += 1
                
                # Extract width from first header row
                if idx == 0 and not col_widths_html:
                    col_widths_html = self._extract_column_widths(cells)
        
        if tbody:
            body_rows = tbody.find_all('tr')
            for idx, tr in enumerate(body_rows):
                cells = tr.find_all(['th', 'td'])
                row_data = []
                col_idx = 0
                for cell in cells:
                    text = cell.get_text(strip=True)
                    colspan = int(cell.get('colspan', 1))
                    rowspan = int(cell.get('rowspan', 1))
                    
                    row_data.append(text)
                    # Add empty cells if colspan exists
                    for _ in range(colspan - 1):
                        row_data.append('')
                    
                    # Save merge info
                    if colspan > 1 or rowspan > 1:
                        merge_info.append((len(rows_data), col_idx, colspan, rowspan))
                    
                    col_idx += colspan
                
                rows_data.append(row_data)
                
                # Extract width from first row if no thead
                if not thead and idx == 0 and not col_widths_html:
                    col_widths_html = self._extract_column_widths(cells)
        
        # Return if table is empty
        if not rows_data:
            return
        
        # Determine column count
        max_cols = max(len(row) for row in rows_data)
        
        # Make all rows have same column count
        for row in rows_data:
            while len(row) < max_cols:
                row.append("")
        
        # Create PowerPoint table
        try:
            ppt_table = slide.shapes.add_table(
                len(rows_data), max_cols,
                left, top, width, height
            ).table
            
            # Fill table data (academic style)
            for i, row_data in enumerate(rows_data):
                for j, cell_data in enumerate(row_data):
                    cell = ppt_table.cell(i, j)
                    cell.text = str(cell_data)
                    
                    # Vertical center alignment (MIDDLE)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Minimize cell margins
                    cell.margin_left = Pt(4)
                    cell.margin_right = Pt(4)
                    cell.margin_top = Pt(2)
                    cell.margin_bottom = Pt(2)
                    
                    # Academic style: no background color
                    cell.fill.background()
                    
                    # Set text format
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)
                        
                        # Header row style
                        if i < header_count or i == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = self.colors['black']
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.text_frame.word_wrap = False  # Prevent header word wrap
                        else:
                            paragraph.font.color.rgb = self.colors['gray_800']
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.text_frame.word_wrap = True  # Data allows word wrap
            
            # Apply academic style borders
            self._apply_academic_table_borders(ppt_table, header_count, len(rows_data), max_cols)
            
            # Adjust column width based on HTML width attribute
            if col_widths_html and any(w is not None for w in col_widths_html):
                self._apply_html_column_widths(ppt_table, col_widths_html, width)
            else:
                # Auto adjustment
                self._adjust_column_widths(ppt_table, rows_data)
            
            # Apply cell merge
            for row_idx, col_idx, colspan, rowspan in merge_info:
                try:
                    start_cell = ppt_table.cell(row_idx, col_idx)
                    end_row = row_idx + rowspan - 1
                    end_col = col_idx + colspan - 1
                    
                    # Range check
                    if end_row < len(rows_data) and end_col < max_cols:
                        end_cell = ppt_table.cell(end_row, end_col)
                        start_cell.merge(end_cell)
                        
                        # Center alignment for merged cells
                        for paragraph in start_cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                except Exception as merge_err:
                    logger.debug(f"Cell merge failed: {merge_err}")
            
        except Exception as e:
            logger.error(f"Table addition failed: {e}")
    
    def _clean_text(self, text: str) -> str:
        """Clean text (remove unnecessary spaces, special characters)"""
        # Convert consecutive spaces to one
        text = re.sub(r'\s+', ' ', text)
        # Remove leading/trailing spaces
        text = text.strip()
        return text
    
    def _capture_element_screenshot(self, selector: str, title: str, index: int = 0) -> None:
        """
        Capture HTML element as screenshot using Playwright and add to slide
        
        Args:
            selector: CSS selector (e.g., '.SeqViewerApp')
            title: Slide title
            index: Index of element to capture among multiple elements (starting from 0)
        """
        if not self.html_path:
            logger.warning(f"Cannot capture screenshot because HTML path is not set: {selector}")
            return
        
        try:
            from playwright.sync_api import sync_playwright
            import tempfile
            
            logger.info(f"Capturing element with Playwright: {selector} (index={index})")
            
            with sync_playwright() as p:
                # Launch Chromium browser
                browser = p.chromium.launch(headless=True)
                
                # Create page (with generous viewport size)
                page = browser.new_page(viewport={'width': 1400, 'height': 900})
                
                # Load HTML file
                file_url = f"file://{self.html_path}"
                page.goto(file_url, wait_until='networkidle')
                
                # Wait for JavaScript rendering
                page.wait_for_timeout(2000)
                
                # Find element
                elements = page.locator(selector)
                count = elements.count()
                
                if count == 0:
                    logger.warning(f"Cannot find element: {selector}")
                    browser.close()
                    return
                
                if index >= count:
                    logger.warning(f"Index out of range: {index} >= {count}")
                    browser.close()
                    return
                
                # Select element at specific index
                element = elements.nth(index)
                
                # Scroll to make element visible
                element.scroll_into_view_if_needed()
                page.wait_for_timeout(500)
                
                # Save screenshot to temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    screenshot_path = tmp_file.name
                
                element.screenshot(path=screenshot_path)
                logger.info(f"Screenshot saved: {screenshot_path}")
                
                browser.close()
                
                # Add screenshot to slide
                self._create_screenshot_slide(screenshot_path, title)
                
                # Delete temporary file
                import os
                os.unlink(screenshot_path)
                
        except ImportError:
            logger.error("Playwright is not installed. Run 'pip install playwright && playwright install chromium'.")
        except Exception as e:
            logger.error(f"Failed to capture element screenshot: {e}")
    
    def _create_screenshot_slide(self, image_path: str, title: str) -> None:
        """
        Add screenshot image to slide
        
        Args:
            image_path: Image file path
            title: Slide title
        """
        try:
            from PIL import Image
            
            # Add new slide
            blank_layout = self.prs.slide_layouts[6]  # Empty layout
            slide = self.prs.slides.add_slide(blank_layout)
            
            # Add title
            title_top = Inches(0.3)
            title_shape = slide.shapes.add_textbox(
                self.margin_left,
                title_top,
                self.content_width,
                Inches(0.5)
            )
            tf = title_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['gray_900']
            
            # Check image size
            with Image.open(image_path) as img:
                img_width, img_height = img.size
            
            # Calculate available area
            content_top = Inches(1.0)
            available_width = self.content_width
            available_height = self.slide_height - content_top - self.margin_bottom
            
            # Resize while maintaining ratio
            width_ratio = available_width / img_width
            height_ratio = available_height / img_height
            scale = min(width_ratio, height_ratio)
            
            # Final size (converted to EMU)
            final_width = int(img_width * scale)
            final_height = int(img_height * scale)
            
            # Center alignment
            left = self.margin_left + (available_width - final_width) / 2
            top = content_top + (available_height - final_height) / 2
            
            # Add image
            slide.shapes.add_picture(
                image_path,
                left,
                top,
                final_width,
                final_height
            )
            
            logger.info(f"Screenshot slide added: {title}")
            
        except Exception as e:
            logger.error(f"Failed to create screenshot slide: {e}")


def convert_html_to_pptx(html_path: Path, output_path: Path) -> None:
    """
    Convenience function to convert HTML file to PPTX
    
    Args:
        html_path: Input HTML file path
        output_path: Output PPTX file path
    
    Example:
        >>> from pathlib import Path
        >>> html_path = Path("input.html")
        >>> output_path = Path("output.pptx")
        >>> convert_html_to_pptx(html_path, output_path)
    """
    converter = HtmlToPptxConverter()
    converter.convert(html_path, output_path)


if __name__ == "__main__":
    # Run test
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python html_to_pptx.py <input.html> <output.pptx>")
        sys.exit(1)
    
    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    
    convert_html_to_pptx(input_path, output_path)
    print(f"Conversion complete: {output_path}")
