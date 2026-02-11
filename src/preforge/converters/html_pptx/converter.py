"""
HTML(.html) to PowerPoint(.pptx) converter

Modularly structured to analyze HTML structure and generate slides by section.
"""
import logging
import tempfile
import os
from pathlib import Path
from typing import List, Optional

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt

from .config import (
    SlideConfig, 
    TableConfig, 
    ColorPalette,
    DEFAULT_SLIDE_CONFIG,
    DEFAULT_TABLE_CONFIG,
    DEFAULT_COLORS
)
from .style_utils import TextUtils
from .table_builder import TableBuilder, TableDataExtractor
from .slide_factory import (
    TitleSlideBuilder,
    ContentSlideBuilder,
    TableSlideBuilder,
    ImageSlideBuilder,
    EvidenceSlideBuilder,
    SectionSlideBuilder,
    ReferenceCardSlideBuilder
)

logger = logging.getLogger(__name__)


class HtmlToPptxConverter:
    """Converter to transform HTML to PowerPoint"""
    
    def __init__(
        self,
        slide_config: SlideConfig = None,
        table_config: TableConfig = None,
        colors: ColorPalette = None
    ):
        """
        Initialize converter
        
        Args:
            slide_config: Slide layout settings
            table_config: Table settings
            colors: Color palette
        """
        self.slide_config = slide_config or DEFAULT_SLIDE_CONFIG
        self.table_config = table_config or DEFAULT_TABLE_CONFIG
        self.colors = colors or DEFAULT_COLORS
        
        self.prs: Optional[Presentation] = None
        self.html_path: Optional[Path] = None
        
        # Slide builders (initialized during convert)
        self._title_builder: Optional[TitleSlideBuilder] = None
        self._content_builder: Optional[ContentSlideBuilder] = None
        self._table_builder: Optional[TableSlideBuilder] = None
        self._image_builder: Optional[ImageSlideBuilder] = None
        self._evidence_builder: Optional[EvidenceSlideBuilder] = None
        self._section_builder: Optional[SectionSlideBuilder] = None
        self._reference_builder: Optional[ReferenceCardSlideBuilder] = None
    
    def convert(self, html_path: Path, output_path: Path) -> None:
        """
        Convert HTML file to PPTX
        
        Args:
            html_path: Input HTML file path
            output_path: Output PPTX file path
        """
        logger.info(f"HTML -> PPTX conversion started: {html_path} -> {output_path}")
        
        self.html_path = Path(html_path).absolute()
        
        # Read HTML file
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'lxml')
        
        # Initialize presentation
        self.prs = Presentation()
        self.prs.slide_width = self.slide_config.width
        self.prs.slide_height = self.slide_config.height
        
        # Initialize slide builders
        self._init_builders()
        
        # Create slides
        self._create_title_slide(soup)
        self._create_analysis_summary_slides(soup)
        self._process_main_content(soup)
        
        # Save
        self.prs.save(str(output_path))
        logger.info(f"Conversion complete: {output_path} (total {len(self.prs.slides)} slides)")
    
    def _init_builders(self) -> None:
        """Initialize slide builders"""
        self._title_builder = TitleSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
        self._content_builder = ContentSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
        self._table_builder = TableSlideBuilder(
            self.prs, self.slide_config, self.colors,
            self.table_config.max_rows_per_slide
        )
        self._image_builder = ImageSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
        self._evidence_builder = EvidenceSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
        self._section_builder = SectionSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
        self._reference_builder = ReferenceCardSlideBuilder(
            self.prs, self.slide_config, self.colors
        )
    
    def _create_title_slide(self, soup: BeautifulSoup) -> None:
        """Create title slide"""
        title_elem = soup.find('div', class_='header-title')
        subtitle_elem = soup.find('div', class_='header-subtitle')
        
        title = title_elem.get_text(strip=True) if title_elem else "GeneSeq Vista AI Agent"
        subtitle = subtitle_elem.get_text(strip=True) if subtitle_elem else ""
        
        self._title_builder.create(title, subtitle)
    
    def _create_analysis_summary_slides(self, soup: BeautifulSoup) -> None:
        """Create Analysis Summary section slides"""
        analysis_div = soup.find('div', class_='analysis-summary')
        if not analysis_div:
            return
        
        summary_sections = analysis_div.find_all('div', class_='summary-section')
        
        for idx, section in enumerate(summary_sections):
            header = section.find('div', class_='section-header')
            header_text = header.get_text(strip=True) if header else f"Summary {idx + 1}"
            
            table_elem = section.find('table')
            if table_elem:
                self._table_builder.create_from_html(table_elem, header_text)
    
    def _process_main_content(self, soup: BeautifulSoup) -> None:
        """Process main content"""
        content_container = soup.find('div', class_='content-container')
        if not content_container:
            return
        
        gene_title_elem = content_container.find('h1', class_='gene-title')
        main_title = gene_title_elem.get_text(strip=True) if gene_title_elem else "Gene Analysis"
        
        # Find sequence-section (Detailed Results section)
        sequence_section = content_container.find('div', class_='sequence-section')
        
        # Create gene-section ID list inside sequence-section (for duplicate prevention)
        sequence_section_ids = set()
        if sequence_section:
            for gs in sequence_section.find_all('div', class_='gene-section'):
                sequence_section_ids.add(id(gs))
        
        gene_sections = content_container.find_all('div', class_='gene-section')
        seq_viewer_index = 0
        
        for idx, gene_section in enumerate(gene_sections, 1):
            # Skip gene-sections inside sequence-section (processed separately later)
            if id(gene_section) in sequence_section_ids:
                continue
            
            section_title = self._get_section_title(gene_section, idx)
            
            # Check if this is an Evidence section (starts with number)
            if section_title and section_title[0].isdigit() and '.' in section_title:
                # Evidence sections are processed in _process_evidence_section
                continue
            
            # Capture SeqViewerApp screenshots
            seq_viewers = gene_section.find_all('div', class_='SeqViewerApp', recursive=False)
            for sv_idx, _ in enumerate(seq_viewers):
                viewer_title = f"{section_title} - Sequence Viewer"
                if len(seq_viewers) > 1:
                    viewer_title += f" ({sv_idx + 1})"
                self._capture_element_screenshot('.SeqViewerApp', viewer_title, seq_viewer_index)
                seq_viewer_index += 1
            
            # Image processing
            self._process_images(gene_section, section_title, main_title)
            
            # Table processing
            self._process_tables(gene_section, section_title, main_title)
            
            # Subsection processing
            self._process_subsections(gene_section, main_title)
        
        # Process standalone h3 sections outside gene-section (e.g., 3.3, 3.4)
        self._process_standalone_h3_sections(content_container, main_title)
        
        # Process Detailed Results section (sequence-section)
        self._process_sequence_section(content_container)
        
        # Process Evidence tables (Source Summary)
        self._process_evidence_section(content_container)
    
    def _get_section_title(self, gene_section: Tag, default_idx: int) -> str:
        """Extract section title"""
        # Check h2.subsection-title
        subsection_title = gene_section.find('h2', class_='subsection-title')
        if subsection_title:
            return subsection_title.get_text(strip=True)
        
        # Check h3.subsection-title
        h3_title = gene_section.find('h3', class_='subsection-title')
        if h3_title:
            return h3_title.get_text(strip=True)
        
        # Check general h3
        h3_elem = gene_section.find('h3')
        if h3_elem:
            return h3_elem.get_text(strip=True)
        
        # Check general h2
        h2_elem = gene_section.find('h2')
        if h2_elem:
            return h2_elem.get_text(strip=True)
        
        # Default value
        if default_idx > 0:
            return f"Section {default_idx}"
        return "Section"
    
    def _process_images(self, gene_section: Tag, section_title: str, main_title: str) -> None:
        """Process images"""
        image_placeholder = gene_section.find('div', class_='image-placeholder')
        if image_placeholder:
            img = image_placeholder.find('img')
            if img and img.get('src', '').startswith('data:image'):
                self._image_builder.create_from_base64(img, section_title)
    
    def _process_tables(self, gene_section: Tag, section_title: str, main_title: str) -> None:
        """Process tables (dynamic grouping)"""
        tables = gene_section.find_all('table', class_='data-table')
        if not tables:
            return
        
        table_infos = [
            {'table': t, 'rows': len(t.find_all('tr'))} 
            for t in tables
        ]
        
        # Dynamic grouping
        title_space = Inches(0.85)
        table_gap = Inches(0.3)
        available_height = (
            self.slide_config.height - 
            self.slide_config.margin_top - 
            self.slide_config.margin_bottom - 
            title_space
        )
        row_height = Inches(0.28)
        
        i = 0
        while i < len(table_infos):
            current_group = [table_infos[i]]
            current_height = row_height * table_infos[i]['rows']
            
            while i + 1 < len(table_infos):
                next_rows = table_infos[i + 1]['rows']
                next_table_height = row_height * next_rows + table_gap
                remaining_space = available_height - current_height
                
                if next_table_height <= remaining_space:
                    i += 1
                    current_group.append(table_infos[i])
                    current_height += next_table_height
                else:
                    break
            
            if len(current_group) == 1:
                self._table_builder.create_from_html(
                    current_group[0]['table'], section_title, main_title
                )
            else:
                self._create_combined_table_slide(
                    [info['table'] for info in current_group],
                    section_title, main_title
                )
            
            i += 1
    
    def _create_combined_table_slide(
        self, 
        tables: List[Tag], 
        section_title: str, 
        main_title: str
    ) -> None:
        """Combine multiple tables into a single slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add title
        if main_title:
            main_box = slide.shapes.add_textbox(
                self.slide_config.margin_left, Inches(0.1),
                self.slide_config.content_width, Inches(0.25)
            )
            main_frame = main_box.text_frame
            main_frame.text = main_title
            main_frame.paragraphs[0].font.size = Pt(12)
            main_frame.paragraphs[0].font.color.rgb = self.colors['gray_600']
        
        title_top = Inches(0.35) if main_title else self.slide_config.margin_top - Inches(0.1)
        title_box = slide.shapes.add_textbox(
            self.slide_config.margin_left, title_top,
            self.slide_config.content_width, Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        # Add tables
        current_top = title_top + Inches(0.5)
        available_height = self.slide_config.height - current_top - self.slide_config.margin_bottom
        table_builder = TableBuilder(colors=self.colors)
        
        for table_idx, table_elem in enumerate(tables):
            extractor = TableDataExtractor(table_elem).extract()
            rows = len(extractor.rows_data)
            
            table_height = min(Inches(0.25) * rows, available_height * 0.4)
            
            if table_idx > 0:
                current_top += Inches(0.2)
            
            table_builder.create_table(
                slide,
                extractor.rows_data,
                len(extractor.header_rows),
                extractor.col_widths_html,
                self.slide_config.margin_left, current_top,
                self.slide_config.content_width, table_height,
                extractor.merge_info,
                extractor.cell_styles
            )
            
            current_top += table_height + Inches(0.1)
    
    def _process_subsections(self, gene_section: Tag, main_title: str) -> None:
        """Process subsections (h3) - Option A: All h3 titles and tables in one slide"""
        h3_sections = gene_section.find_all('h3')
        if not h3_sections:
            return
        
        # Collect h3 and related tables
        h3_table_pairs = []
        for h3 in h3_sections:
            h3_title = h3.get_text(strip=True)
            next_table = h3.find_next('table')
            if next_table:
            # Check if table is within current gene_section
            table_parent = next_table.parent
            while table_parent and table_parent != gene_section:
                table_parent = table_parent.parent
            if table_parent == gene_section:
                h3_table_pairs.append({'h3_title': h3_title, 'table': next_table})
        
        if not h3_table_pairs:
            return
        
        # Display all h3 and tables in one slide (Option A)
        self._create_h3_combined_slide(h3_table_pairs, main_title)
    
    def _create_h3_combined_slide(
        self,
        h3_table_pairs: List[dict],
        main_title: str
    ) -> None:
        """Combine h3 titles and tables into one slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Main title
        if main_title:
            main_box = slide.shapes.add_textbox(
                self.slide_config.margin_left, Inches(0.1),
                self.slide_config.content_width, Inches(0.25)
            )
            main_frame = main_box.text_frame
            main_frame.text = main_title
            main_frame.paragraphs[0].font.size = Pt(12)
            main_frame.paragraphs[0].font.color.rgb = self.colors['gray_600']
        
        current_top = Inches(0.35) if main_title else self.slide_config.margin_top
        table_builder = TableBuilder(colors=self.colors)
        
        # Add each h3+table
        for idx, pair in enumerate(h3_table_pairs):
            h3_title = pair['h3_title']
            table_elem = pair['table']
            
            # Add h3 title
            h3_box = slide.shapes.add_textbox(
                self.slide_config.margin_left, current_top,
                self.slide_config.content_width, Inches(0.3)
            )
            h3_frame = h3_box.text_frame
            h3_frame.text = h3_title
            h3_para = h3_frame.paragraphs[0]
            h3_para.font.size = Pt(14)
            h3_para.font.bold = True
            h3_para.font.color.rgb = self.colors['primary_red']
            
            current_top += Inches(0.35)
            
            # Add table
            extractor = TableDataExtractor(table_elem).extract()
            rows = len(extractor.rows_data)
            
            # Calculate remaining space
            remaining_height = self.slide_config.height - current_top - self.slide_config.margin_bottom
            # Maximum height limit
            table_height = min(Inches(0.25) * rows, remaining_height * 0.4)
            
            table_builder.create_table(
                slide,
                extractor.rows_data,
                len(extractor.header_rows),
                extractor.col_widths_html,
                self.slide_config.margin_left, current_top,
                self.slide_config.content_width, table_height,
                extractor.merge_info,
                extractor.cell_styles
            )
            
            current_top += table_height + Inches(0.3)
    
    def _process_standalone_h3_sections(self, content_container: Tag, main_title: str) -> None:
        """
        Process standalone h3 sections outside gene-section (e.g., sections 3.3, 3.4)
        Find h3 elements that are direct children of content-container
        """
        # Find h3 elements that are direct children of content-container
        standalone_h3_elements = content_container.find_all('h3', recursive=False)
        
        for h3 in standalone_h3_elements:
            h3_title = h3.get_text(strip=True)
            
            # Collect elements following h3
            next_sibling = h3.find_next_sibling()
            
            if next_sibling:
                # Table processing (e.g., 3.3 Key Agency Recommendations)
                if next_sibling.name == 'table':
                    self._table_builder.create_from_html(next_sibling, h3_title, main_title)
                
                # Multiple reference cards (e.g., 3.4 Key Literature)
                elif next_sibling.name == 'div' and 'reference-card' in next_sibling.get('class', []):
                    # Create section separator slide
                    self._section_builder.create(h3_title)
                    
                    # Create individual slides for all reference-cards
                    current = next_sibling
                    while current and current.name == 'div' and 'reference-card' in current.get('class', []):
                        self._reference_builder.create(current, h3_title)
                        current = current.find_next_sibling()
    
    def _process_sequence_section(self, content_container: Tag) -> None:
        """
        Process Detailed Results of the AI-based sequence analysis section
        Find div with sequence-section class
        """
        # Find sequence-section
        sequence_section = content_container.find('div', class_='sequence-section')
        if not sequence_section:
            return
        
        # Find h2.sequence-title
        sequence_title_elem = sequence_section.find('h2', class_='sequence-title')
        if sequence_title_elem:
            sequence_title = sequence_title_elem.get_text(strip=True)
        else:
            sequence_title = "Detailed Results of the AI-based sequence analysis"
        
        # Create title slide
        self._title_builder.create(sequence_title, "AI-based sequence analysis detailed results")
        
        # Process gene-sections inside sequence-section
        gene_sections = sequence_section.find_all('div', class_='gene-section')
        seq_viewer_index = 0
        
        for gene_section in gene_sections:
            section_title = self._get_section_title(gene_section, 0)
            
            # Capture SeqViewerApp screenshots
            seq_viewers = gene_section.find_all('div', class_='SeqViewerApp', recursive=False)
            for sv_idx, _ in enumerate(seq_viewers):
                viewer_title = f"{section_title} - Sequence Viewer"
                if len(seq_viewers) > 1:
                    viewer_title += f" ({sv_idx + 1})"
                self._capture_element_screenshot('.SeqViewerApp', viewer_title, seq_viewer_index)
                seq_viewer_index += 1
            
            # Image processing
            self._process_images(gene_section, section_title, sequence_title)
            
            # Process tables and background-text in HTML order
            tables = gene_section.find_all('table', class_='data-table')
            background_texts = gene_section.find_all('div', class_='background-text')
            
            if tables and not background_texts:
                # Only tables
                for table in tables:
                    self._table_builder.create_from_html(table, section_title, sequence_title)
            elif background_texts and not tables:
                # Only background-text
                for bg_text in background_texts:
                    text_content = bg_text.get_text(strip=True)
                    if text_content:
                        self._content_builder.create_with_text(section_title, text_content, sequence_title)
            elif tables and background_texts:
                # Both: include background-text in table slide
                for table in tables:
                    # Use background-text as subtitle
                    bg_text = background_texts[0].get_text(strip=True) if background_texts else ""
                    self._table_builder.create_from_html(table, section_title, sequence_title)
    
    def _process_evidence_section(self, content_container: Tag) -> None:
        """
        Process Source Summary section (Evidence tables)
        Process Evidence sections starting with h2.subsection-title
        """
        # Find Evidence-related h2 sections
        evidence_titles = [
            '1. gene name', '2. protein name', '3. gene_synonyms',
            '4. related gene name', '5. feature reference', '6. gc content reference',
            '7. copy number reference', '8. application usage reference',
            '9. mutation variation reference', '10. species discrimination reference',
            '11. pros cons reference', '12. product institution reference'
        ]
        
        all_subsections = content_container.find_all('h2', class_='subsection-title')
        evidence_subsections = []
        
        for subsection in all_subsections:
            section_title = subsection.get_text(strip=True)
            # Check if this is an Evidence-related section (starts with number or contains specific keywords)
            if any(section_title.lower().startswith(t.split('.')[0] + '.') for t in evidence_titles) or \
               'reference' in section_title.lower() or \
               section_title.lower().startswith(tuple(str(i) + '.' for i in range(1, 13))):
                evidence_subsections.append(subsection)
        
        if not evidence_subsections:
            # Process Evidence tables using existing method
            self._process_evidence_tables_legacy(content_container)
            return
        
        # Create Source Summary section slide
        self._section_builder.create("Source Summary", "Evidence References")
        
        # Process each Evidence section
        for subsection in evidence_subsections:
            section_title = subsection.get_text(strip=True)
            
            next_elem = subsection.find_next_sibling()
            while next_elem:
                if next_elem.name == 'h2':
                    break
                if next_elem.get('class') and 'gene-section' in next_elem.get('class', []):
                    break
                
                if next_elem.name == 'div' and 'evidence-table' in next_elem.get('class', []):
                    self._evidence_builder.create(next_elem, section_title)
                
                next_elem = next_elem.find_next_sibling()
    
    def _process_evidence_tables_legacy(self, content_container: Tag) -> None:
        """Process Evidence tables (legacy method)"""
        all_subsections = content_container.find_all('h2', class_='subsection-title')
        
        for subsection in all_subsections:
            section_title = subsection.get_text(strip=True)
            
            next_elem = subsection.find_next_sibling()
            while next_elem:
                if next_elem.name == 'h2':
                    break
                if next_elem.get('class') and 'gene-section' in next_elem.get('class', []):
                    break
                
                if next_elem.name == 'div' and 'evidence-table' in next_elem.get('class', []):
                    self._evidence_builder.create(next_elem, section_title)
                
                next_elem = next_elem.find_next_sibling()
    
    def _capture_element_screenshot(self, selector: str, title: str, index: int = 0) -> None:
        """Capture HTML element screenshot using Playwright"""
        if not self.html_path:
            logger.warning(f"Cannot capture screenshot: HTML path not set")
            return
        
        try:
            from playwright.sync_api import sync_playwright
            
            logger.info(f"Capturing element with Playwright: {selector} (index={index})")
            
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page(viewport={'width': 1400, 'height': 900})
                
                file_url = f"file://{self.html_path}"
                page.goto(file_url, wait_until='networkidle')
                page.wait_for_timeout(2000)
                
                elements = page.locator(selector)
                count = elements.count()
                
                if count == 0 or index >= count:
                    browser.close()
                    return
                
                element = elements.nth(index)
                element.scroll_into_view_if_needed()
                page.wait_for_timeout(500)
                
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    screenshot_path = tmp_file.name
                
                element.screenshot(path=screenshot_path)
                browser.close()
                
                self._image_builder.create_from_file(screenshot_path, title)
                os.unlink(screenshot_path)
                
        except ImportError:
            logger.error("Playwright is not installed.")
        except Exception as e:
            logger.error(f"Failed to capture element screenshot: {e}")


def convert_html_to_pptx(html_path: Path, output_path: Path) -> None:
    """
    Convenience function to convert HTML file to PPTX
    
    Args:
        html_path: Input HTML file path
        output_path: Output PPTX file path
    """
    converter = HtmlToPptxConverter()
    converter.convert(html_path, output_path)


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python converter.py <input.html> <output.pptx>")
        sys.exit(1)
    
    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    
    convert_html_to_pptx(input_path, output_path)
    print(f"Conversion complete: {output_path}")
