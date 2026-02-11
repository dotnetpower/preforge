"""
Table creation and styling module

Handles PowerPoint table creation, border styling, column width adjustment, etc.
"""
import logging
from typing import List, Optional, Dict, Any, Tuple
from bs4 import Tag
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
from pptx.oxml.ns import qn

from .config import (
    DEFAULT_TABLE_CONFIG, 
    DEFAULT_BORDER_CONFIG, 
    DEFAULT_COLORS,
    TableConfig,
    BorderConfig,
    ColorPalette
)
from .style_utils import StyleExtractor, TextUtils

logger = logging.getLogger(__name__)


class TableDataExtractor:
    """Class that extracts data from HTML tables"""
    
    def __init__(self, table_elem: Tag):
        self.table_elem = table_elem
        self.rows_data: List[List[str]] = []
        self.header_rows: List[List[str]] = []
        self.body_rows: List[List[str]] = []
        self.col_widths_html: List[Optional[int]] = []
        self.merge_info: List[Tuple[int, int, int, int]] = []  # (row, col, colspan, rowspan)
        self.cell_styles: Dict[Tuple[int, int], Dict[str, Any]] = {}
        self.has_header = False
        self.max_cols = 0
    
    def extract(self) -> 'TableDataExtractor':
        """Extract table data"""
        thead = self.table_elem.find('thead')
        tbody = self.table_elem.find('tbody')
        
        # Process thead
        if thead:
            self.has_header = True
            header_trs = thead.find_all('tr')
            for tr in header_trs:
                row_data = self._extract_row_data(tr, len(self.rows_data))
                self.header_rows.append(row_data)
                self.rows_data.append(row_data)
                
                if not self.col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    self.col_widths_html = StyleExtractor.extract_column_widths(cells)
        
        # Process tbody
        if tbody:
            body_trs = tbody.find_all('tr')
            for idx, tr in enumerate(body_trs):
                row_data = self._extract_row_data(tr, len(self.rows_data))
                self.body_rows.append(row_data)
                self.rows_data.append(row_data)
                
                if not self.has_header and idx == 0 and not self.col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    self.col_widths_html = StyleExtractor.extract_column_widths(cells)
        
        # If neither thead nor tbody exists (use tr directly)
        if not self.has_header and not tbody:
            all_rows = self.table_elem.find_all('tr')
            for idx, tr in enumerate(all_rows):
                row_data = self._extract_row_data(tr, len(self.rows_data))
                self.body_rows.append(row_data)
                self.rows_data.append(row_data)
                
                if idx == 0 and not self.col_widths_html:
                    cells = tr.find_all(['th', 'td'])
                    self.col_widths_html = StyleExtractor.extract_column_widths(cells)
        
        # Determine and normalize column count
        if self.rows_data:
            self.max_cols = max(len(row) for row in self.rows_data)
            for row in self.rows_data:
                while len(row) < self.max_cols:
                    row.append("")
        
        return self
    
    def _extract_row_data(self, tr: Tag, row_idx: int) -> List[str]:
        """Extract row data (including colspan handling)"""
        cells = tr.find_all(['th', 'td'])
        row_data = []
        col_idx = 0
        
        for cell in cells:
            # Extract text while preserving bullets and line breaks
            text = TextUtils.extract_cell_text_with_formatting(cell)
            colspan = int(cell.get('colspan', 1))
            rowspan = int(cell.get('rowspan', 1))
            
            # Extract styles
            styles = StyleExtractor.extract_cell_styles(cell)
            if styles['bold'] or styles['color'] or styles['link']:
                self.cell_styles[(row_idx, col_idx)] = styles
            
            row_data.append(text)
            
            # Add empty cells if colspan exists
            for _ in range(colspan - 1):
                row_data.append('')
            
            # Save merge information
            if colspan > 1 or rowspan > 1:
                self.merge_info.append((row_idx, col_idx, colspan, rowspan))
            
            col_idx += colspan
        
        return row_data
    
    def is_key_value_table(self) -> bool:
        """Check if the table is a key-value table"""
        if self.has_header:
            return False
        if len(self.body_rows) > 5:
            return False
        if not self.body_rows:
            return False
        
        first_row = self.body_rows[0]
        return len(first_row) == 2


class TableBorderStyler:
    """Class that applies table border styles"""
    
    def __init__(
        self, 
        border_config: BorderConfig = None,
        colors: ColorPalette = None
    ):
        self.border_config = border_config or DEFAULT_BORDER_CONFIG
        self.colors = colors or DEFAULT_COLORS
    
    def apply_academic_borders(
        self, 
        ppt_table, 
        header_count: int, 
        row_count: int, 
        col_count: int
    ) -> None:
        """Apply academic paper style borders (thick lines at top/bottom, thick line below header)"""
        thick_line = self.border_config.thick_line
        thin_line = self.border_config.thin_line
        no_line = self.border_config.no_line
        
        black = self.colors['black']
        gray_line = self.colors['gray_line']
        
        for i in range(row_count):
            for j in range(col_count):
                try:
                    cell = ppt_table.cell(i, j)
                    
                    # Top line
                    if i == 0:
                        self._set_cell_border(cell, 'top', thick_line, black)
                    elif i == header_count and header_count > 0:
                        self._set_cell_border(cell, 'top', no_line, black)
                    else:
                        self._set_cell_border(cell, 'top', thin_line, gray_line)
                    
                    # Bottom line
                    if i == row_count - 1:
                        self._set_cell_border(cell, 'bottom', thick_line, black)
                    elif i == header_count - 1 and header_count > 0:
                        self._set_cell_border(cell, 'bottom', thick_line, black)
                    else:
                        self._set_cell_border(cell, 'bottom', thin_line, gray_line)
                    
                    # No left/right borders
                    self._set_cell_border(cell, 'left', no_line, black)
                    self._set_cell_border(cell, 'right', no_line, black)
                    
                except Exception:
                    pass
    
    def _set_cell_border(self, cell, side: str, width: int, color: RGBColor) -> None:
        """Set specific border of a cell"""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        border_map = {
            'top': 'a:lnT',
            'bottom': 'a:lnB',
            'left': 'a:lnL',
            'right': 'a:lnR'
        }
        
        border_elem_name = border_map.get(side)
        if not border_elem_name:
            return
        
        # Remove existing border elements
        for existing in list(tcPr):
            if existing.tag == qn(border_elem_name):
                tcPr.remove(existing)
        
        width_emu = int(width) if width > 0 else 0
        
        # Create new border element
        ln = etree.Element(qn(border_elem_name))
        
        if width_emu > 0:
            ln.set('w', str(width_emu))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')
            
            solidFill = etree.SubElement(ln, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
            
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'solid')
        else:
            ln.set('w', '0')
            etree.SubElement(ln, qn('a:noFill'))
        
        tcPr.insert(0, ln)


class TableColumnAdjuster:
    """Class that adjusts table column widths"""
    
    @staticmethod
    def apply_html_widths(
        ppt_table, 
        col_widths_html: List[Optional[int]], 
        total_width: float
    ) -> None:
        """Apply width attributes extracted from HTML"""
        try:
            col_count = len(col_widths_html)
            if col_count == 0:
                return
            
            specified_widths = [w for w in col_widths_html if w is not None]
            unspecified_count = col_widths_html.count(None)
            
            if not specified_widths:
                return
            
            html_to_ppt_ratio = total_width / 900
            
            specified_total_ppt = sum(
                int(w * html_to_ppt_ratio) 
                for w in col_widths_html if w is not None
            )
            
            remaining_width = total_width - specified_total_ppt
            
            if remaining_width < 0 or (
                unspecified_count > 0 and remaining_width < total_width * 0.3
            ):
                specified_portion = 0.3 if unspecified_count > 0 else 1.0
                total_specified_html = sum(specified_widths)
                
                for j, html_width in enumerate(col_widths_html):
                    if html_width is not None:
                        proportion = html_width / total_specified_html
                        ppt_table.columns[j].width = int(
                            total_width * specified_portion * proportion
                        )
                
                if unspecified_count > 0:
                    remaining = total_width * (1 - specified_portion)
                    equal_width = int(remaining / unspecified_count)
                    for j, html_width in enumerate(col_widths_html):
                        if html_width is None:
                            ppt_table.columns[j].width = equal_width
            else:
                for j, html_width in enumerate(col_widths_html):
                    if html_width is not None:
                        ppt_table.columns[j].width = int(html_width * html_to_ppt_ratio)
                
                if unspecified_count > 0:
                    equal_width = int(remaining_width / unspecified_count)
                    for j, html_width in enumerate(col_widths_html):
                        if html_width is None:
                            ppt_table.columns[j].width = equal_width
        
        except Exception as e:
            logger.debug(f"Failed to apply HTML width: {e}")
    
    @staticmethod
    def auto_adjust(ppt_table, rows_data: List[List[str]]) -> None:
        """Auto-adjust column widths based on text length"""
        try:
            col_count = len(rows_data[0]) if rows_data else 0
            if col_count == 0:
                return
            
            total_table_width = sum(col.width for col in ppt_table.columns)
            
            max_lengths = [0] * col_count
            for row in rows_data:
                for j, cell in enumerate(row):
                    cell_text = str(cell)
                    korean_count = len([
                        c for c in cell_text 
                        if ord(c) >= 0xAC00 and ord(c) <= 0xD7A3
                    ])
                    english_count = len(cell_text) - korean_count
                    weighted_length = english_count + (korean_count * 1.8)
                    max_lengths[j] = max(max_lengths[j], weighted_length)
            
            min_proportion = 0.05
            total_length = sum(max_lengths)
            
            if total_length == 0:
                equal_width = total_table_width // col_count
                for j in range(col_count):
                    ppt_table.columns[j].width = equal_width
                return
            
            for j in range(col_count):
                proportion = max(max_lengths[j] / total_length, min_proportion)
                ppt_table.columns[j].width = int(total_table_width * proportion)
        
        except Exception as e:
            logger.debug(f"Failed to adjust column widths: {e}")


class TableBuilder:
    """Class that creates PowerPoint tables"""
    
    def __init__(
        self,
        table_config: TableConfig = None,
        colors: ColorPalette = None
    ):
        self.table_config = table_config or DEFAULT_TABLE_CONFIG
        self.colors = colors or DEFAULT_COLORS
        self.border_styler = TableBorderStyler(colors=colors)
    
    def create_table(
        self,
        slide,
        rows_data: List[List[str]],
        header_count: int,
        col_widths_html: List[Optional[int]],
        left: float,
        top: float,
        width: float,
        height: float,
        merge_info: List[Tuple[int, int, int, int]] = None,
        cell_styles: Dict[Tuple[int, int], Dict[str, Any]] = None
    ):
        """Create PowerPoint table"""
        if merge_info is None:
            merge_info = []
        if cell_styles is None:
            cell_styles = {}
        
        if not rows_data:
            return None
        
        max_cols = len(rows_data[0])
        row_count = len(rows_data)
        
        # Determine font size
        if row_count > 20 or max_cols > 6:
            base_font_size = Pt(7)
            header_font_size = Pt(8)
        elif row_count > 15:
            base_font_size = Pt(8)
            header_font_size = Pt(9)
        else:
            base_font_size = self.table_config.body_font_size
            header_font_size = self.table_config.header_font_size
        
        # Calculate table height
        min_row_height = self.table_config.min_row_height
        required_height = min_row_height * row_count
        height = min(required_height, height)
        
        try:
            ppt_table = slide.shapes.add_table(
                row_count, max_cols,
                left, top, width, height
            ).table
            
            # Fill data
            for i, row_data in enumerate(rows_data):
                for j, cell_data in enumerate(row_data):
                    if j >= max_cols:
                        continue
                    
                    cell = ppt_table.cell(i, j)
                    cell.text = str(cell_data) if j < len(row_data) else ""
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    cell.margin_left = self.table_config.cell_margin
                    cell.margin_right = self.table_config.cell_margin
                    cell.margin_top = self.table_config.cell_margin_vertical
                    cell.margin_bottom = self.table_config.cell_margin_vertical
                    
                    cell.fill.background()
                    
                    html_style = cell_styles.get((i, j), {})
                    has_custom_bold = html_style.get('bold', False)
                    custom_color = html_style.get('color')
                    has_link = html_style.get('link')
                    
                    for paragraph in cell.text_frame.paragraphs:
                        if i < header_count:
                            paragraph.font.size = header_font_size
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = self.colors['black']
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.text_frame.word_wrap = False
                        else:
                            paragraph.font.size = base_font_size
                            
                            if has_custom_bold:
                                paragraph.font.bold = True
                            
                            if custom_color:
                                paragraph.font.color.rgb = custom_color
                            else:
                                paragraph.font.color.rgb = self.colors['gray_800']
                            
                            if has_link:
                                paragraph.font.color.rgb = self.colors['link_blue']
                                paragraph.font.underline = True
                            
                            # Left-align if bullet(•) present, otherwise center-align
                            if '•' in cell_data or '\n' in cell_data:
                                paragraph.alignment = PP_ALIGN.LEFT
                            else:
                                paragraph.alignment = PP_ALIGN.CENTER
                            
                            cell.text_frame.word_wrap = True
                        
                        paragraph.line_spacing = 1.1
            
            # Apply borders
            self.border_styler.apply_academic_borders(
                ppt_table, header_count, row_count, max_cols
            )
            
            # Apply cell merge
            for row_idx, col_idx, colspan, rowspan in merge_info:
                try:
                    if row_idx < row_count and col_idx < max_cols:
                        start_cell = ppt_table.cell(row_idx, col_idx)
                        end_row = min(row_idx + rowspan - 1, row_count - 1)
                        end_col = min(col_idx + colspan - 1, max_cols - 1)
                        
                        end_cell = ppt_table.cell(end_row, end_col)
                        start_cell.merge(end_cell)
                        
                        for paragraph in start_cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                except Exception:
                    pass
            
            # Adjust column widths
            if col_widths_html and any(w is not None for w in col_widths_html):
                TableColumnAdjuster.apply_html_widths(ppt_table, col_widths_html, width)
            else:
                TableColumnAdjuster.auto_adjust(ppt_table, rows_data)
            
            return ppt_table
            
        except Exception as e:
            logger.error(f"Failed to create table: {e}")
            return None
