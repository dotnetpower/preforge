"""
슬라이드 생성 팩토리 모듈

다양한 유형의 슬라이드(타이틀, 테이블, 이미지 등)를 생성하는 기능을 제공합니다.
"""
import logging
import base64
from io import BytesIO
from typing import List, Optional, Any
from bs4 import Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

from .config import SlideConfig, ColorPalette, DEFAULT_SLIDE_CONFIG, DEFAULT_COLORS
from .table_builder import TableBuilder, TableDataExtractor
from .style_utils import TextUtils

logger = logging.getLogger(__name__)


class SlideFactory:
    """슬라이드 생성 팩토리"""
    
    def __init__(
        self,
        presentation: Presentation,
        slide_config: SlideConfig = None,
        colors: ColorPalette = None
    ):
        self.prs = presentation
        self.config = slide_config or DEFAULT_SLIDE_CONFIG
        self.colors = colors or DEFAULT_COLORS
        self.table_builder = TableBuilder(colors=colors)
    
    def _get_blank_slide(self):
        """빈 레이아웃 슬라이드 생성"""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])
    
    def _add_title(
        self, 
        slide, 
        text: str, 
        font_size: int = Pt(20),
        top: float = None,
        bold: bool = True,
        color: RGBColor = None
    ) -> float:
        """슬라이드에 제목 추가"""
        top = top if top is not None else self.config.margin_top - Inches(0.2)
        color = color or self.colors['primary_red']
        
        title_box = slide.shapes.add_textbox(
            self.config.margin_left, top,
            self.config.content_width, Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = font_size
        title_para.font.bold = bold
        title_para.font.color.rgb = color
        
        return top + Inches(0.5)
    
    def _add_subtitle(
        self, 
        slide, 
        text: str, 
        font_size: int = Pt(12),
        top: float = None,
        color: RGBColor = None
    ) -> float:
        """슬라이드에 부제목 추가"""
        top = top if top is not None else Inches(0.1)
        color = color or self.colors['gray_600']
        
        subtitle_box = slide.shapes.add_textbox(
            self.config.margin_left, top,
            self.config.content_width, Inches(0.25)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = font_size
        subtitle_para.font.color.rgb = color
        
        return top + Inches(0.3)


class TitleSlideBuilder(SlideFactory):
    """타이틀 슬라이드 생성"""
    
    def create(self, title: str, subtitle: str = "") -> Any:
        """타이틀 슬라이드 생성"""
        slide = self._get_blank_slide()
        
        # 배경
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.config.width, self.config.height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['gray_50']
        background.line.fill.background()
        
        # 상단 빨간 박스
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(2),
            self.config.width, Inches(1.5)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['primary_red']
        header_box.line.fill.background()
        
        # 타이틀 텍스트
        title_frame = header_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']
        
        # 부제목
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4),
                self.config.width - Inches(2), Inches(1.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = self.colors['gray_800']
            subtitle_frame.word_wrap = True
        
        return slide


class ContentSlideBuilder(SlideFactory):
    """일반 컨텐츠 슬라이드 생성"""
    
    def create_with_text(
        self, 
        title: str, 
        content: str,
        subtitle: str = None
    ) -> Any:
        """텍스트 컨텐츠 슬라이드 생성"""
        slide = self._get_blank_slide()
        
        y_position = self.config.margin_top - Inches(0.2)
        
        # 제목
        y_position = self._add_title(slide, title, Pt(32), y_position)
        
        # 부제목
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.config.margin_left, y_position,
                self.config.content_width, Inches(0.3)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.bold = True
            subtitle_para.font.color.rgb = self.colors['gray_800']
            y_position += Inches(0.4)
        
        # 본문
        text_box = slide.shapes.add_textbox(
            self.config.margin_left, y_position + Inches(0.2),
            self.config.content_width, Inches(4)
        )
        text_frame = text_box.text_frame
        text_frame.text = content
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = self.colors['gray_800']
            paragraph.line_spacing = 1.5
        
        return slide


class TableSlideBuilder(SlideFactory):
    """테이블 슬라이드 생성"""
    
    def __init__(
        self,
        presentation: Presentation,
        slide_config: SlideConfig = None,
        colors: ColorPalette = None,
        max_rows_per_slide: int = 8
    ):
        super().__init__(presentation, slide_config, colors)
        self.max_rows_per_slide = max_rows_per_slide
    
    def create_from_html(
        self, 
        table_elem: Tag, 
        title: str,
        main_title: str = ""
    ) -> List[Any]:
        """HTML 테이블에서 슬라이드 생성 (자동 분할)"""
        # 테이블 데이터 먼저 추출
        extractor = TableDataExtractor(table_elem).extract()
        
        # 분할이 필요한 경우 - 슬라이드 생성 전에 체크
        header_count = len(extractor.header_rows)
        body_count = len(extractor.body_rows)
        
        y_position = self.config.margin_top - Inches(0.2)
        if main_title:
            y_position = Inches(0.35)
        table_top = y_position + Inches(0.5)
        table_height = self.config.height - table_top - self.config.margin_bottom
        
        if body_count > self.max_rows_per_slide:
            # 분할 시 별도 슬라이드들 생성
            return self._create_split_table_slides(
                extractor, title, main_title, table_top, table_height
            )
        
        # 단일 슬라이드 생성
        slide = self._get_blank_slide()
        
        # 메인 타이틀
        if main_title:
            self._add_subtitle(slide, main_title, Pt(12), Inches(0.1))
        
        # 섹션 타이틀
        self._add_title(slide, title, Pt(18), y_position)
        
        # 키-값 테이블인 경우 카드 스타일로 표시
        if extractor.is_key_value_table():
            self._add_key_value_cards(
                slide, table_elem,
                self.config.margin_left, table_top,
                self.config.content_width, table_height
            )
            return [slide]
        
        # 단일 슬라이드
        self.table_builder.create_table(
            slide,
            extractor.rows_data,
            header_count,
            extractor.col_widths_html,
            self.config.margin_left, table_top,
            self.config.content_width, table_height,
            extractor.merge_info,
            extractor.cell_styles
        )
        
        return [slide]
    
    def _create_split_table_slides(
        self,
        extractor: TableDataExtractor,
        title: str,
        main_title: str,
        table_top: float,
        table_height: float
    ) -> List[Any]:
        """분할된 테이블 슬라이드 생성"""
        slides = []
        header_count = len(extractor.header_rows)
        body_rows = extractor.body_rows
        
        num_chunks = (len(body_rows) + self.max_rows_per_slide - 1) // self.max_rows_per_slide
        
        for chunk_idx in range(num_chunks):
            start_idx = chunk_idx * self.max_rows_per_slide
            end_idx = min(start_idx + self.max_rows_per_slide, len(body_rows))
            
            chunk_data = extractor.header_rows + body_rows[start_idx:end_idx]
            
            # 해당 청크의 merge_info 필터링
            chunk_merge_info = []
            for row_idx, col_idx, colspan, rowspan in extractor.merge_info:
                if row_idx < header_count:
                    chunk_merge_info.append((row_idx, col_idx, colspan, rowspan))
                elif row_idx - header_count >= start_idx and row_idx - header_count < end_idx:
                    new_row_idx = header_count + (row_idx - header_count - start_idx)
                    chunk_merge_info.append((new_row_idx, col_idx, colspan, rowspan))
            
            # 해당 청크의 cell_styles 필터링
            chunk_cell_styles = {}
            for (r, c), styles in extractor.cell_styles.items():
                if r < header_count:
                    chunk_cell_styles[(r, c)] = styles
                elif r - header_count >= start_idx and r - header_count < end_idx:
                    new_r = header_count + (r - header_count - start_idx)
                    chunk_cell_styles[(new_r, c)] = styles
            
            # 슬라이드 생성
            slide = self._get_blank_slide()
            
            # 제목 (첫 슬라이드 이후에는 "(계속 N)" 추가)
            slide_title = title if chunk_idx == 0 else f"{title} (계속 {chunk_idx + 1})"
            self._add_title(slide, slide_title, Pt(18), Inches(0.1))
            
            # 테이블 생성
            self.table_builder.create_table(
                slide,
                chunk_data,
                header_count,
                extractor.col_widths_html,
                self.config.margin_left, Inches(0.6),
                self.config.content_width, self.config.height - Inches(0.9),
                chunk_merge_info,
                chunk_cell_styles
            )
            
            slides.append(slide)
        
        return slides
    
    def _add_key_value_cards(
        self,
        slide,
        table_elem: Tag,
        left: float,
        top: float,
        width: float,
        height: float
    ) -> List[Any]:
        """키-값 형태의 테이블을 카드 스타일로 표시"""
        tbody = table_elem.find('tbody')
        if not tbody:
            return []
        
        rows = tbody.find_all('tr')
        if not rows:
            return []
        
        label_bg_color = RGBColor(55, 65, 81)
        value_bg_color = RGBColor(249, 250, 251)
        border_color = RGBColor(209, 213, 219)
        
        y_position = top
        card_height = Inches(1.3)
        card_spacing = Inches(0.12)
        label_width = Inches(1.6)
        
        for tr in rows:
            cells = tr.find_all(['th', 'td'])
            if len(cells) < 2:
                continue
            
            label = TextUtils.clean_text(cells[0].get_text(strip=True))
            value = TextUtils.clean_text(cells[1].get_text(strip=True))
            
            # 라벨 영역
            label_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, y_position,
                label_width, card_height
            )
            label_shape.fill.solid()
            label_shape.fill.fore_color.rgb = label_bg_color
            label_shape.line.fill.background()
            
            label_tf = label_shape.text_frame
            label_tf.word_wrap = True
            label_tf.paragraphs[0].text = label
            label_tf.paragraphs[0].font.size = Pt(13)
            label_tf.paragraphs[0].font.bold = True
            label_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            label_shape.text_frame.margin_left = Pt(10)
            label_shape.text_frame.margin_right = Pt(10)
            label_shape.text_frame.margin_top = Pt(10)
            label_shape.text_frame.margin_bottom = Pt(10)
            
            # 값 영역
            value_left = left + label_width
            value_width = width - label_width
            
            value_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                value_left, y_position,
                value_width, card_height
            )
            value_shape.fill.solid()
            value_shape.fill.fore_color.rgb = value_bg_color
            value_shape.line.color.rgb = border_color
            value_shape.line.width = Pt(1)
            
            value_tf = value_shape.text_frame
            value_tf.word_wrap = True
            value_tf.paragraphs[0].text = value
            value_tf.paragraphs[0].font.size = Pt(11)
            value_tf.paragraphs[0].font.color.rgb = self.colors['gray_800']
            value_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            value_shape.text_frame.margin_left = Pt(12)
            value_shape.text_frame.margin_right = Pt(12)
            value_shape.text_frame.margin_top = Pt(10)
            value_shape.text_frame.margin_bottom = Pt(10)
            
            y_position += card_height + card_spacing
        
        return [slide]


class ImageSlideBuilder(SlideFactory):
    """이미지 슬라이드 생성"""
    
    def create_from_base64(self, img_tag: Tag, title: str) -> Optional[Any]:
        """Base64 이미지에서 슬라이드 생성"""
        if not HAS_PIL:
            logger.warning("PIL이 설치되지 않아 이미지 슬라이드를 생성할 수 없습니다")
            return None
        
        src = img_tag.get('src', '')
        if not src.startswith('data:image'):
            return None
        
        try:
            header, data = src.split(',', 1)
            img_bytes = base64.b64decode(data)
            pil_img = Image.open(BytesIO(img_bytes))
            img_width, img_height = pil_img.size
            
            slide = self._get_blank_slide()
            
            # 제목
            self._add_title(slide, title or "Analysis Chart", Pt(20), Inches(0.2))
            
            # 이미지 크기 계산
            available_width = self.config.content_width
            available_height = self.config.height - Inches(1.2)
            
            img_ratio = img_width / img_height
            available_ratio = available_width / available_height
            
            if img_ratio > available_ratio:
                final_width = available_width
                final_height = available_width / img_ratio
            else:
                final_height = available_height
                final_width = available_height * img_ratio
            
            # 중앙 정렬
            img_left = self.config.margin_left + (available_width - final_width) / 2
            img_top = Inches(0.7) + (available_height - final_height) / 2
            
            # 이미지 저장 및 추가
            img_stream = BytesIO()
            pil_img.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            slide.shapes.add_picture(
                img_stream,
                img_left, img_top,
                final_width, final_height
            )
            
            logger.info(f"이미지 슬라이드 생성: {title} ({img_width}x{img_height})")
            return slide
            
        except Exception as e:
            logger.error(f"이미지 슬라이드 생성 실패: {e}")
            return None
    
    def create_from_file(self, image_path: str, title: str) -> Optional[Any]:
        """파일에서 이미지 슬라이드 생성"""
        if not HAS_PIL:
            return None
        
        try:
            pil_img = Image.open(image_path)
            img_width, img_height = pil_img.size
            
            slide = self._get_blank_slide()
            
            self._add_title(slide, title, Pt(20), Inches(0.2))
            
            available_width = self.config.content_width
            available_height = self.config.height - Inches(1.2)
            
            img_ratio = img_width / img_height
            available_ratio = available_width / available_height
            
            if img_ratio > available_ratio:
                final_width = available_width
                final_height = available_width / img_ratio
            else:
                final_height = available_height
                final_width = available_height * img_ratio
            
            img_left = self.config.margin_left + (available_width - final_width) / 2
            img_top = Inches(0.7) + (available_height - final_height) / 2
            
            slide.shapes.add_picture(
                image_path,
                img_left, img_top,
                final_width, final_height
            )
            
            return slide
            
        except Exception as e:
            logger.error(f"이미지 슬라이드 생성 실패: {e}")
            return None


class SectionSlideBuilder(SlideFactory):
    """섹션 구분 슬라이드 생성 (중간 타이틀)"""
    
    def create(self, title: str, subtitle: str = "") -> Any:
        """섹션 구분 슬라이드 생성"""
        slide = self._get_blank_slide()
        
        # 배경
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.config.width, self.config.height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['gray_50']
        background.line.fill.background()
        
        # 중앙 섹션 제목 박스
        title_height = Inches(1.2)
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, (self.config.height - title_height) / 2,
            self.config.width, title_height
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = self.colors['primary_red']
        title_box.line.fill.background()
        
        # 섹션 제목 텍스트
        title_tf = title_box.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_tf.paragraphs[0].font.size = Pt(32)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = self.colors['white']
        
        # 부제목
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), (self.config.height + title_height) / 2 + Inches(0.2),
                self.config.width - Inches(2), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.color.rgb = self.colors['gray_600']
            subtitle_frame.word_wrap = True
        
        return slide


class ReferenceCardSlideBuilder(SlideFactory):
    """Reference 카드 슬라이드 생성"""
    
    def create(self, reference_card: Tag, section_title: str = "") -> Optional[Any]:
        """Reference 카드에서 슬라이드 생성"""
        slide = self._get_blank_slide()
        
        y_position = self.config.margin_top - Inches(0.2)
        
        # 섹션 타이틀 (작은 글씨)
        if section_title:
            self._add_subtitle(slide, section_title, Pt(11), Inches(0.1))
            y_position = Inches(0.35)
        
        # Reference 제목 추출 (reference-number에서)
        ref_number = reference_card.find('div', class_='reference-number')
        if ref_number:
            ref_title = ref_number.get_text(strip=True)
        else:
            ref_title = "Reference"
        
        # 제목 (Reference 제목)
        title_box = slide.shapes.add_textbox(
            self.config.margin_left, y_position,
            self.config.content_width, Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = ref_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary_red']
        
        y_position += Inches(0.7)
        
        # 메타 정보 (저널, 날짜 등)
        ref_meta = reference_card.find('div', class_='reference-meta')
        if ref_meta:
            meta_text = ref_meta.get_text(strip=True)
            meta_box = slide.shapes.add_textbox(
                self.config.margin_left, y_position,
                self.config.content_width, Inches(0.3)
            )
            meta_frame = meta_box.text_frame
            meta_frame.text = meta_text
            meta_para = meta_frame.paragraphs[0]
            meta_para.font.size = Pt(10)
            meta_para.font.color.rgb = self.colors['gray_600']
            y_position += Inches(0.4)
        
        # 요약 내용
        ref_summary = reference_card.find('div', class_='reference-summary')
        if ref_summary:
            # 요약 텍스트 정리
            summary_text = ref_summary.get_text(separator='\n', strip=True)
            # 길이 제한
            if len(summary_text) > 1500:
                summary_text = summary_text[:1500] + "..."
            
            summary_box = slide.shapes.add_textbox(
                self.config.margin_left, y_position,
                self.config.content_width,
                self.config.height - y_position - self.config.margin_bottom
            )
            summary_frame = summary_box.text_frame
            summary_frame.word_wrap = True
            summary_frame.text = summary_text
            
            for paragraph in summary_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = self.colors['gray_800']
                paragraph.line_spacing = 1.3
        
        return slide


class EvidenceSlideBuilder(SlideFactory):
    """Evidence 테이블 슬라이드 생성"""
    
    def create(self, evidence_div: Tag, title: str) -> Optional[Any]:
        """Evidence 테이블 슬라이드 생성"""
        slide = self._get_blank_slide()
        
        self._add_title(slide, title, Pt(18))
        
        # evidence-row 내부 또는 직접 evidence-cell 찾기
        evidence_row = evidence_div.find('div', class_='evidence-row')
        if evidence_row:
            # evidence-row 내부의 evidence-cell들을 데이터 행으로 사용
            data_rows = evidence_row.find_all('div', class_='evidence-cell')
        else:
            # evidence-row가 없으면 직접 evidence-cell 찾기
            data_rows = evidence_div.find_all('div', class_='evidence-cell')
        
        if not data_rows:
            return slide
        
        max_rows = min(len(data_rows), 10)
        
        # 테이블 데이터 추출
        table_data = []
        link_data = []
        
        # 헤더 추출
        header_div = evidence_div.find('div', class_='evidence-header')
        if header_div:
            headers = [elem.strip() for elem in header_div.stripped_strings]
            if headers:
                table_data.append(headers[:8])
        
        # 데이터 행 추출
        for row_idx, row in enumerate(data_rows[:max_rows]):
            row_texts = []
            text_elements = row.find_all('div', class_='evidence-text')
            
            for col_idx, elem in enumerate(text_elements[:8]):
                link = elem.find('a')
                if link:
                    link_text = link.get_text(strip=True)
                    link_url = link.get('href', '')
                    row_texts.append(link_text)
                    if link_url:
                        link_data.append((len(table_data), col_idx, link_url))
                else:
                    text = TextUtils.clean_text(elem.get_text(strip=True))
                    text = TextUtils.truncate_text(text, 80)
                    row_texts.append(text)
            
            if row_texts:
                table_data.append(row_texts)
        
        if len(table_data) <= 1:
            return slide
        
        # 열 수 통일
        max_cols = max(len(row) for row in table_data)
        for row in table_data:
            while len(row) < max_cols:
                row.append('')
        
        # 테이블 생성 - 행 수에 따라 높이 조정
        table_top = self.config.margin_top + Inches(0.4)
        max_table_height = self.config.height - table_top - self.config.margin_bottom
        
        # 행당 높이 계산 (헤더: 0.3인치, 데이터: 0.4인치)
        num_rows = len(table_data)
        row_height = Inches(0.4)
        calculated_height = row_height * num_rows
        
        # 계산된 높이와 최대 높이 중 작은 값 사용
        table_height = min(calculated_height, max_table_height)
        
        try:
            ppt_table = slide.shapes.add_table(
                len(table_data), max_cols,
                self.config.margin_left, table_top,
                self.config.content_width, table_height
            ).table
            
            # 데이터 채우기
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    if j >= max_cols:
                        continue
                    
                    cell = ppt_table.cell(i, j)
                    cell.text = str(cell_data)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    cell.margin_left = Pt(3)
                    cell.margin_right = Pt(3)
                    cell.margin_top = Pt(2)
                    cell.margin_bottom = Pt(2)
                    
                    cell.fill.background()
                    
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(7)
                        
                        if i == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = self.colors['black']
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.text_frame.word_wrap = False
                        else:
                            paragraph.font.color.rgb = self.colors['gray_800']
                            cell.text_frame.word_wrap = True
                            
                            if j in [2, 5]:
                                paragraph.alignment = PP_ALIGN.LEFT
                            else:
                                paragraph.alignment = PP_ALIGN.CENTER
                            
                            if cell_data == 'Link':
                                paragraph.font.color.rgb = self.colors['link_blue']
                                paragraph.font.underline = True
            
            # 링크 스타일 적용
            for row_idx, col_idx, url in link_data:
                try:
                    cell = ppt_table.cell(row_idx, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = self.colors['link_blue']
                            run.font.underline = True
                except Exception:
                    pass
            
            # 테두리 적용
            from .table_builder import TableBorderStyler
            border_styler = TableBorderStyler(colors=self.colors)
            border_styler.apply_academic_borders(ppt_table, 1, len(table_data), max_cols)
            
        except Exception as e:
            logger.error(f"Evidence 테이블 생성 실패: {e}")
        
        return slide
