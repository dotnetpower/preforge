"""
PowerPoint(.pptx) to Word(.docx) converter

Converts to Word document format while maintaining slide structure.
Analyzes slide structure through preprocessing parsing, distinguishing title/TOC/body.
"""
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
from io import BytesIO
from dataclasses import dataclass, field
import logging
import re

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

from docx import Document as DocxDocument
from docx.shared import Inches, Pt, Cm, RGBColor, Twips
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsmap
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

logger = logging.getLogger(__name__)

# Invalid XML control character pattern
INVALID_XML_CHARS_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]')

# Special symbol replacement map
SPECIAL_CHAR_MAP = {
    '\uf0d8': '▶',  # arrow
    '\uf0fc': '✓',  # checkmark
    '\uf0a7': '•',  # bullet
    '\uf0b7': '•',  # bullet
    '\uf076': '★',  # star
    '\uf0e0': '→',  # arrow
    '\uf0e8': '●',  # circle
    '\uf0ab': '◆',  # diamond
    '\uf0a8': '○',  # empty circle
    '\uf02d': '–',  # dash
    '\uf0b2': '■',  # square
    '\uf06c': '◎',  # double circle
    '\uf0d7': '▼',  # down arrow
    '\uf0de': '▲',  # up arrow
    '\uf0a0': ' ',  # special space
    '\uf020': ' ',  # special space 2
    '\uf06e': '♦',  # diamond 2
}

# Keywords to highlight
HIGHLIGHT_KEYWORDS = [
    'Pathogen', 'Pathogens',
    'Symptom', 'Symptoms', 
    'Diagnosis', 'Diagnostic',
    'Treatment', 'Treatments',
    'Disease', 'Diseases',
    'Epidemiology',
    'Prevention',
    'Transmission',
    'Clinical Features',
]


@dataclass
class GridCell:
    """Grid cell information"""
    row: int
    col: int
    rowspan: int = 1
    colspan: int = 1
    content_type: str = 'empty'  # 'text', 'table', 'image', 'mixed', 'empty'
    left: int = 0  # EMU unit coordinate
    top: int = 0
    width: int = 0
    height: int = 0
    shapes: List[Any] = field(default_factory=list)  # shapes contained in cell


@dataclass
class GridLayout:
    """Slide grid layout"""
    rows: int = 1
    cols: int = 1
    cells: List[GridCell] = field(default_factory=list)
    row_heights: List[int] = field(default_factory=list)  # EMU units
    col_widths: List[int] = field(default_factory=list)  # EMU units
    
    def get_cell(self, row: int, col: int) -> Optional[GridCell]:
        """Return cell at specific position"""
        for cell in self.cells:
            if cell.row == row and cell.col == col:
                return cell
        return None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert grid info to dictionary (for saving parsing results)"""
        return {
            'rows': self.rows,
            'cols': self.cols,
            'row_heights': self.row_heights,
            'col_widths': self.col_widths,
            'cells': [
                {
                    'row': cell.row,
                    'col': cell.col,
                    'rowspan': cell.rowspan,
                    'colspan': cell.colspan,
                    'content_type': cell.content_type,
                    'shape_count': len(cell.shapes),
                }
                for cell in self.cells
            ]
        }


@dataclass
class SlideContent:
    """Parsed slide content"""
    slide_index: int
    slide_type: str  # 'title', 'toc', 'content'
    slide: Any = None  # Slide object reference (for table cell image processing)
    title: Optional[str] = None
    subtitle: Optional[str] = None
    author: Optional[str] = None
    date: Optional[str] = None
    texts: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    toc_items: List[str] = field(default_factory=list)
    section_title: Optional[str] = None  # TOC main title (e.g., "1. Disease")
    grid_layout: Optional[GridLayout] = None  # Grid layout information


@dataclass
class ParsedPresentation:
    """Parsed presentation structure"""
    title: Optional[str] = None
    author: Optional[str] = None
    date: Optional[str] = None
    slides: List[SlideContent] = field(default_factory=list)
    toc_slides: List[int] = field(default_factory=list)  # TOC slide indices
    section_titles: Dict[int, str] = field(default_factory=dict)  # Section title by slide index


def sanitize_text(text: str) -> str:
    """
    Remove XML-incompatible characters and replace special symbols.
    """
    if not text:
        return text
    
    # Replace special symbols
    for old_char, new_char in SPECIAL_CHAR_MAP.items():
        text = text.replace(old_char, new_char)
    
    # Remove control characters
    text = INVALID_XML_CHARS_RE.sub('', text)
    
    return text


def is_highlight_keyword(text: str) -> bool:
    """Check if keyword needs highlighting (check after removing line breaks)"""
    # Check by replacing line breaks with spaces
    text_normalized = ' '.join(text.split()).lower()
    for keyword in HIGHLIGHT_KEYWORDS:
        if keyword.lower() in text_normalized:
            return True
    return False


def normalize_text_for_highlighting(text: str) -> str:
    """Normalize text for highlighting (remove line breaks)"""
    return ' '.join(text.split())


def is_page_number(text: str) -> bool:
    """
    Check if text is a page number
    
    Page number pattern:
    - Only digits (1, 2, 3, ...)
    - Short number (1-3 digits)
    """
    text = text.strip()
    if not text:
        return False
    
    # If only digits and 1-3 digits, consider as page number
    if text.isdigit() and len(text) <= 3:
        return True
    
    return False


class PptxToDocxConverter:
    """Class that converts PowerPoint to Word document"""
    
    def __init__(
        self,
        include_images: bool = True,
        include_tables: bool = True,
        include_notes: bool = False,
        landscape_after_toc: bool = True,
        image_max_width_inches: float = 8.0,
        highlight_keywords: bool = True,
    ):
        """
        Args:
            include_images: Whether to include images
            include_tables: Whether to include tables
            include_notes: Whether to include presenter notes
            landscape_after_toc: Apply landscape layout after TOC
            image_max_width_inches: Maximum image width (inches)
            highlight_keywords: Whether to highlight keywords
        """
        self.include_images = include_images
        self.include_tables = include_tables
        self.include_notes = include_notes
        self.landscape_after_toc = landscape_after_toc
        self.image_max_width_inches = image_max_width_inches
        self.highlight_keywords = highlight_keywords
        
        # Track current section title
        self._current_section_title = None
        self._processed_section_titles = set()
    
    def convert(
        self,
        pptx_path: Path,
        output_path: Optional[Path] = None,
    ) -> Path:
        """
        Convert PPTX file to DOCX
        """
        pptx_path = Path(pptx_path)
        
        if not pptx_path.exists():
            raise FileNotFoundError(f"File not found: {pptx_path}")
        
        if pptx_path.suffix.lower() not in [".pptx", ".ppt"]:
            raise ValueError(f"Unsupported file format: {pptx_path.suffix}")
        
        if output_path is None:
            output_path = pptx_path.with_suffix(".docx")
        else:
            output_path = Path(output_path)
        
        # Initialize state
        self._processed_section_titles = set()
        
        # Step 1: PPTX preprocessing parsing
        prs = Presentation(pptx_path)
        parsed = self._preprocess_presentation(prs)
        
        # Step 2: Create DOCX
        doc = DocxDocument()
        self._setup_document_styles(doc)
        self._copy_metadata(prs, doc)
        
        # Step 3: Convert content
        self._convert_parsed_content(doc, parsed, prs)
        
        # Save
        doc.save(output_path)
        logger.info(f"Conversion complete: {pptx_path} -> {output_path}")
        
        return output_path
    
    def _preprocess_presentation(self, prs: Presentation) -> ParsedPresentation:
        """Preprocess presentation parsing"""
        parsed = ParsedPresentation()
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_content = self._parse_slide(slide, slide_idx)
            parsed.slides.append(slide_content)
            
            # Save title slide information
            if slide_content.slide_type == 'title':
                parsed.title = slide_content.title
                parsed.author = slide_content.author
                parsed.date = slide_content.date
            
            # Save TOC slide index
            if slide_content.slide_type == 'toc':
                parsed.toc_slides.append(slide_idx)
            
            # Save section title
            if slide_content.section_title:
                parsed.section_titles[slide_idx] = slide_content.section_title
        
        return parsed
    
    def _parse_slide(self, slide: Any, slide_idx: int) -> SlideContent:
        """Parse single slide"""
        title_text = self._get_slide_title(slide)
        
        # Determine slide type
        slide_type = self._determine_slide_type(slide, slide_idx, title_text)
        
        content = SlideContent(
            slide_index=slide_idx,
            slide_type=slide_type,
            slide=slide,  # Save slide reference
            title=title_text,
        )
        
        if slide_type == 'title':
            self._parse_title_slide(slide, content)
        elif slide_type == 'toc':
            self._parse_toc_slide(slide, content)
        else:
            self._parse_content_slide(slide, content)
        
        return content
    
    def _determine_slide_type(
        self, 
        slide: Any, 
        slide_idx: int, 
        title_text: Optional[str]
    ) -> str:
        """Determine slide type"""
        # First slide is usually title
        if slide_idx == 1:
            return 'title'
        
        # If title contains 'TOC', it's a TOC slide
        if title_text:
            title_lower = title_text.lower()
            if '목차' in title_lower or 'contents' in title_lower or 'index' in title_lower:
                return 'toc'
        
        # If no title or second slide, check for TOC keyword in body
        if slide_idx == 2 or not title_text:
            all_text = self._get_all_slide_text(slide).lower()
            if '[목차]' in all_text or '목차' in all_text[:20]:
                return 'toc'
        
        return 'content'
    
    def _get_all_slide_text(self, slide: Any) -> str:
        """Extract all text from slide"""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                texts.append(shape.text_frame.text)
            elif hasattr(shape, 'text'):
                texts.append(shape.text)
        return ' '.join(texts)
    
    def _parse_title_slide(self, slide: Any, content: SlideContent):
        """Parse title slide"""
        texts = []
        
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            
            text = sanitize_text(shape.text_frame.text.strip())
            if not text:
                continue
            
            # Title already extracted
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            
            texts.append(text)
        
        # Try to extract date and author from text
        for text in texts:
            # Date pattern (e.g., 2020.10.27, 2020-10-27)
            date_pattern = r'\d{4}[.\-/]\s?\d{1,2}[.\-/]\s?\d{1,2}'
            if re.search(date_pattern, text):
                content.date = text
            # Team/author pattern
            elif '팀' in text or 'Team' in text or any(c in text for c in ['김', '이', '박', '최', '정']):
                content.author = text
            else:
                content.subtitle = text
    
    def _parse_toc_slide(self, slide: Any, content: SlideContent):
        """Parse TOC slide"""
        toc_items = []
        
        for shape in slide.shapes:
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    text = sanitize_text(para.text.strip())
                    if text and len(text) > 1:
                        # Use "[TOC]" text as title and exclude from list
                        text_lower = text.lower()
                        if '[목차]' in text_lower or text_lower == '목차':
                            if not content.title:
                                content.title = text
                        else:
                            toc_items.append(text)
        
        # Set default if no title
        if not content.title:
            content.title = "[TOC]"
        
        content.toc_items = toc_items
    
    def _parse_content_slide(self, slide: Any, content: SlideContent):
        """Parse general content slide"""
        title_shape = slide.shapes.title
        
        # Extract section title (match with TOC main title)
        if content.title:
            # Title starting with number is considered section title
            if re.match(r'^\d+[\.\s]', content.title):
                content.section_title = content.title
        
        # Collect shapes
        shapes_data = []
        for shape in slide.shapes:
            if title_shape and shape == title_shape:
                continue
            
            shape_info = self._get_shape_info(shape)
            if shape_info:
                shapes_data.append(shape_info)
        
        # Sort by position
        shapes_data.sort(key=lambda x: (x['top'], x['left']))
        
        for shape_info in shapes_data:
            if shape_info['type'] == 'text':
                content.texts.append(shape_info)
            elif shape_info['type'] == 'table':
                content.tables.append(shape_info)
            elif shape_info['type'] == 'image':
                content.images.append(shape_info)
            elif shape_info['type'] == 'group':
                # Process group interior
                self._parse_group_shape(shape_info['shape'], content)
        
        # Analyze grid layout
        content.grid_layout = self._analyze_grid_layout(slide, title_shape)
    
    def _parse_group_shape(self, group_shape: Any, content: SlideContent):
        """Parse group shape"""
        for sub_shape in group_shape.shapes:
            shape_info = self._get_shape_info(sub_shape)
            if shape_info:
                if shape_info['type'] == 'text':
                    content.texts.append(shape_info)
                elif shape_info['type'] == 'table':
                    content.tables.append(shape_info)
                elif shape_info['type'] == 'image':
                    content.images.append(shape_info)
                elif shape_info['type'] == 'group':
                    self._parse_group_shape(sub_shape, content)
    
    def _analyze_grid_layout(
        self, 
        slide: Any, 
        title_shape: Any = None
    ) -> GridLayout:
        """
        Analyzes the grid layout of the slide.
        
        Analyzes row/column structure based on shape positions
        and maps what content is in each cell.
        
        Args:
            slide: Slide object
            title_shape: Title shape (excluded from grid)
            
        Returns:
            GridLayout: Analyzed grid layout
        """
        # Collect all shapes (exclude title)
        shapes = []
        for shape in slide.shapes:
            if title_shape and shape == title_shape:
                continue
            
            # Exclude logo, Confidential, and other top-right corner elements
            shape_info = self._extract_shape_bounds(shape)
            if shape_info:
                shapes.append(shape_info)
        
        if not shapes:
            return GridLayout(rows=1, cols=1)
        
        # 1. Find row boundaries based on Y coordinates
        y_coords = sorted(set(
            [s['top'] for s in shapes] + [s['bottom'] for s in shapes]
        ))
        
        # Calculate dynamic threshold based on slide size (~15% of total size)
        # Typical slide: 9144000 x 6858000 EMU (10" x 7.5")
        y_range = max(y_coords) - min(y_coords) if y_coords else 1
        x_range = 0
        
        # 2. Find column boundaries based on X coordinates
        x_coords = sorted(set(
            [s['left'] for s in shapes] + [s['right'] for s in shapes]
        ))
        x_range = max(x_coords) - min(x_coords) if x_coords else 1
        
        # Threshold: set to about 15% of slide size
        # To be divided into 2+ columns, need at least ~50% gap
        y_threshold = max(y_range // 6, 500000)  # Minimum ~0.5 inches
        x_threshold = max(x_range // 6, 500000)  # Minimum ~0.5 inches
        
        row_boundaries = self._find_cluster_boundaries(y_coords, y_threshold)
        col_boundaries = self._find_cluster_boundaries(x_coords, x_threshold)
        
        # 3. Determine row/column count
        num_rows = len(row_boundaries) - 1 if len(row_boundaries) > 1 else 1
        num_cols = len(col_boundaries) - 1 if len(col_boundaries) > 1 else 1
        
        # 4. Create grid cells and map shapes
        cells = []
        row_heights = []
        col_widths = []
        
        # Calculate row heights
        for i in range(num_rows):
            if i + 1 < len(row_boundaries):
                row_heights.append(row_boundaries[i + 1] - row_boundaries[i])
            else:
                row_heights.append(0)
        
        # Calculate column widths
        for i in range(num_cols):
            if i + 1 < len(col_boundaries):
                col_widths.append(col_boundaries[i + 1] - col_boundaries[i])
            else:
                col_widths.append(0)
        
        # Map shapes to each cell
        for row_idx in range(num_rows):
            row_top = row_boundaries[row_idx] if row_idx < len(row_boundaries) else 0
            row_bottom = row_boundaries[row_idx + 1] if row_idx + 1 < len(row_boundaries) else row_top
            
            for col_idx in range(num_cols):
                col_left = col_boundaries[col_idx] if col_idx < len(col_boundaries) else 0
                col_right = col_boundaries[col_idx + 1] if col_idx + 1 < len(col_boundaries) else col_left
                
                cell = GridCell(
                    row=row_idx,
                    col=col_idx,
                    left=col_left,
                    top=row_top,
                    width=col_right - col_left,
                    height=row_bottom - row_top,
                )
                
                # Find shapes belonging to this cell
                cell_shapes = []
                content_types = set()
                
                for s in shapes:
                    # Check if shape center is in cell area
                    center_x = (s['left'] + s['right']) // 2
                    center_y = (s['top'] + s['bottom']) // 2
                    
                    if (col_left <= center_x < col_right and 
                        row_top <= center_y < row_bottom):
                        cell_shapes.append(s['shape'])
                        content_types.add(s['type'])
                
                cell.shapes = cell_shapes
                
                # Determine content type
                if not cell_shapes:
                    cell.content_type = 'empty'
                elif len(content_types) == 1:
                    cell.content_type = list(content_types)[0]
                else:
                    cell.content_type = 'mixed'
                
                cells.append(cell)
        
        # 5. Analyze cell merging (when adjacent cells share the same shape)
        cells = self._detect_cell_spans(cells, num_rows, num_cols)
        
        return GridLayout(
            rows=num_rows,
            cols=num_cols,
            cells=cells,
            row_heights=row_heights,
            col_widths=col_widths,
        )
    
    def _extract_shape_bounds(self, shape: Any) -> Optional[Dict[str, Any]]:
        """Extract shape bounds and type information"""
        try:
            left = shape.left if hasattr(shape, 'left') else 0
            top = shape.top if hasattr(shape, 'top') else 0
            width = shape.width if hasattr(shape, 'width') else 0
            height = shape.height if hasattr(shape, 'height') else 0
            
            # Exclude very small shapes (logos, page numbers, etc.)
            if width < 100000 and height < 100000:  # Less than ~1cm
                return None
            
            # Determine type
            shape_type = 'text'
            if shape.has_table:
                shape_type = 'table'
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = 'image'
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_type = 'group'
            elif hasattr(shape, 'text_frame'):
                shape_type = 'text'
            
            return {
                'shape': shape,
                'left': left,
                'top': top,
                'right': left + width,
                'bottom': top + height,
                'type': shape_type,
            }
        except Exception:
            return None
    
    def _find_cluster_boundaries(
        self, 
        coords: List[int], 
        threshold: int = 200000  # ~2cm EMU
    ) -> List[int]:
        """
        Cluster coordinates to find grid boundaries
        
        Nearby coordinates are merged into the same boundary
        """
        if not coords:
            return [0]
        
        coords = sorted(set(coords))
        boundaries = [coords[0]]
        
        for coord in coords[1:]:
            # Add as new boundary if sufficiently far from previous boundary
            if coord - boundaries[-1] > threshold:
                boundaries.append(coord)
        
        # Add if last coordinate is not in boundaries
        if coords[-1] - boundaries[-1] > threshold:
            boundaries.append(coords[-1])
        
        return boundaries
    
    def _detect_cell_spans(
        self, 
        cells: List[GridCell], 
        num_rows: int, 
        num_cols: int
    ) -> List[GridCell]:
        """
        Detect cell merging (span)
        
        If a single shape spans multiple cells, treat as cell merge
        """
        # Track which cells each shape appears in
        shape_cells = {}  # shape_id -> [(row, col), ...]
        
        for cell in cells:
            for shape in cell.shapes:
                shape_id = id(shape)
                if shape_id not in shape_cells:
                    shape_cells[shape_id] = []
                shape_cells[shape_id].append((cell.row, cell.col))
        
        # Handle shapes that span multiple cells
        merged_cells = set()  # (row, col) -> merged cells
        
        for shape_id, positions in shape_cells.items():
            if len(positions) > 1:
                # Find min/max row, col
                min_row = min(p[0] for p in positions)
                max_row = max(p[0] for p in positions)
                min_col = min(p[1] for p in positions)
                max_col = max(p[1] for p in positions)
                
                # Mark cells in merge range
                for r in range(min_row, max_row + 1):
                    for c in range(min_col, max_col + 1):
                        if (r, c) != (min_row, min_col):
                            merged_cells.add((r, c))
        
        # Update merged cells
        result_cells = []
        for cell in cells:
            if (cell.row, cell.col) in merged_cells:
                # This cell is merged into another cell - skip
                continue
            
            # Calculate rowspan, colspan
            rowspan = 1
            colspan = 1
            
            # Find adjacent cells with same shape
            for shape in cell.shapes:
                shape_id = id(shape)
                if shape_id in shape_cells:
                    positions = shape_cells[shape_id]
                    if len(positions) > 1:
                        min_row = min(p[0] for p in positions)
                        max_row = max(p[0] for p in positions)
                        min_col = min(p[1] for p in positions)
                        max_col = max(p[1] for p in positions)
                        
                        if cell.row == min_row and cell.col == min_col:
                            rowspan = max(rowspan, max_row - min_row + 1)
                            colspan = max(colspan, max_col - min_col + 1)
            
            cell.rowspan = rowspan
            cell.colspan = colspan
            result_cells.append(cell)
        
        return result_cells

    def _convert_parsed_content(
        self, 
        doc: DocxDocument, 
        parsed: ParsedPresentation,
        prs: Presentation
    ):
        """Convert parsed content to DOCX"""
        toc_end_idx = max(parsed.toc_slides) if parsed.toc_slides else 1
        
        for slide_content in parsed.slides:
            # Skip slides with no content (prevent blank pages)
            if self._is_empty_slide(slide_content):
                continue
            
            if slide_content.slide_type == 'title':
                self._create_title_page(doc, slide_content, parsed)
                doc.add_page_break()
                
            elif slide_content.slide_type == 'toc':
                self._create_toc_page(doc, slide_content)
                
                # If last TOC, add section break then landscape layout
                if (self.landscape_after_toc and 
                    slide_content.slide_index == toc_end_idx):
                    self._add_landscape_section(doc)
                else:
                    doc.add_page_break()
                
            else:
                self._create_content_page(doc, slide_content, prs)
                
                # Add page break if not the last slide
                if slide_content.slide_index < len(parsed.slides):
                    doc.add_page_break()
    
    def _is_empty_slide(self, content: SlideContent) -> bool:
        """Check if slide has no content to display"""
        # Title or TOC slides are always included
        if content.slide_type in ('title', 'toc'):
            return False
        
        # Empty slide if no title, text, table, or image
        has_title = bool(content.title and content.title.strip())
        has_texts = any(t.get('shape') and hasattr(t['shape'], 'text_frame') 
                       and t['shape'].text_frame.text.strip() for t in content.texts)
        has_tables = len(content.tables) > 0
        has_images = len(content.images) > 0
        
        return not (has_title or has_texts or has_tables or has_images)
    
    def _create_title_page(
        self, 
        doc: DocxDocument, 
        content: SlideContent,
        parsed: ParsedPresentation
    ):
        """Create title page (fill entire page)"""
        # Empty paragraphs for top margin
        for _ in range(8):
            doc.add_paragraph()
        
        # Main title
        if content.title:
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(sanitize_text(content.title))
            title_run.font.size = Pt(36)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 51, 102)
        
        # Subtitle
        if content.subtitle:
            doc.add_paragraph()
            subtitle_para = doc.add_paragraph()
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            subtitle_run = subtitle_para.add_run(sanitize_text(content.subtitle))
            subtitle_run.font.size = Pt(18)
            subtitle_run.font.color.rgb = RGBColor(100, 100, 100)
        
        # Bottom margin
        for _ in range(6):
            doc.add_paragraph()
        
        # Separator line
        self._add_horizontal_line(doc, color='003366')
        
        doc.add_paragraph()
        
        # Date
        if content.date:
            date_para = doc.add_paragraph()
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            date_run = date_para.add_run(sanitize_text(content.date))
            date_run.font.size = Pt(14)
            date_run.font.color.rgb = RGBColor(80, 80, 80)
        
        # Author
        if content.author:
            author_para = doc.add_paragraph()
            author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            author_run = author_para.add_run(sanitize_text(content.author))
            author_run.font.size = Pt(12)
            author_run.font.color.rgb = RGBColor(100, 100, 100)
    
    def _create_toc_page(self, doc: DocxDocument, content: SlideContent):
        """Create TOC page"""
        # TOC title
        if content.title:
            heading = doc.add_heading(sanitize_text(content.title), level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        # TOC items
        for idx, item in enumerate(content.toc_items, 1):
            item_text = sanitize_text(item)
            
            # Check if number already exists
            if re.match(r'^\d+[\.\s]', item_text):
                # Use as is if already numbered
                para = doc.add_paragraph(item_text, style='List Number')
            else:
                # Add number
                para = doc.add_paragraph(f"{idx}. {item_text}")
            
            para.paragraph_format.space_after = Pt(8)
            para.paragraph_format.left_indent = Cm(1)
            
            for run in para.runs:
                run.font.size = Pt(12)
    
    def _create_content_page(
        self, 
        doc: DocxDocument, 
        content: SlideContent,
        prs: Presentation
    ):
        """Create general content page"""
        # Section title handling (main title from TOC only once)
        if content.section_title:
            if content.section_title not in self._processed_section_titles:
                # First appearance: show section title
                self._processed_section_titles.add(content.section_title)
                heading = doc.add_heading(sanitize_text(content.section_title), level=1)
                self._add_horizontal_line(doc)
            else:
                # Already appeared section: only show subtitle (if any)
                pass
        elif content.title:
            # Regular title
            heading = doc.add_heading(sanitize_text(content.title), level=1)
            self._add_horizontal_line(doc)
        
        # Grid-based rendering if layout exists and has 2+ columns
        grid = content.grid_layout
        if grid and grid.cols >= 2 and self._should_use_grid_layout(grid):
            self._create_grid_based_content(doc, content, grid, prs)
            return
        
        # Original method: sort all content by position and output in order
        all_items = []
        
        for text_info in content.texts:
            all_items.append(('text', text_info['top'], text_info['left'], text_info))
        
        if self.include_tables:
            for table_info in content.tables:
                all_items.append(('table', table_info['top'], table_info['left'], table_info))
        
        if self.include_images:
            for image_info in content.images:
                all_items.append(('image', image_info['top'], image_info['left'], image_info))
        
        # Collect table internal image coordinates (for duplicate prevention)
        table_image_positions = set()
        # Collect table regions (for preventing caption text duplicates within tables)
        table_regions = []
        if self.include_tables and content.slide:
            for table_info in content.tables:
                table_shape = table_info['shape']
                # Save table region
                table_regions.append({
                    'left': table_shape.left,
                    'right': table_shape.left + table_shape.width,
                    'top': table_shape.top,
                    'bottom': table_shape.top + table_shape.height,
                })
                
                img_map = self._find_table_cell_images(
                    content.slide, table_shape, table_shape.table
                )
                for img_list in img_map.values():
                    for img_data in img_list:
                        # Save image position (to exclude from standalone images later)
                        table_image_positions.add(id(img_data.get('blob', b'')))
        
        # Position-based sorting (top to bottom, left to right)
        all_items.sort(key=lambda x: (x[1], x[2]))
        
        for item_type, _, _, item_info in all_items:
            if item_type == 'text':
                # Skip text overlapping with table region (already added as table caption)
                shape = item_info['shape']
                if self._is_shape_in_table_region(shape, table_regions):
                    continue
                self._add_text_from_shape(doc, shape)
            elif item_type == 'table':
                self._add_table_from_shape(
                    doc, item_info['shape'], prs, content.slide
                )
            elif item_type == 'image':
                self._add_image_from_shape(doc, item_info['shape'])
    
    def _should_use_grid_layout(self, grid: GridLayout) -> bool:
        """
        Determine if grid layout should be used
        
        Only use when 2+ columns and meaningful content is arranged side by side
        """
        if grid.cols < 2:
            return False
        
        # Need content in different columns within the same row
        for row_idx in range(grid.rows):
            cols_with_content = []
            for cell in grid.cells:
                if cell.row == row_idx and cell.content_type != 'empty':
                    cols_with_content.append(cell.col)
            
            # If 2+ columns have content in the same row, use grid
            if len(set(cols_with_content)) >= 2:
                return True
        
        return False
    
    def _create_grid_based_content(
        self, 
        doc: DocxDocument, 
        content: SlideContent, 
        grid: GridLayout,
        prs: Presentation
    ):
        """
        Create DOCX content based on grid layout
        
        Uses DOCX tables to implement 2-column layouts, etc.
        """
        # Add grid info comment (for debugging, optional)
        # doc.add_paragraph(f"[Grid: {grid.rows}x{grid.cols}]")
        
        # Process by row
        for row_idx in range(grid.rows):
            row_cells = [c for c in grid.cells if c.row == row_idx]
            row_cells.sort(key=lambda c: c.col)
            
            # Number of columns in this row (considering colspan)
            non_empty_cells = [c for c in row_cells if c.content_type != 'empty']
            
            if not non_empty_cells:
                continue
            
            # Single cell (full width)
            if len(non_empty_cells) == 1:
                cell = non_empty_cells[0]
                self._render_grid_cell_content(doc, cell, prs, content.slide)
            else:
                # Multiple columns: use DOCX table for layout
                self._create_layout_table(doc, non_empty_cells, prs, content.slide, grid)
        
        doc.add_paragraph()
    
    def _create_layout_table(
        self, 
        doc: DocxDocument, 
        cells: List[GridCell], 
        prs: Presentation,
        slide: Any,
        grid: GridLayout
    ):
        """
        Create layout table (no borders)
        
        Map grid cells to DOCX table cells
        """
        num_cols = len(cells)
        layout_table = doc.add_table(rows=1, cols=num_cols)
        layout_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        
        # Remove borders (for layout)
        self._remove_table_borders(layout_table)
        
        # Set column widths (based on grid ratio)
        total_width = sum(grid.col_widths) if grid.col_widths else 1
        
        for idx, cell in enumerate(cells):
            doc_cell = layout_table.rows[0].cells[idx]
            
            # Set cell width
            if grid.col_widths and cell.col < len(grid.col_widths):
                # Width considering colspan
                cell_width = 0
                for c in range(cell.col, min(cell.col + cell.colspan, len(grid.col_widths))):
                    cell_width += grid.col_widths[c]
                
                # Approximate width based on landscape layout (about 10 inches total)
                width_ratio = cell_width / total_width if total_width > 0 else 0.5
                doc_cell.width = Inches(10 * width_ratio)
            
            # Render cell content
            self._render_cell_shapes(doc_cell, cell.shapes, prs, slide)
        
        doc.add_paragraph()
    
    def _remove_table_borders(self, table: Any):
        """Remove table borders (for layout tables)"""
        try:
            tbl = table._tbl
            tbl_pr = tbl.tblPr if tbl.tblPr is not None else OxmlElement('w:tblPr')
            tbl_borders = OxmlElement('w:tblBorders')
            
            for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                border = OxmlElement(f'w:{border_name}')
                border.set(qn('w:val'), 'nil')
                tbl_borders.append(border)
            
            tbl_pr.append(tbl_borders)
            if tbl.tblPr is None:
                tbl.insert(0, tbl_pr)
        except Exception as e:
            logger.debug(f"Failed to remove table borders: {e}")
    
    def _render_grid_cell_content(
        self, 
        doc: DocxDocument, 
        cell: GridCell, 
        prs: Presentation,
        slide: Any
    ):
        """Render single grid cell content (full width)"""
        for shape in cell.shapes:
            if shape.has_table:
                self._add_table_from_shape(doc, shape, prs, slide)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._add_image_from_shape(doc, shape)
            elif hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                self._add_text_from_shape(doc, shape)
    
    def _render_cell_shapes(
        self, 
        doc_cell: Any, 
        shapes: List[Any], 
        prs: Presentation,
        slide: Any
    ):
        """Render shapes in DOCX table cell"""
        for shape in shapes:
            try:
                if shape.has_table:
                    # Nesting table in cell is complex, so just extract text
                    self._add_table_text_to_cell(doc_cell, shape.table)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._add_image_to_cell(doc_cell, shape)
                elif hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    self._add_text_to_cell(doc_cell, shape)
            except Exception as e:
                logger.debug(f"Failed to render cell shape: {e}")
    
    def _add_table_text_to_cell(self, doc_cell: Any, ppt_table: Any):
        """Add PPT table content to DOCX cell as text"""
        for row in ppt_table.rows:
            row_texts = []
            for cell in row.cells:
                text = sanitize_text(cell.text.strip())
                if text:
                    row_texts.append(text)
            
            if row_texts:
                para = doc_cell.add_paragraph(' | '.join(row_texts))
                para.paragraph_format.space_after = Pt(2)
    
    def _add_image_to_cell(self, doc_cell: Any, shape: Any):
        """Add image to DOCX cell"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Apply crop info
            crop_info = self._get_image_crop_info(shape)
            if crop_info and HAS_PIL:
                image_bytes = self._apply_image_crop(image_bytes, crop_info)
            
            # Limit image size within cell (max 3 inches)
            width = shape.width
            width_inches = min(width / 914400, 3.0)
            
            para = doc_cell.add_paragraph()
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = para.add_run()
            run.add_picture(BytesIO(image_bytes), width=Inches(width_inches))
        except Exception as e:
            logger.debug(f"Failed to add cell image: {e}")
    
    def _add_text_to_cell(self, doc_cell: Any, shape: Any):
        """Add text to DOCX cell"""
        if not hasattr(shape, 'text_frame'):
            return
        
        for paragraph in shape.text_frame.paragraphs:
            text = sanitize_text(paragraph.text.strip())
            if not text:
                continue
            
            # Exclude page number
            if is_page_number(text):
                continue
            
            para = doc_cell.add_paragraph(text)
            para.paragraph_format.space_after = Pt(2)
            
            # Keyword highlighting
            if self.highlight_keywords and is_highlight_keyword(text):
                for run in para.runs:
                    self._apply_highlight_style(run)

    def _is_shape_in_table_region(
        self, 
        shape: Any, 
        table_regions: List[Dict[str, int]]
    ) -> bool:
        """
        Check if shape overlaps with table region
        (only short caption text is excluded)
        """
        if not table_regions:
            return False
        
        # Only check shapes with text
        if not hasattr(shape, 'text_frame'):
            return False
        
        text = shape.text_frame.text.strip()
        # Only exclude short caption text (less than 100 characters)
        if len(text) >= 100:
            return False
        
        shape_left = shape.left if hasattr(shape, 'left') else 0
        shape_top = shape.top if hasattr(shape, 'top') else 0
        shape_right = shape_left + (shape.width if hasattr(shape, 'width') else 0)
        shape_bottom = shape_top + (shape.height if hasattr(shape, 'height') else 0)
        
        for region in table_regions:
            # Check if shape overlaps with table region
            if (shape_left < region['right'] and shape_right > region['left'] and
                shape_top < region['bottom'] and shape_bottom > region['top']):
                return True
        
        return False

    def _add_text_from_shape(self, doc: DocxDocument, shape: Any):
        """Extract text from shape and add to document"""
        if not hasattr(shape, 'text_frame'):
            if hasattr(shape, 'text') and shape.text.strip():
                text = sanitize_text(shape.text.strip())
                # Exclude page number
                if text and not is_page_number(text):
                    para = doc.add_paragraph(text)
                    self._apply_keyword_highlighting(para)
            return
        
        for paragraph in shape.text_frame.paragraphs:
            text = sanitize_text(paragraph.text.strip())
            if not text:
                continue
            
            # Exclude page number
            if is_page_number(text):
                continue
            
            level = paragraph.level if hasattr(paragraph, 'level') else 0
            has_bullet = self._has_bullet(paragraph)
            
            # Check if keyword should be highlighted
            is_keyword = self.highlight_keywords and is_highlight_keyword(text)
            
            # If has bullet or indentation, use list style
            if has_bullet or level > 0:
                # Add bullet character if not in text
                if not text.startswith(('•', '-', '▶', '●', '○', '■', '◆')):
                    para = doc.add_paragraph(f"• {text}")
                else:
                    para = doc.add_paragraph(text)
                para.paragraph_format.left_indent = Cm(max(level, 1) * 0.75)
                if is_keyword:
                    for run in para.runs:
                        self._apply_highlight_style(run)
            else:
                para = doc.add_paragraph()
                run = para.add_run(text)
                if is_keyword:
                    self._apply_highlight_style(run)
            
            self._copy_paragraph_style(paragraph, para)
            para.paragraph_format.space_after = Pt(4)
    
    def _has_bullet(self, paragraph: Any) -> bool:
        """Check if PPT paragraph has bullet"""
        try:
            # Check bullet in pptx paragraph object
            if hasattr(paragraph, '_pPr') and paragraph._pPr is not None:
                # Has bullet if not buNone in XML
                pPr = paragraph._pPr
                # Has bullet if buChar, buAutoNum, etc. exist
                bu_elements = pPr.findall('.//' + qn('a:buChar'))
                bu_elements += pPr.findall('.//' + qn('a:buAutoNum'))
                bu_elements += pPr.findall('.//' + qn('a:buBlip'))
                return len(bu_elements) > 0
            
            # If text starts with bullet character, consider it as bullet
            text = paragraph.text.strip() if hasattr(paragraph, 'text') else ''
            bullet_chars = ('•', '-', '–', '▶', '●', '○', '■', '◆', '★', '✓', '*')
            return text.startswith(bullet_chars)
        except Exception:
            return False
    
    def _add_table_from_shape(
        self, 
        doc: DocxDocument, 
        shape: Any, 
        prs: Presentation,
        slide: Any = None
    ):
        """Add table (with cell merge and images)"""
        ppt_table = shape.table
        row_count = len(ppt_table.rows)
        col_count = len(ppt_table.columns)
        
        if row_count == 0 or col_count == 0:
            return
        
        # Find images within table region
        cell_image_map = {}  # (row, col) -> image_bytes
        if slide and self.include_images:
            cell_image_map = self._find_table_cell_images(slide, shape, ppt_table)
        
        # Collect cell merge information
        merge_info = self._get_table_merge_info(ppt_table)
        
        doc_table = doc.add_table(rows=row_count, cols=col_count)
        doc_table.style = 'Table Grid'
        doc_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        
        # First handle cell merges
        self._apply_cell_merges(doc_table, merge_info)
        
        # Track merged cells (skip spanned cells)
        processed_cells = set()
        
        for row_idx, row in enumerate(ppt_table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_key = (row_idx, col_idx)
                
                # Skip already processed cells (merged sub-cells)
                if cell_key in processed_cells:
                    continue
                
                # Check if merged cell
                if cell.is_spanned and not cell.is_merge_origin:
                    continue
                
                doc_cell = doc_table.rows[row_idx].cells[col_idx]
                
                # Check and add images in cell
                has_image = False
                if cell_key in cell_image_map:
                    for img_data in cell_image_map[cell_key]:
                        try:
                            para = doc_cell.add_paragraph()
                            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = para.add_run()
                            run.add_picture(BytesIO(img_data['blob']), width=Inches(1.5))
                            has_image = True
                            
                            # Add caption
                            caption = img_data.get('caption')
                            if caption:
                                cap_para = doc_cell.add_paragraph()
                                cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                cap_run = cap_para.add_run(caption)
                                cap_run.font.size = Pt(9)
                                cap_run.font.italic = True
                                
                        except Exception as e:
                            logger.debug(f"Failed to add cell image: {e}")
                
                # Add text
                cell_text = sanitize_text(cell.text.strip())
                if cell_text:
                    if has_image:
                        # If image exists, add text in new paragraph
                        para = doc_cell.add_paragraph(cell_text)
                    else:
                        doc_cell.text = cell_text
                    
                    # Keyword highlighting
                    if self.highlight_keywords and is_highlight_keyword(cell_text):
                        for para in doc_cell.paragraphs:
                            for run in para.runs:
                                self._apply_highlight_style(run)
                
                # First row as header, make bold
                if row_idx == 0:
                    for para in doc_cell.paragraphs:
                        for run in para.runs:
                            run.font.bold = True
        
        doc.add_paragraph()
    
    def _get_table_merge_info(self, ppt_table: Any) -> List[Dict[str, Any]]:
        """Extract cell merge information from PPTX table"""
        merge_info = []
        
        for row_idx in range(len(ppt_table.rows)):
            for col_idx in range(len(ppt_table.columns)):
                cell = ppt_table.rows[row_idx].cells[col_idx]
                
                if cell.is_merge_origin:
                    # Merge origin cell - calculate colspan and rowspan
                    colspan = 1
                    rowspan = 1
                    
                    # Calculate colspan (number of cells merged to right)
                    for c in range(col_idx + 1, len(ppt_table.columns)):
                        if ppt_table.rows[row_idx].cells[c].is_spanned:
                            # Check if merged from same origin cell
                            colspan += 1
                        else:
                            break
                    
                    # Calculate rowspan (number of cells merged down)
                    for r in range(row_idx + 1, len(ppt_table.rows)):
                        if ppt_table.rows[r].cells[col_idx].is_spanned:
                            rowspan += 1
                        else:
                            break
                    
                    merge_info.append({
                        'row': row_idx,
                        'col': col_idx,
                        'rowspan': rowspan,
                        'colspan': colspan,
                    })
        
        return merge_info
    
    def _apply_cell_merges(self, doc_table: Any, merge_info: List[Dict[str, Any]]):
        """Apply cell merges to DOCX table"""
        for merge in merge_info:
            row = merge['row']
            col = merge['col']
            rowspan = merge['rowspan']
            colspan = merge['colspan']
            
            try:
                # Start cell
                start_cell = doc_table.rows[row].cells[col]
                
                # End cell (bottom-right of merge range)
                end_row = row + rowspan - 1
                end_col = col + colspan - 1
                
                if end_row < len(doc_table.rows) and end_col < len(doc_table.rows[0].cells):
                    end_cell = doc_table.rows[end_row].cells[end_col]
                    start_cell.merge(end_cell)
                    
            except Exception as e:
                logger.debug(f"Cell merge failed ({row},{col}): {e}")
    
    def _find_merge_origin(
        self, 
        ppt_table: Any, 
        row: int, 
        col: int
    ) -> Tuple[int, int]:
        """
        Return merge_origin cell coordinates if the given cell is merged
        
        Args:
            ppt_table: PPTX table object
            row: Row index
            col: Column index
            
        Returns:
            (origin_row, origin_col): merge_origin cell coordinates (original coordinates if not merged)
        """
        try:
            cell = ppt_table.rows[row].cells[col]
            
            # Already merge_origin or non-merged cell
            if cell.is_merge_origin or not cell.is_spanned:
                return (row, col)
            
            # If is_spanned, search up/left for merge_origin
            # First search upward in same column
            for r in range(row - 1, -1, -1):
                check_cell = ppt_table.rows[r].cells[col]
                if check_cell.is_merge_origin:
                    return (r, col)
                elif not check_cell.is_spanned:
                    break
            
            # Search left in same row
            for c in range(col - 1, -1, -1):
                check_cell = ppt_table.rows[row].cells[c]
                if check_cell.is_merge_origin:
                    return (row, c)
                elif not check_cell.is_spanned:
                    break
            
            # Search diagonally (upper left)
            for offset in range(1, max(row, col) + 1):
                r, c = row - offset, col - offset
                if r >= 0 and c >= 0:
                    check_cell = ppt_table.rows[r].cells[c]
                    if check_cell.is_merge_origin:
                        return (r, c)
            
        except Exception as e:
            logger.debug(f"Failed to find merge_origin ({row},{col}): {e}")
        
        return (row, col)

    def _find_table_cell_images(
        self, 
        slide: Any, 
        table_shape: Any, 
        table: Any,
        include_side_images: bool = True
    ) -> Dict[Tuple[int, int], List[Dict[str, Any]]]:
        """
        Find images within table region and map to cells
        
        Args:
            slide: Slide object
            table_shape: Table shape
            table: Table object
            include_side_images: Whether to include images beside (right of) table
        """
        cell_image_map = {}
        
        try:
            # Calculate absolute position of each column
            col_positions = [table_shape.left]
            for i in range(len(table.columns)):
                col_positions.append(col_positions[-1] + table.columns[i].width)
            
            # Calculate absolute position of each row
            row_positions = [table_shape.top]
            for i in range(len(table.rows)):
                row_positions.append(row_positions[-1] + table.rows[i].height)
            
            # Table region
            table_left = table_shape.left
            table_right = col_positions[-1]
            table_top = table_shape.top
            table_bottom = row_positions[-1]
            
            # Find all images from slide
            images_to_check = self._collect_images_from_slide(slide)
            
            # Collect caption info for images beside table
            side_image_captions = {}
            if include_side_images:
                side_image_captions = self._find_side_image_captions(
                    slide, table_right, table_top, table_bottom
                )
            
            # Collect captions in table interior and adjacent area (text below images)
            inner_captions = self._find_table_inner_captions(
                slide, table_left, table_right, table_top, table_bottom
            )
            # Merge side_image_captions and inner_captions
            all_captions = {**side_image_captions, **inner_captions}
            
            # Check if each image belongs to a table cell
            for img_shape in images_to_check:
                try:
                    img_center_x = img_shape.left + img_shape.width // 2
                    img_center_y = img_shape.top + img_shape.height // 2
                    
                    # Table internal image
                    col = -1
                    for i in range(len(col_positions) - 1):
                        if col_positions[i] <= img_center_x < col_positions[i + 1]:
                            col = i
                            break
                    
                    row = -1
                    for i in range(len(row_positions) - 1):
                        if row_positions[i] <= img_center_y < row_positions[i + 1]:
                            row = i
                            break
                    
                    if row >= 0 and col >= 0:
                        # If merged cell, redirect to merge_origin
                        actual_row, actual_col = self._find_merge_origin(
                            table, row, col
                        )
                        
                        # Table internal image
                        cell_key = (actual_row, actual_col)
                        if cell_key not in cell_image_map:
                            cell_image_map[cell_key] = []
                        
                        # Handle ImageWithPosition or regular shape
                        actual_shape = img_shape.shape if hasattr(img_shape, 'shape') else img_shape
                        crop_info = self._get_image_crop_info(actual_shape)
                        image_blob = actual_shape.image.blob
                        if crop_info and HAS_PIL:
                            image_blob = self._apply_image_crop(image_blob, crop_info)
                        
                        # Find caption for table internal images too
                        caption = self._find_caption_for_image(img_shape, all_captions)
                        
                        cell_image_map[cell_key].append({
                            'blob': image_blob,
                            'ext': actual_shape.image.ext,
                            'width': img_shape.width,
                            'height': img_shape.height,
                            'caption': caption,
                        })
                    
                    elif include_side_images and img_center_x > table_right:
                        # Image to right of table - map to row at same height
                        if table_top <= img_center_y <= table_bottom:
                            # Check which row the image center overlaps with
                            for i in range(len(row_positions) - 1):
                                if row_positions[i] <= img_center_y < row_positions[i + 1]:
                                    row = i
                                    break
                            
                            if row >= 0:
                                # Map to last column (column for right-side images)
                                col = len(table.columns) - 1
                                cell_key = (row, col)
                                
                                if cell_key not in cell_image_map:
                                    cell_image_map[cell_key] = []
                                
                                actual_shape = img_shape.shape if hasattr(img_shape, 'shape') else img_shape
                                crop_info = self._get_image_crop_info(actual_shape)
                                image_blob = actual_shape.image.blob
                                if crop_info and HAS_PIL:
                                    image_blob = self._apply_image_crop(image_blob, crop_info)
                                
                                # Find caption
                                caption = self._find_caption_for_image(
                                    img_shape, all_captions
                                )
                                
                                cell_image_map[cell_key].append({
                                    'blob': image_blob,
                                    'ext': actual_shape.image.ext,
                                    'width': img_shape.width,
                                    'height': img_shape.height,
                                    'caption': caption,
                                })
                                
                except Exception as e:
                    logger.debug(f"Failed to check image position: {e}")
                    
        except Exception as e:
            logger.debug(f"Failed to find table cell images: {e}")
        
        return cell_image_map
    
    def _collect_images_from_slide(self, slide: Any) -> List[Any]:
        """Collect all image shapes from slide (including groups, calculate absolute coordinates)"""
        
        @dataclass
        class ImageWithPosition:
            """Image with absolute position information"""
            shape: Any
            abs_left: int
            abs_top: int
            abs_width: int
            abs_height: int
            
            @property
            def left(self):
                return self.abs_left
            
            @property
            def top(self):
                return self.abs_top
            
            @property
            def width(self):
                return self.abs_width
            
            @property
            def height(self):
                return self.abs_height
            
            @property
            def image(self):
                return self.shape.image
        
        images = []
        
        def collect_recursive(shape, parent_left=0, parent_top=0):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Calculate absolute coordinates
                abs_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
                abs_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
                
                images.append(ImageWithPosition(
                    shape=shape,
                    abs_left=abs_left,
                    abs_top=abs_top,
                    abs_width=shape.width if hasattr(shape, 'width') else 0,
                    abs_height=shape.height if hasattr(shape, 'height') else 0,
                ))
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    # Accumulate group coordinates and pass to sub-shapes
                    group_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
                    group_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
                    for sub_shape in shape.shapes:
                        collect_recursive(sub_shape, group_left, group_top)
                except Exception:
                    pass
        
        for shape in slide.shapes:
            collect_recursive(shape)
        
        return images
    
    def _find_side_image_captions(
        self, 
        slide: Any, 
        table_right: int, 
        table_top: int, 
        table_bottom: int
    ) -> Dict[Tuple[int, int], str]:
        """
        Collect image caption text from right side of table (including text in groups)
        
        Returns:
            Dict[(top, bottom), caption_text]: Caption text by position
        """
        captions = {}
        
        def collect_text_recursive(shape, parent_left=0, parent_top=0):
            """Collect text recursively"""
            abs_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
            abs_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
            
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                # Exclude page numbers, short captions only
                if not is_page_number(text) and len(text) < 100:
                    shape_height = shape.height if hasattr(shape, 'height') else 0
                    captions[(abs_top, abs_top + shape_height)] = sanitize_text(text)
            
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub_shape in shape.shapes:
                        collect_text_recursive(sub_shape, abs_left, abs_top)
                except Exception:
                    pass
        
        for shape in slide.shapes:
            # Skip tables
            if shape.has_table:
                continue
            
            shape_left = shape.left if hasattr(shape, 'left') else 0
            # Only elements to right of table
            if shape_left >= table_right * 0.8:
                collect_text_recursive(shape)
        
        return captions
    
    def _find_table_inner_captions(
        self, 
        slide: Any, 
        table_left: int, 
        table_right: int, 
        table_top: int, 
        table_bottom: int
    ) -> Dict[Tuple[int, int], str]:
        """
        Collect image caption text inside table region
        (TEXT_BOX floating above table, not table cell content)
        
        Returns:
            Dict[(top, bottom), caption_text]: Caption text by position
        """
        captions = {}
        
        def collect_text_recursive(shape, parent_left=0, parent_top=0):
            """Collect text recursively"""
            abs_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
            abs_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
            abs_height = shape.height if hasattr(shape, 'height') else 0
            
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                # Exclude page numbers, short captions only
                if not is_page_number(text) and len(text) < 100:
                    captions[(abs_top, abs_top + abs_height)] = sanitize_text(text)
            
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub_shape in shape.shapes:
                        collect_text_recursive(sub_shape, abs_left, abs_top)
                except Exception:
                    pass
        
        for shape in slide.shapes:
            # Skip tables
            if shape.has_table:
                continue
            
            shape_left = shape.left if hasattr(shape, 'left') else 0
            shape_top = shape.top if hasattr(shape, 'top') else 0
            shape_right = shape_left + (shape.width if hasattr(shape, 'width') else 0)
            shape_bottom = shape_top + (shape.height if hasattr(shape, 'height') else 0)
            
            # Elements overlapping table region (inside or adjacent)
            if (shape_left < table_right and shape_right > table_left and
                shape_top < table_bottom and shape_bottom > table_top):
                collect_text_recursive(shape)
        
        return captions
    
    def _find_caption_for_image(
        self, 
        img_shape: Any, 
        captions: Dict[Tuple[int, int], str]
    ) -> Optional[str]:
        """
        Find the closest caption to an image (below or slightly overlapping position)
        """
        if not captions:
            return None
        
        img_top = img_shape.top
        img_bottom = img_shape.top + img_shape.height
        img_center_x = img_shape.left + img_shape.width // 2
        img_height = img_shape.height
        
        best_caption = None
        best_distance = float('inf')
        
        for (cap_top, cap_bottom), text in captions.items():
            # Check if caption is near the bottom of image
            # - Caption starts at or below image lower 50%
            # - Caption is directly below image
            img_lower_half = img_top + img_height // 2
            
            if cap_top >= img_lower_half:
                # Distance between image bottom and caption top
                distance = abs(cap_top - img_bottom)
                
                # If distance is within 50% of image height
                if distance < img_height // 2 and distance < best_distance:
                    best_distance = distance
                    best_caption = text
        
        return best_caption
    
    def _add_image_from_shape(self, doc: DocxDocument, shape: Any):
        """Add image (apply crop info)"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Get crop info
            crop_info = self._get_image_crop_info(shape)
            
            # Crop image if crop info exists
            if crop_info and HAS_PIL:
                image_bytes = self._apply_image_crop(image_bytes, crop_info)
            
            width = shape.width
            height = shape.height
            
            # EMU to Inches
            width_inches = width / 914400
            height_inches = height / 914400
            
            # Limit maximum width
            if width_inches > self.image_max_width_inches:
                scale = self.image_max_width_inches / width_inches
                width_inches = self.image_max_width_inches
                height_inches = height_inches * scale
            
            para = doc.add_paragraph()
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = para.add_run()
            run.add_picture(BytesIO(image_bytes), width=Inches(width_inches))
            
            doc.add_paragraph()
            
        except Exception as e:
            logger.warning(f"Failed to add image: {e}")
    
    def _get_image_crop_info(self, shape: Any) -> Optional[Dict[str, float]]:
        """
        Extract image crop information from PPTX shape
        
        Image crop in PPTX is stored in a:srcRect element:
        - l (left): Left crop ratio (in 1/100000 units)
        - t (top): Top crop ratio
        - r (right): Right crop ratio
        - b (bottom): Bottom crop ratio
        """
        try:
            # Access XML element of shape
            spTree = shape._element
            
            # Find a:srcRect element (inside blipFill)
            nsmap_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            srcRect = spTree.find('.//a:srcRect', nsmap_a)
            
            if srcRect is None:
                return None
            
            # Extract crop values (1/100000 units -> 0~1 ratio)
            crop = {
                'left': int(srcRect.get('l', '0')) / 100000,
                'top': int(srcRect.get('t', '0')) / 100000,
                'right': int(srcRect.get('r', '0')) / 100000,
                'bottom': int(srcRect.get('b', '0')) / 100000,
            }
            
            # Return None if no crop
            if all(v == 0 for v in crop.values()):
                return None
            
            logger.debug(f"Image crop info: {crop}")
            return crop
            
        except Exception as e:
            logger.debug(f"Failed to extract crop info: {e}")
            return None
    
    def _apply_image_crop(
        self, 
        image_bytes: bytes, 
        crop: Dict[str, float]
    ) -> bytes:
        """
        Apply crop to image using PIL
        
        Args:
            image_bytes: Original image bytes
            crop: Crop info (left, top, right, bottom ratios)
        
        Returns:
            Cropped image bytes
        """
        try:
            # Open image
            img = Image.open(BytesIO(image_bytes))
            orig_width, orig_height = img.size
            
            # Calculate crop region (in pixels)
            left = int(orig_width * crop['left'])
            top = int(orig_height * crop['top'])
            right = int(orig_width * (1 - crop['right']))
            bottom = int(orig_height * (1 - crop['bottom']))
            
            # Validation
            if left >= right or top >= bottom:
                logger.debug("Crop region is invalid")
                return image_bytes
            
            # Crop image
            cropped_img = img.crop((left, top, right, bottom))
            
            # Convert to bytes
            output = BytesIO()
            # Keep original format (PNG if unsupported)
            img_format = img.format if img.format else 'PNG'
            if img_format.upper() == 'JPEG':
                # Convert RGBA to RGB
                if cropped_img.mode == 'RGBA':
                    cropped_img = cropped_img.convert('RGB')
                cropped_img.save(output, format='JPEG', quality=95)
            else:
                cropped_img.save(output, format=img_format)
            
            logger.debug(
                f"Image crop complete: {orig_width}x{orig_height} -> "
                f"{right-left}x{bottom-top}"
            )
            
            return output.getvalue()
            
        except Exception as e:
            logger.warning(f"Image crop failed, using original: {e}")
            return image_bytes
    
    def _apply_highlight_style(self, run: Any):
        """Apply highlight style"""
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 102, 153)  # Dark teal
        run.font.size = Pt(12)
    
    def _apply_keyword_highlighting(self, para: Any):
        """Highlight keywords in paragraph"""
        if not self.highlight_keywords:
            return
        
        for run in para.runs:
            if is_highlight_keyword(run.text):
                self._apply_highlight_style(run)
    
    def _add_landscape_section(self, doc: DocxDocument):
        """Add new landscape section"""
        # Start new section
        new_section = doc.add_section()
        
        # Set landscape orientation
        new_section.orientation = WD_ORIENT.LANDSCAPE
        
        # Adjust page size (A4 landscape)
        new_section.page_width = Cm(29.7)
        new_section.page_height = Cm(21.0)
        
        # Set margins
        new_section.top_margin = Cm(1.5)
        new_section.bottom_margin = Cm(1.5)
        new_section.left_margin = Cm(2)
        new_section.right_margin = Cm(2)
    
    def _setup_document_styles(self, doc: DocxDocument):
        """Set up document styles"""
        # Default section margins
        for section in doc.sections:
            section.top_margin = Cm(2)
            section.bottom_margin = Cm(2)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(2.5)
        
        styles = doc.styles
        
        # Heading 1 style
        heading1 = styles['Heading 1']
        heading1.font.size = Pt(18)
        heading1.font.bold = True
        heading1.font.color.rgb = RGBColor(0, 51, 102)
        heading1.paragraph_format.space_before = Pt(12)
        heading1.paragraph_format.space_after = Pt(6)
        
        # Heading 2 style
        heading2 = styles['Heading 2']
        heading2.font.size = Pt(14)
        heading2.font.bold = True
        heading2.font.color.rgb = RGBColor(0, 102, 153)
        heading2.paragraph_format.space_before = Pt(10)
        heading2.paragraph_format.space_after = Pt(4)
        
        # Emphasis style (Intense Emphasis)
        try:
            emphasis = styles['Intense Emphasis']
            emphasis.font.bold = True
            emphasis.font.color.rgb = RGBColor(0, 102, 153)
        except KeyError:
            pass
    
    def _copy_metadata(self, prs: Presentation, doc: DocxDocument):
        """Copy metadata"""
        src = prs.core_properties
        dst = doc.core_properties
        
        if src.title:
            dst.title = src.title
        if src.author:
            dst.author = src.author
        if src.subject:
            dst.subject = src.subject
        if src.keywords:
            dst.keywords = src.keywords
        if src.comments:
            dst.comments = src.comments
    
    def _get_slide_title(self, slide: Any) -> Optional[str]:
        """Extract slide title"""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return sanitize_text(slide.shapes.title.text.strip())
        return None
    
    def _get_shape_info(self, shape: Any, parent_top: int = 0, parent_left: int = 0) -> Optional[dict]:
        """Extract shape type and position information"""
        try:
            top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
            left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
        except Exception:
            top = parent_top
            left = parent_left
        
        shape_type = None
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_type = 'group'
        elif shape.has_table:
            shape_type = 'table'
        elif hasattr(shape, 'image'):
            shape_type = 'image'
        elif hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
            shape_type = 'text'
        elif hasattr(shape, 'text') and shape.text.strip():
            shape_type = 'text'
        
        if shape_type:
            return {
                'shape': shape,
                'type': shape_type,
                'top': top,
                'left': left,
            }
        
        return None
    
    def _copy_paragraph_style(self, src_paragraph: Any, dst_paragraph: Any):
        """Copy paragraph style"""
        try:
            if src_paragraph.runs:
                src_run = src_paragraph.runs[0]
                if dst_paragraph.runs:
                    dst_run = dst_paragraph.runs[0]
                    
                    if src_run.font.bold:
                        dst_run.font.bold = True
                    if src_run.font.italic:
                        dst_run.font.italic = True
                    if src_run.font.size:
                        size_pt = src_run.font.size.pt if src_run.font.size else 11
                        dst_run.font.size = Pt(min(size_pt, 14))
        except Exception as e:
            logger.debug(f"Failed to copy style: {e}")
    
    def _add_horizontal_line(self, doc: DocxDocument, color: str = '003366'):
        """Add horizontal line"""
        para = doc.add_paragraph()
        para.paragraph_format.space_before = Pt(0)
        para.paragraph_format.space_after = Pt(6)
        
        p = para._p
        pPr = p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bottom = OxmlElement('w:bottom')
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '6')
        bottom.set(qn('w:space'), '1')
        bottom.set(qn('w:color'), color)
        pBdr.append(bottom)
        pPr.append(pBdr)
