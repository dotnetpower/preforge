"""
PowerPoint document (.pptx) parser
"""
import json
from pathlib import Path
from typing import Any, Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches

from ..core.document import (
    Document,
    DocumentType,
    DocumentMetadata,
    TextContent,
    TableContent,
    ImageContent,
    CellMerge,
    PageLayout,
    GridCell,
)
from ..core.parser import BaseParser


class PptxParser(BaseParser):
    """PowerPoint document parser"""

    def __init__(self, layout_overrides_path: Optional[Path] = None) -> None:
        self._layout_overrides = self._load_layout_overrides(layout_overrides_path)
    
    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]
    
    @property
    def document_type(self) -> DocumentType:
        return DocumentType.PPTX
    
    def parse(self, file_path: Path) -> Document:
        """Parse PowerPoint document"""
        self.validate_file(file_path)
        
        prs = Presentation(file_path)
        
        # Extract metadata
        metadata = self._extract_metadata(prs)
        
        # Extract text
        text_contents = self._extract_text(prs)
        
        # Extract tables
        tables = self._extract_tables(prs)
        
        # Extract images
        images = self._extract_images(prs)
        
        # Analyze page layouts
        page_layouts = self._analyze_page_layouts(prs, text_contents, tables, images)
        
        return Document(
            file_path=file_path,
            doc_type=self.document_type,
            metadata=metadata,
            text_contents=text_contents,
            tables=tables,
            images=images,
            page_layouts=page_layouts,
            raw_content=prs,
        )

    def _load_layout_overrides(
        self, layout_overrides_path: Optional[Path]
    ) -> Dict[int, Dict[str, Any]]:
        if layout_overrides_path is None:
            layout_overrides_path = Path.cwd() / "private" / "layout_overrides.json"

        if not layout_overrides_path.exists():
            return {}

        try:
            payload = json.loads(layout_overrides_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            return {}

        pages = payload.get("pages", {}) if isinstance(payload, dict) else {}
        overrides: Dict[int, Dict[str, Any]] = {}
        for page_str, config in pages.items():
            try:
                page_num = int(page_str)
            except (TypeError, ValueError):
                continue
            if isinstance(config, dict) and "rows" in config and "cols" in config:
                overrides[page_num] = config
        return overrides
    
    def _extract_metadata(self, prs: Presentation) -> DocumentMetadata:
        """Extract metadata"""
        core_props = prs.core_properties
        
        return DocumentMetadata(
            title=core_props.title,
            author=core_props.author,
            created_at=core_props.created,
            modified_at=core_props.modified,
            subject=core_props.subject,
            keywords=core_props.keywords.split(",") if core_props.keywords else None,
            page_count=len(prs.slides),
            properties={
                "category": core_props.category,
                "comments": core_props.comments,
                "language": core_props.language,
                "slide_count": len(prs.slides),
            }
        )
    
    def _extract_text_from_shape(self, shape, slide_idx: int, text_contents: list, is_title: bool = False, parent_top: int = 0, parent_left: int = 0):
        """Recursively extract text from shapes (GROUP support, absolute coordinate calculation)"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # Title shape is already processed, so skip
        if is_title:
            return
        
        # Current shape's top + parent's cumulative top = absolute position
        try:
            shape_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
            shape_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
        except:
            shape_top = parent_top
            shape_left = parent_left
        
        # For GROUP, recursively process inner shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                self._extract_text_from_shape(sub_shape, slide_idx, text_contents, False, shape_top, shape_left)
        # Process shape with text
        elif hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                text_contents.append(
                    TextContent(
                        text=text,
                        level=0,
                        page_number=slide_idx,
                        position=shape_top,
                        left=shape_left,
                    )
                )
    
    def _extract_text(self, prs: Presentation) -> List[TextContent]:
        """Extract text"""
        text_contents = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Extract slide title
            title_shape = None
            if slide.shapes.title:
                title_shape = slide.shapes.title
                text_contents.append(
                    TextContent(
                        text=title_shape.text,
                        level=1,  # Slide title is level 1
                        style="Title",
                        page_number=slide_idx,
                    )
                )
            
            # Sort shapes by position (top first, then left for same line)
            # Only sort shapes excluding title
            shapes_to_process = []
            for shape in slide.shapes:
                is_title = (title_shape is not None and shape == title_shape)
                if not is_title:
                    shapes_to_process.append(shape)
            
            # Sort shapes with position info by top, left
            def get_position(shape):
                try:
                    return (shape.top, shape.left)
                except:
                    # For shapes without position info, use large value (put at end)
                    return (999999999, 999999999)
            
            shapes_to_process.sort(key=get_position)
            
            # Extract text in sorted order (including GROUP, recursive)
            for shape in shapes_to_process:
                self._extract_text_from_shape(shape, slide_idx, text_contents, False, parent_top=0, parent_left=0)
        
        return text_contents
    
    def _extract_tables(self, prs: Presentation) -> List[TableContent]:
        """Extract tables"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from preforge.core.document import CellImage
        
        tables = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Find all table shapes in the slide
            table_shapes = [s for s in slide.shapes if s.has_table]
            
            for table_shape in table_shapes:
                table = table_shape.table
                
                # First row is considered header (including merged cells)
                headers = []
                for col_idx, cell in enumerate(table.rows[0].cells):
                    if cell.is_spanned:
                        # For merged cells, find merge_origin going left
                        for prev_col_idx in range(col_idx - 1, -1, -1):
                            prev_cell = table.rows[0].cells[prev_col_idx]
                            if prev_cell.is_merge_origin or not prev_cell.is_spanned:
                                # Mark as merged cell or empty string
                                headers.append("")
                                break
                        else:
                            headers.append("")
                    else:
                        headers.append(cell.text.strip())
                
                # Extract remaining rows as data (handle merged cells)
                rows = []
                for row_idx in range(1, len(table.rows)):
                    row = table.rows[row_idx]
                    row_data = []
                    last_value = {}  # Track last merge start value for each column
                    
                    for col_idx, cell in enumerate(row.cells):
                        if cell.is_spanned:
                            # For merged cells, find merge_origin value from same column
                            # Search upward to find merge_origin cell
                            for prev_row_idx in range(row_idx - 1, -1, -1):
                                prev_cell = table.rows[prev_row_idx].cells[col_idx]
                                if prev_cell.is_merge_origin or not prev_cell.is_spanned:
                                    row_data.append(prev_cell.text.strip())
                                    break
                            else:
                                row_data.append("")
                        else:
                            row_data.append(cell.text.strip())
                    
                    rows.append(row_data)
                
                # Find images within table cells
                cell_images = self._find_images_in_table(slide, table_shape, table)
                
                # Extract cell merge information
                cell_merges = []
                for row_idx in range(len(table.rows)):
                    for col_idx in range(len(table.columns)):
                        cell = table.rows[row_idx].cells[col_idx]
                        if cell.is_merge_origin:
                            # Merge origin cell - calculate colspan and rowspan
                            colspan = 1
                            rowspan = 1
                            
                            # Calculate colspan (number of merged cells to the right)
                            for c in range(col_idx + 1, len(table.columns)):
                                if table.rows[row_idx].cells[c].is_spanned:
                                    # Check if spanned from same row
                                    colspan += 1
                                else:
                                    break
                            
                            # Calculate rowspan (number of merged cells downward)
                            for r in range(row_idx + 1, len(table.rows)):
                                if table.rows[r].cells[col_idx].is_spanned:
                                    rowspan += 1
                                else:
                                    break
                            
                            cell_merges.append(CellMerge(
                                row=row_idx,
                                col=col_idx,
                                colspan=colspan,
                                rowspan=rowspan,
                                is_merged=False
                            ))
                        elif cell.is_spanned:
                            # Part of merged cell (do not display)
                            cell_merges.append(CellMerge(
                                row=row_idx,
                                col=col_idx,
                                colspan=1,
                                rowspan=1,
                                is_merged=True
                            ))
                
                tables.append(
                    TableContent(
                        headers=headers,
                        rows=rows,
                        page_number=slide_idx,
                        cell_images=cell_images,
                        cell_merges=cell_merges,
                    )
                )
        
        return tables
    
    def _find_images_in_table(self, slide, table_shape, table) -> List:
        """Find images inside table"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from preforge.core.document import CellImage
        
        cell_images = []
        
        # Calculate absolute position of each column
        col_positions = [table_shape.left]
        for i in range(len(table.columns)):
            col_positions.append(col_positions[-1] + table.columns[i].width)
        
        # Calculate absolute position of each row
        row_positions = [table_shape.top]
        for i in range(len(table.rows)):
            row_positions.append(row_positions[-1] + table.rows[i].height)
        
        # Find all images in slide (direct images + images in groups)
        images_to_check = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images_to_check.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images_to_check.append(sub_shape)
        
        # Check if each image belongs to a table cell
        for img in images_to_check:
            img_center_x = img.left + img.width // 2
            img_center_y = img.top + img.height // 2
            
            # Find which column it belongs to
            col = -1
            for i in range(len(col_positions) - 1):
                if col_positions[i] <= img_center_x < col_positions[i + 1]:
                    col = i
                    break
            
            # Find which row it belongs to
            row = -1
            for i in range(len(row_positions) - 1):
                if row_positions[i] <= img_center_y < row_positions[i + 1]:
                    row = i
                    break
            
            # Only add if inside table
            if row >= 0 and col >= 0:
                try:
                    cell_images.append(
                        CellImage(
                            row=row,
                            col=col,
                            data=img.image.blob,
                            format=img.image.ext,
                            width=img.width,
                            height=img.height,
                        )
                    )
                except Exception:
                    # Ignore if image extraction fails
                    pass
        
        return cell_images
    
    def _is_image_in_table(self, img, tables_info):
        """Check if image is inside table"""
        img_center_x = img.left + img.width // 2
        img_center_y = img.top + img.height // 2
        
        for table_info in tables_info:
            table_shape = table_info['shape']
            table = table_info['table']
            
            # Calculate absolute position of each column
            col_positions = [table_shape.left]
            for i in range(len(table.columns)):
                col_positions.append(col_positions[-1] + table.columns[i].width)
            
            # Calculate absolute position of each row
            row_positions = [table_shape.top]
            for i in range(len(table.rows)):
                row_positions.append(row_positions[-1] + table.rows[i].height)
            
            # Check if image is inside table area
            if (col_positions[0] <= img_center_x < col_positions[-1] and
                row_positions[0] <= img_center_y < row_positions[-1]):
                return True
        
        return False
    
    def _extract_images(self, prs: Presentation) -> List[ImageContent]:
        """Extract images (recursively traverse nested groups, exclude table images)"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        images = []
        
        def extract_from_shape(shape, slide_idx, tables_info, parent_top=0, parent_left=0):
            """Recursively extract images from shape (absolute coordinate calculation)"""
            # Current shape's top + parent's cumulative top = absolute position
            shape_top = (shape.top if hasattr(shape, 'top') else 0) + parent_top
            shape_left = (shape.left if hasattr(shape, 'left') else 0) + parent_left
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Exclude images inside tables
                    if not self._is_image_in_table(shape, tables_info):
                        image = shape.image
                        images.append(
                            ImageContent(
                                data=image.blob,
                                format=image.ext,
                                width=shape.width,
                                height=shape.height,
                                page_number=slide_idx,
                                position=shape_top,
                                left=shape_left,
                            )
                        )
                except Exception:
                    pass
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # For GROUP, recursively traverse sub-shapes
                # Pass cumulative top to sub-shapes
                try:
                    for sub_shape in shape.shapes:
                        extract_from_shape(sub_shape, slide_idx, tables_info, shape_top, shape_left)
                except Exception:
                    pass
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Collect all table info from current slide
            tables_info = []
            for shape in slide.shapes:
                if shape.has_table:
                    tables_info.append({'shape': shape, 'table': shape.table})
            
            # Extract images (pass table info)
            for shape in slide.shapes:
                extract_from_shape(shape, slide_idx, tables_info, parent_top=0)
        
        return images
    
    def _analyze_page_layouts(
        self, 
        prs: Presentation, 
        text_contents: List[TextContent],
        tables: List[TableContent],
        images: List[ImageContent]
    ) -> List[PageLayout]:
        """Analyze grid layout per page"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        page_layouts = []
        
        # Color palette (for visualization)
        colors = [
            '#FFE5E5', '#E5F3FF', '#E5FFE5', '#FFF5E5', '#F5E5FF', '#E5FFFF',
            '#FFD4D4', '#D4E8FF', '#D4FFD4', '#FFEBD4', '#EBD4FF', '#D4FFFF'
        ]
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Slide size (EMU units)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Bottom 90% or more is considered page number/footer
            footer_threshold = slide_height * 90 // 100
            # Top 15% or less is considered title area
            header_threshold = slide_height * 15 // 100
            
            # Collect all content from page (directly analyze shapes)
            content_items = []
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            for shape in slide.shapes:
                top = shape.top
                # Exclude footer area
                if top >= footer_threshold:
                    continue
                
                if shape.has_table:
                    content_items.append({
                        'type': 'table',
                        'id': f'table_{len(content_items)}',
                        'top': top,
                        'left': shape.left,
                        'width': shape.width,
                        'height': shape.height,
                    })
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    content_items.append({
                        'type': 'image',
                        'id': f'image_{len(content_items)}',
                        'top': top,
                        'left': shape.left,
                        'width': shape.width,
                        'height': shape.height,
                    })
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    # Exclude text in header area from layout detection
                    if top < header_threshold:
                        continue
                    content_items.append({
                        'type': 'text',
                        'id': f'text_{len(content_items)}',
                        'top': top,
                        'left': shape.left,
                        'width': shape.width,
                        'height': shape.height,
                    })
            
            override = self._layout_overrides.get(slide_idx)
            if override:
                rows, cols, grid_cells = self._build_layout_from_override(
                    override,
                    content_items,
                    slide_width,
                    slide_height,
                    colors,
                )
            elif not content_items:
                # If no content, set as 1x1 grid
                rows, cols = 1, 1
                grid_cells = [
                    GridCell(
                        row=0, col=0,
                        top=0, left=0,
                        width=slide_width, height=slide_height,
                        color=colors[0]
                    )
                ]
            else:
                # Grid analysis: determine rows/columns based on content positions
                rows, cols, grid_cells = self._detect_grid_layout(
                    content_items, slide_width, slide_height, colors
                )
            
            layout = PageLayout(
                page_number=slide_idx,
                rows=rows,
                cols=cols,
                slide_width=slide_width,
                slide_height=slide_height,
                grid_cells=grid_cells
            )
            page_layouts.append(layout)
        
        return page_layouts

    def _build_layout_from_override(
        self,
        override: Dict[str, Any],
        content_items: List[dict],
        slide_width: int,
        slide_height: int,
        colors: List[str],
    ) -> tuple:
        rows = int(override.get("rows", 1))
        cols = int(override.get("cols", 1))
        row_colspans = override.get("row_colspans")

        if not row_colspans or len(row_colspans) != rows:
            row_colspans = [[1] * cols for _ in range(rows)]

        row_height = slide_height / rows
        col_width = slide_width / cols

        row_boundaries = [int(round(row_height * r)) for r in range(rows + 1)]
        col_boundaries = [int(round(col_width * c)) for c in range(cols + 1)]

        grid_cells: List[GridCell] = []
        color_idx = 0

        for r in range(rows):
            row_top = row_boundaries[r]
            row_bottom = row_boundaries[r + 1]
            row_height_actual = row_bottom - row_top
            col_index = 0

            for span in row_colspans[r]:
                span = int(span)
                left = col_boundaries[col_index]
                right = col_boundaries[min(col_index + span, cols)]
                width = right - left

                cell = GridCell(
                    row=r,
                    col=col_index,
                    top=row_top,
                    left=left,
                    width=width,
                    height=row_height_actual,
                    content_ids=[],
                    color=colors[color_idx % len(colors)],
                    colspan=span,
                )
                grid_cells.append(cell)
                color_idx += 1
                col_index += span

        for item in content_items:
            item_center_x = item['left'] + item['width'] // 2
            item_center_y = item['top'] + item['height'] // 2

            row_idx = min(max(int(item_center_y // row_height), 0), rows - 1)
            col_idx = min(max(int(item_center_x // col_width), 0), cols - 1)

            for cell in grid_cells:
                if cell.row != row_idx:
                    continue
                if cell.col <= col_idx < cell.col + cell.colspan:
                    cell.content_ids.append(item['id'])
                    break

        return rows, cols, grid_cells
    
    def _detect_grid_layout(
        self, 
        content_items: List[dict], 
        slide_width: int, 
        slide_height: int,
        colors: List[str]
    ) -> tuple:
        """Detect grid layout based on content placement
        
        Principles:
        1. Minimize rows (usually 1 row)
        2. Only 2 columns when elements are clearly separated left and right
        3. Symmetric layouts are treated as 1 column
        """
        
        if not content_items:
            return 1, 1, []
        
        # Exclude header area (top 15%), analyze only body elements
        header_threshold = slide_height * 15 // 100
        body_items = [item for item in content_items if item['top'] > header_threshold]
        
        if not body_items:
            return self._create_single_cell_layout(content_items, slide_width, slide_height, colors)
        
        mid_x = slide_width // 2
        
        # Left/right classification
        left_items = [item for item in body_items 
                     if (item['left'] + item['width'] // 2) < mid_x]
        right_items = [item for item in body_items 
                      if (item['left'] + item['width'] // 2) >= mid_x]
        
        # Determine columns
        cols = 1
        
        if left_items and right_items:
            # Elements on both sides
            
            # Check for strict symmetry pattern (same number of elements on both sides with symmetric positions)
            if self._is_symmetric_layout(left_items, right_items, slide_height):
                cols = 1
            else:
                # By default, 2 columns if elements exist on both sides
                cols = 2
        
        # Determine rows: default is 1 row
        rows = 1
        
        return self._build_grid_cells(
            content_items, rows, cols, slide_width, slide_height, colors
        )
    
    def _is_symmetric_layout(
        self, 
        left_items: List[dict], 
        right_items: List[dict], 
        slide_height: int
    ) -> bool:
        """Check if layout is left-right symmetric (table of contents, grid, etc.)
        
        Conditions:
        1. Must have 3 or more elements on each side
        2. Difference in element count between sides must be 2 or less
        3. Many pairs should have similar y positions
        """
        # Need sufficient elements on both sides to determine symmetry
        if len(left_items) < 3 or len(right_items) < 3:
            return False
        
        # If difference in element count is too large, not symmetric
        if abs(len(left_items) - len(right_items)) > 2:
            return False
        
        # Check y position matching
        left_tops = sorted([item['top'] for item in left_items])
        right_tops = sorted([item['top'] for item in right_items])
        
        # Find pairs with similar y positions
        matches = 0
        used_right = set()
        for lt in left_tops:
            for i, rt in enumerate(right_tops):
                if i not in used_right and abs(lt - rt) < slide_height * 0.08:
                    matches += 1
                    used_right.add(i)
                    break
        
        # If 70% or more match, it's symmetric
        min_items = min(len(left_tops), len(right_tops))
        return matches >= min_items * 0.7
        for lt, rt in zip(left_tops, right_tops):
            if abs(lt - rt) > slide_height * 0.08:  # 8% tolerance
                return False
        
        return True
    
    def _create_single_cell_layout(
        self,
        content_items: List[dict],
        slide_width: int,
        slide_height: int,
        colors: List[str]
    ) -> tuple:
        """Create 1x1 grid"""
        cell = GridCell(
            row=0,
            col=0,
            top=0,
            left=0,
            width=slide_width,
            height=slide_height,
            content_ids=[item['id'] for item in content_items],
            color=colors[0]
        )
        return 1, 1, [cell]
    
    def _build_grid_cells(
        self,
        content_items: List[dict],
        rows: int,
        cols: int,
        slide_width: int,
        slide_height: int,
        colors: List[str]
    ) -> tuple:
        """Create grid cells matching row/column count"""
        row_height = slide_height // rows
        col_width = slide_width // cols
        
        grid_cells = []
        color_idx = 0
        
        for r in range(rows):
            row_top = r * row_height
            row_bottom = (r + 1) * row_height if r < rows - 1 else slide_height
            actual_row_height = row_bottom - row_top
            
            for c in range(cols):
                col_left = c * col_width
                col_right = (c + 1) * col_width if c < cols - 1 else slide_width
                actual_col_width = col_right - col_left
                
                # Find content belonging to this cell
                cell_contents = []
                for item in content_items:
                    item_center_x = item['left'] + item['width'] // 2
                    item_center_y = item['top'] + item['height'] // 2
                    
                    if (row_top <= item_center_y < row_bottom and
                        col_left <= item_center_x < col_right):
                        cell_contents.append(item['id'])
                
                cell = GridCell(
                    row=r,
                    col=c,
                    top=row_top,
                    left=col_left,
                    width=actual_col_width,
                    height=actual_row_height,
                    content_ids=cell_contents,
                    color=colors[color_idx % len(colors)]
                )
                grid_cells.append(cell)
                color_idx += 1
        
        return rows, cols, grid_cells
